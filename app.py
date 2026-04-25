import os
import base64
import json
import re
import threading
import time
from datetime import datetime, timedelta, timezone
from flask import Flask, request, jsonify, send_file, Response, Blueprint
from flask_sqlalchemy import SQLAlchemy
from flask_cors import CORS
from flask_jwt_extended import JWTManager, jwt_required, create_access_token, get_jwt_identity
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
import google.generativeai as genai
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import arabic_reshaper
from bidi.algorithm import get_display
from dotenv import load_dotenv
from collections import Counter
import assemblyai as aai
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
import backoff
from sendgrid import SendGridAPIClient
from sendgrid.helpers.mail import Mail, Attachment, FileContent, FileName, FileType, Disposition
import logging
import socket
import yt_dlp
from concurrent.futures import ThreadPoolExecutor
import ssl
import certifi
import PyPDF2
from pptx import Presentation
import pytesseract
from PIL import Image
from bs4 import BeautifulSoup
import requests
import io

# Fix for SSL: CERTIFICATE_VERIFY_FAILED
# This resolves issues with self-signed certificates in local corporate/proxy environments
try:
    _create_unverified_https_context = ssl._create_unverified_context
except AttributeError:
    pass
else:
    ssl._create_default_https_context = _create_unverified_https_context
from queue import Queue, Empty
from flask import request, jsonify, Response
from flask_jwt_extended import jwt_required, get_jwt_identity
try:
    import browser_cookie3
    HAS_BROWSER_COOKIE3 = True
except ImportError:
    HAS_BROWSER_COOKIE3 = False
    logger.warning("browser_cookie3 not found, browser cookie extraction will be disabled.")
import random
import secrets
# Removed Twilio and Firebase - using terminal-based OTP instead
from meeting import init_meeting_system
# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Configure APIs
aai.settings.api_key = os.getenv("ASSEMBLYAI_API_KEY")
gemini_api_key = os.getenv("GEMINI_API_KEY")
if not gemini_api_key:
    raise ValueError("Missing GEMINI_API_KEY in .env file")
genai.configure(api_key=gemini_api_key)

# Flask Configuration
app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*"}}, supports_credentials=False)

@app.after_request
def add_cors_headers(response):
    response.headers['Access-Control-Allow-Origin'] = '*'
    response.headers['Access-Control-Allow-Methods'] = 'GET, POST, PUT, DELETE, OPTIONS'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type, Authorization'
    return response
app.config['UPLOAD_FOLDER'] = "uploads"
app.config['OUTPUT_FOLDER'] = "outputs"
import os
from pathlib import Path

# Fix for deployment data loss: Ensure SQLite uses a persistent data directory if PostgreSQL is not provided.
DATA_DIR = os.getenv('DATA_DIR', os.path.join(os.path.dirname(os.path.abspath(__file__)), 'data'))
os.makedirs(DATA_DIR, exist_ok=True)
default_sqlite = f"sqlite:///{os.path.join(DATA_DIR, 'database.db')}"

# Use DATABASE_URL if available (e.g. Render, Heroku Postgres), otherwise use the persistent SQLite
db_uri = os.getenv('DATABASE_URL', os.getenv('SQLALCHEMY_DATABASE_URI', default_sqlite))
if db_uri.startswith("postgres://"):
    db_uri = db_uri.replace("postgres://", "postgresql://", 1)

app.config['SQLALCHEMY_DATABASE_URI'] = db_uri
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['JWT_SECRET_KEY'] = os.getenv('JWT_SECRET_KEY', 'your-super-secret-jwt-key-change-in-prod')
app.config['JWT_ACCESS_TOKEN_EXPIRES'] = timedelta(days=30)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB
app.config['JWT_ALGORITHM'] = 'HS256'
app.config['JWT_TOKEN_LOCATION'] = ['headers']
app.config['JWT_HEADER_NAME'] = 'Authorization'
app.config['JWT_HEADER_TYPE'] = 'Bearer'

# Define UPLOAD_DIR for compatibility with legacy code
UPLOAD_DIR = app.config['UPLOAD_FOLDER']
UPLOAD_DIR = os.getenv("UPLOAD_DIR", "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

db = SQLAlchemy(app)
jwt = JWTManager(app)

# Thread pool for processing
executor = ThreadPoolExecutor(max_workers=5)  # Adjust max_workers based on your server resources

# Register Unicode Fonts for PDF (Try standard Windows paths)
DEFAULT_FONTS = {}

font_configs = [
    # (Font Name, Filename, Purpose)
    ('Arial', 'arial.ttf', 'Standard/Arabic'),
    ('Nirmala', 'Nirmala.ttf', 'Indic (Hindi/Urdu/etc)'),
    ('MSYH', 'msyh.ttc', 'Chinese (Microsoft YaHei)'),
    ('MSJH', 'msjh.ttc', 'Traditional Chinese'),
    ('MSGothic', 'msgothic.ttc', 'Japanese'),
    ('Gulim', 'gulim.ttc', 'Korean'),
    ('SimSun', 'simsun.ttc', 'Chinese Sim')
]

for font_name, filename, purpose in font_configs:
    try:
        path = os.path.join("C:\\Windows\\Fonts", filename)
        if os.path.exists(path):
            # ReportLab handles some .ttc files if indexed, but it's safer to try
            pdfmetrics.registerFont(TTFont(font_name, path))
            DEFAULT_FONTS[font_name] = font_name
            logger.info(f"Registered font: {font_name} from {filename}")
        else:
            logger.debug(f"Optional font not found: {filename}")
    except Exception as e:
        logger.warning(f"Failed to register font {font_name}: {e}")

DEFAULT_FONT = DEFAULT_FONTS.get('Arial', 'Helvetica')

def get_font_for_text(text):
    """Detect script and return the best registered font name."""
    if not text:
        return DEFAULT_FONT
    
    # Chinese/Japanese/Korean
    if any('\u4e00' <= c <= '\u9fff' or '\u3040' <= c <= '\u30ff' or '\uac00' <= c <= '\ud7af' for c in text):
        return DEFAULT_FONTS.get('MSYH') or DEFAULT_FONTS.get('SimSun') or DEFAULT_FONTS.get('MSGothic') or DEFAULT_FONT
    
    # Indic (Hindi, Bengali, etc.)
    if any('\u0900' <= c <= '\u0d7f' for c in text):
        return DEFAULT_FONTS.get('Nirmala') or DEFAULT_FONT
        
    # Default is Arial or Helvetica
    return DEFAULT_FONT

# Helper for RTL and Unicode shaping
def fix_text_direction(text):
    if not text:
        return ""
    try:
        # Check if text contains RTL characters (Arabic/Urdu/Hebrew range)
        if any('\u0600' <= c <= '\u06FF' or '\u0750' <= c <= '\u077F' or '\u08A0' <= c <= '\u08FF' or '\u0590' <= c <= '\u05FF' for c in text):
            reshaped_text = arabic_reshaper.reshape(text)
            bidi_text = get_display(reshaped_text)
            return bidi_text
        return text
    except Exception as e:
        logger.error(f"Error in fix_text_direction: {e}")
        return text

# Language Mapping for AI Prompt
LANGUAGE_CODE_TO_NAME = {
    'af': 'Afrikaans', 'sq': 'Albanian', 'am': 'Amharic', 'ar': 'Arabic', 'hy': 'Armenian',
    'as': 'Assamese', 'az': 'Azerbaijani', 'eu': 'Basque', 'be': 'Belarusian', 'bn': 'Bengali',
    'bs': 'Bosnian', 'bg': 'Bulgarian', 'my': 'Burmese', 'ca': 'Catalan', 'ceb': 'Cebuano',
    'ny': 'Chichewa', 'zh': 'Chinese (Simplified)', 'zh-tw': 'Chinese (Traditional)',
    'co': 'Corsican', 'hr': 'Croatian', 'cs': 'Czech', 'da': 'Danish', 'nl': 'Dutch',
    'en': 'English', 'eo': 'Esperanto', 'et': 'Estonian', 'fil': 'Filipino', 'fi': 'Finnish',
    'fr': 'French', 'gl': 'Galician', 'ka': 'Georgian', 'de': 'German', 'el': 'Greek',
    'gu': 'Gujarati', 'ht': 'Haitian Creole', 'ha': 'Hausa', 'haw': 'Hawaiian', 'he': 'Hebrew',
    'hi': 'Hindi', 'hmn': 'Hmong', 'hu': 'Hungarian', 'is': 'Icelandic', 'ig': 'Igbo',
    'id': 'Indonesian', 'ga': 'Irish', 'it': 'Italian', 'ja': 'Japanese', 'jv': 'Javanese',
    'kn': 'Kannada', 'kk': 'Kazakh', 'km': 'Khmer', 'ko': 'Korean', 'ku': 'Kurdish',
    'ky': 'Kyrgyz', 'lo': 'Lao', 'la': 'Latin', 'lv': 'Latvian', 'lt': 'Lithuanian',
    'lb': 'Luxembourgish', 'mk': 'Macedonian', 'mg': 'Malagasy', 'ms': 'Malay', 'ml': 'Malayalam',
    'mt': 'Maltese', 'mi': 'Maori', 'mr': 'Marathi', 'mn': 'Mongolian', 'ne': 'Nepali',
    'no': 'Norwegian', 'or': 'Odia (Oriya)', 'ps': 'Pashto', 'fa': 'Persian', 'pl': 'Polish',
    'pt': 'Portuguese', 'pa': 'Punjabi', 'ro': 'Romanian', 'ru': 'Russian', 'sm': 'Samoan',
    'gd': 'Scots Gaelic', 'sr': 'Serbian', 'st': 'Sesotho', 'sn': 'Shona', 'sd': 'Sindhi',
    'si': 'Sinhala', 'sk': 'Slovak', 'sl': 'Slovenian', 'so': 'Somali', 'es': 'Spanish',
    'su': 'Sundanese', 'sw': 'Swahili', 'sv': 'Swedish', 'tg': 'Tajik', 'ta': 'Tamil',
    'tt': 'Tatar', 'te': 'Telugu', 'th': 'Thai', 'tr': 'Turkish', 'tk': 'Turkmen',
    'uk': 'Ukrainian', 'ur': 'Urdu', 'ug': 'Uyghur', 'uz': 'Uzbek', 'vi': 'Vietnamese',
    'cy': 'Welsh', 'xh': 'Xhosa', 'yi': 'Yiddish', 'yo': 'Yoruba', 'zu': 'Zulu'
}

# Database Models
class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    full_name = db.Column(db.String(100), nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    phone_number = db.Column(db.String(20), unique=True, nullable=True)
    password_hash = db.Column(db.String(128), nullable=False)
    bio = db.Column(db.Text, nullable=True)
    reset_token = db.Column(db.String(100), nullable=True)
    reset_token_expiry = db.Column(db.DateTime, nullable=True)
    email_otp = db.Column(db.String(6), nullable=True)
    email_verified = db.Column(db.Boolean, default=False)
    email_token = db.Column(db.String(100), nullable=True)
    email_token_expires = db.Column(db.DateTime, nullable=True)
    email_otp_expires = db.Column(db.DateTime, nullable=True)
    phone_verified = db.Column(db.Boolean, default=False)
    phone_otp = db.Column(db.String(6), nullable=True)
    phone_otp_expires = db.Column(db.DateTime, nullable=True)
    otp_expiry = db.Column(db.DateTime, nullable=True)
    is_verified = db.Column(db.Boolean, default=False)
    password_updated_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    two_factor_pin_hash = db.Column(db.String(128), nullable=True)
    two_factor_enabled = db.Column(db.Boolean, default=False)
    is_active = db.Column(db.Boolean, default=True)
    deactivated_at = db.Column(db.DateTime, nullable=True)
    status = db.Column(db.String(20), default='offline')
    last_seen = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    image = db.Column(db.Text, nullable=True)
    # Privacy settings
    privacy_last_seen     = db.Column(db.String(20), default='everyone')  # everyone|contacts|nobody
    privacy_profile_photo = db.Column(db.String(20), default='everyone')
    privacy_about         = db.Column(db.String(20), default='everyone')
    privacy_read_receipts = db.Column(db.Boolean, default=True)

class Meeting(db.Model):
    __tablename__ = 'meetings'

    id = db.Column(db.Integer, primary_key=True)

    # Existing fields (kept exactly as they were)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    title = db.Column(db.String(200), nullable=False)
    filename = db.Column(db.String(200), nullable=False)
    upload_date = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    status = db.Column(db.String(50), default='uploaded')
    transcription = db.Column(db.Text, default='{}')
    notes = db.Column(db.Text, default='{}')
    language = db.Column(db.String(10), default='en')
    transcript_language = db.Column(db.String(10), default='en')
    has_transcription = db.Column(db.Boolean, default=False)
    has_notes = db.Column(db.Boolean, default=False)
    processing_steps = db.Column(db.Text, default='[]')
    current_step_progress = db.Column(db.Integer, default=0)
    is_favorite = db.Column(db.Boolean, default=False)

    # âœ… Newly added fields from the other model
    filepath = db.Column(db.String(500))                 # Path to audio file
    transcript = db.Column(db.Text)                      # Short preview transcript
    duration = db.Column(db.Float, default=0.0)          # Duration in seconds
    participants_count = db.Column(db.Integer, default=0)  # Total people in meeting
    source = db.Column(db.String(50), default='upload')  # 'upload', 'url', 'live'
    participant_mapping = db.Column(db.Text, default='{}') # JSON mapping of speaker IDs to real names

    # Optional extra fields from your new version (if you want)
    # Not required unless you want them
    # full transcription JSON kept in transcription field

    # Indexes for performance
    __table_args__ = (
        db.Index('idx_user_id', 'user_id'),
        db.Index('idx_upload_date', 'upload_date'),
    )

    def to_dict(self):
        return {
            'id': self.id,
            'user_id': self.user_id,
            'title': self.title,
            'filename': self.filename,
            'filepath': self.filepath,
            'language': self.language,
            'transcript_language': self.transcript_language,
            'transcript': self.transcript[:500] if self.transcript else None,
            'status': self.status,
            'upload_date': self.upload_date.isoformat() if self.upload_date else None,
            'duration': self.duration,
            'participants_count': self.participants_count,
            'has_recording': bool(self.filepath),
            'is_favorite': self.is_favorite
        }


# New Database Models
class Activity(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    type = db.Column(db.String(50), nullable=False)  # 'upload', 'processing', 'completed', 'export', 'share'
    title = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text)
    timestamp = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    meeting_id = db.Column(db.Integer, db.ForeignKey('meetings.id'))  # Fixed: 'meetings' not 'meeting'
    activity_metadata = db.Column(db.Text, default='{}')  # Renamed from 'metadata'

class UserMetrics(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    date = db.Column(db.Date, nullable=False)
    uploads_count = db.Column(db.Integer, default=0)
    processing_time_total = db.Column(db.Integer, default=0)  # in seconds
    exports_count = db.Column(db.Integer, default=0)
    active_minutes = db.Column(db.Integer, default=0)
    languages_used = db.Column(db.Text, default='[]')  # JSON array

class Leaderboard(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    period = db.Column(db.String(20), nullable=False)  # 'weekly', 'monthly', 'all_time'
    score = db.Column(db.Integer, default=0)
    rank = db.Column(db.Integer, default=0)
    updated_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))

class Friendship(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    friend_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    created_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    is_blocked = db.Column(db.Boolean, default=False)
    blocked_by_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    is_pinned = db.Column(db.Boolean, default=False)
    is_muted = db.Column(db.Boolean, default=False)
    is_archived = db.Column(db.Boolean, default=False)
    is_favourite = db.Column(db.Boolean, default=False)
    is_deleted = db.Column(db.Boolean, default=False)

class FriendRequest(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    sender_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    receiver_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    status = db.Column(db.String(20), default='pending')  # 'pending', 'accepted', 'rejected'
    created_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))

class Message(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    sender_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    receiver_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    group_id = db.Column(db.Integer, db.ForeignKey('group.id'), nullable=True)
    broadcast_id = db.Column(db.Integer, db.ForeignKey('broadcast_list.id'), nullable=True)
    text = db.Column(db.Text, nullable=True)
    type = db.Column(db.String(20), default='text')  # 'text', 'voice', 'image', 'file', 'video'
    media_url = db.Column(db.String(500), nullable=True)
    timestamp = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    is_read = db.Column(db.Boolean, default=False)
    reply_to_id = db.Column(db.Integer, db.ForeignKey('message.id'), nullable=True)
    is_deleted = db.Column(db.Boolean, default=False)
    deleted_for = db.Column(db.Text, default='[]')
    is_starred_by = db.Column(db.Text, default='[]')
    reaction = db.Column(db.String(32), nullable=True)  # emoji reaction

class MessageReceipt(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    message_id = db.Column(db.Integer, db.ForeignKey('message.id'), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    delivered_at = db.Column(db.DateTime, nullable=True)
    read_at = db.Column(db.DateTime, nullable=True)
    played_at = db.Column(db.DateTime, nullable=True)
    
class Group(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    image = db.Column(db.String(500), nullable=True)
    description = db.Column(db.Text, nullable=True)
    created_by_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    created_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))

class GroupMember(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    group_id = db.Column(db.Integer, db.ForeignKey('group.id'), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    role = db.Column(db.String(20), default='member')  # 'admin', 'member'
    joined_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    is_pinned = db.Column(db.Boolean, default=False)
    is_muted = db.Column(db.Boolean, default=False)
    is_archived = db.Column(db.Boolean, default=False)
    is_favourite = db.Column(db.Boolean, default=False)
    is_exited = db.Column(db.Boolean, default=False)

class GroupInvite(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    group_id = db.Column(db.Integer, db.ForeignKey('group.id'), nullable=False)
    inviter_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    invitee_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    status = db.Column(db.String(20), default='pending')  # 'pending', 'accepted', 'rejected'
    created_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))

class BroadcastList(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    name = db.Column(db.String(100), nullable=False)
    image = db.Column(db.Text, nullable=True)
    description = db.Column(db.Text, nullable=True)
    created_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    last_used = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))

class BroadcastRecipient(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    list_id = db.Column(db.Integer, db.ForeignKey('broadcast_list.id'), nullable=False)
    recipient_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    role = db.Column(db.String(20), default='member')  # 'owner','admin','member'
    added_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))

class CallLog(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    other_user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True) # Nullable if unknown
    type = db.Column(db.String(20), nullable=False) # 'incoming', 'outgoing', 'missed'
    is_video = db.Column(db.Boolean, default=False)
    timestamp = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    duration = db.Column(db.Integer, default=0) # in seconds
    
# Create tables and handle schema updates
with app.app_context():
    db.create_all()
    
    # Auto-migration for existing databases
    try:
        from sqlalchemy import text
        with db.engine.connect() as conn:
            # Check if reset_token column exists
            result = conn.execute(text("PRAGMA table_info(user)")).fetchall()
            columns = [row[1] for row in result]
            
            if 'reset_token' not in columns:
                logger.info("Migrating database: Adding reset_token column")
                conn.execute(text("ALTER TABLE user ADD COLUMN reset_token TEXT"))
                
            if 'reset_token_expiry' not in columns:
                logger.info("Migrating database: Adding reset_token_expiry column")
                conn.execute(text("ALTER TABLE user ADD COLUMN reset_token_expiry DATETIME"))

            if 'phone_number' not in columns:
                logger.info("Migrating database: Adding phone_number column")
                conn.execute(text("ALTER TABLE user ADD COLUMN phone_number TEXT"))

            if 'email_otp' not in columns:
                logger.info("Migrating database: Adding email_otp column")
                conn.execute(text("ALTER TABLE user ADD COLUMN email_otp TEXT"))

            if 'phone_otp' not in columns:
                logger.info("Migrating database: Adding phone_otp column")
                conn.execute(text("ALTER TABLE user ADD COLUMN phone_otp TEXT"))

            if 'is_verified' not in columns:
                logger.info("Migrating database: Adding is_verified column")
                conn.execute(text("ALTER TABLE user ADD COLUMN is_verified BOOLEAN DEFAULT 0"))

            if 'otp_expiry' not in columns:
                logger.info("Migrating database: Adding otp_expiry column")
                conn.execute(text("ALTER TABLE user ADD COLUMN otp_expiry DATETIME"))

            if 'email_verified' not in columns:
                logger.info("Migrating database: Adding email_verified column")
                conn.execute(text("ALTER TABLE user ADD COLUMN email_verified BOOLEAN DEFAULT 0"))

            if 'email_token' not in columns:
                logger.info("Migrating database: Adding email_token column")
                conn.execute(text("ALTER TABLE user ADD COLUMN email_token TEXT"))

            if 'email_token_expires' not in columns:
                logger.info("Migrating database: Adding email_token_expires column")
                conn.execute(text("ALTER TABLE user ADD COLUMN email_token_expires DATETIME"))

            if 'phone_verified' not in columns:
                logger.info("Migrating database: Adding phone_verified column")
                conn.execute(text("ALTER TABLE user ADD COLUMN phone_verified BOOLEAN DEFAULT 0"))

            if 'phone_otp_expires' not in columns:
                logger.info("Migrating database: Adding phone_otp_expires column")
                conn.execute(text("ALTER TABLE user ADD COLUMN phone_otp_expires DATETIME"))

            if 'email_otp_expires' not in columns:
                logger.info("Migrating database: Adding email_otp_expires column")
                conn.execute(text("ALTER TABLE user ADD COLUMN email_otp_expires DATETIME"))

            if 'bio' not in columns:
                logger.info("Adding bio column to user table")
                conn.execute(text("ALTER TABLE user ADD COLUMN bio TEXT"))

            if 'password_updated_at' not in columns:
                logger.info("Migrating database: Adding password_updated_at column")
                conn.execute(text("ALTER TABLE user ADD COLUMN password_updated_at DATETIME"))
                # Set initial value for existing users
                conn.execute(text("UPDATE user SET password_updated_at = :now WHERE password_updated_at IS NULL"), {"now": datetime.now(timezone.utc)})

            # Migrate broadcast_list table
            bl_info = conn.execute(text("PRAGMA table_info(broadcast_list)")).fetchall()
            bl_cols = [r[1] for r in bl_info]
            if 'image' not in bl_cols:
                conn.execute(text("ALTER TABLE broadcast_list ADD COLUMN image TEXT"))
            if 'description' not in bl_cols:
                conn.execute(text("ALTER TABLE broadcast_list ADD COLUMN description TEXT"))

            # Migrate broadcast_recipient table
            br_info = conn.execute(text("PRAGMA table_info(broadcast_recipient)")).fetchall()
            br_cols = [r[1] for r in br_info]
            if 'role' not in br_cols:
                conn.execute(text("ALTER TABLE broadcast_recipient ADD COLUMN role VARCHAR(20) DEFAULT 'member'"))
            if 'added_at' not in br_cols:
                conn.execute(text("ALTER TABLE broadcast_recipient ADD COLUMN added_at DATETIME"))

            # Migrate message table for broadcast_id
            msg_result = conn.execute(text("PRAGMA table_info(message)")).fetchall()
            msg_columns = [row[1] for row in msg_result]
            if 'broadcast_id' not in msg_columns:
                logger.info("Migrating database: Adding broadcast_id to message table")
                conn.execute(text("ALTER TABLE message ADD COLUMN broadcast_id INTEGER REFERENCES broadcast_list(id)"))
                
            if 'two_factor_pin_hash' not in columns:
                logger.info("Migrating database: Adding two_factor_pin_hash column")
                conn.execute(text("ALTER TABLE user ADD COLUMN two_factor_pin_hash VARCHAR(128)"))

            if 'two_factor_enabled' not in columns:
                logger.info("Migrating database: Adding two_factor_enabled column")
                conn.execute(text("ALTER TABLE user ADD COLUMN two_factor_enabled BOOLEAN DEFAULT 0"))
                
            if 'is_active' not in columns:
                logger.info("Migrating database: Adding is_active column")
                conn.execute(text("ALTER TABLE user ADD COLUMN is_active BOOLEAN DEFAULT 1"))

            if 'deactivated_at' not in columns:
                logger.info("Migrating database: Adding deactivated_at column")
                conn.execute(text("ALTER TABLE user ADD COLUMN deactivated_at DATETIME"))

            if 'status' not in columns:
                logger.info("Migrating database: Adding status column")
                conn.execute(text("ALTER TABLE user ADD COLUMN status VARCHAR(20) DEFAULT 'offline'"))

            if 'last_seen' not in columns:
                logger.info("Migrating database: Adding last_seen column")
                conn.execute(text("ALTER TABLE user ADD COLUMN last_seen DATETIME DEFAULT CURRENT_TIMESTAMP"))

            if 'image' not in columns:
                logger.info("Migrating database: Adding image column to user")
                conn.execute(text("ALTER TABLE user ADD COLUMN image TEXT"))

            if 'privacy_last_seen' not in columns:
                conn.execute(text("ALTER TABLE user ADD COLUMN privacy_last_seen VARCHAR(20) DEFAULT 'everyone'"))
            if 'privacy_profile_photo' not in columns:
                conn.execute(text("ALTER TABLE user ADD COLUMN privacy_profile_photo VARCHAR(20) DEFAULT 'everyone'"))
            if 'privacy_about' not in columns:
                conn.execute(text("ALTER TABLE user ADD COLUMN privacy_about VARCHAR(20) DEFAULT 'everyone'"))
            if 'privacy_read_receipts' not in columns:
                conn.execute(text("ALTER TABLE user ADD COLUMN privacy_read_receipts BOOLEAN DEFAULT 1"))
                
            # Migration for meetings table
            meeting_info = conn.execute(text("PRAGMA table_info(meetings)")).fetchall()
            meeting_columns = [row[1] for row in meeting_info]
            
            if 'source' not in meeting_columns:
                logger.info("Migrating database: Adding source column to meetings")
                conn.execute(text("ALTER TABLE meetings ADD COLUMN source VARCHAR(50) DEFAULT 'upload'"))
                
            if 'participant_mapping' not in meeting_columns:
                logger.info("Migrating database: Adding participant_mapping column to meetings")
                conn.execute(text("ALTER TABLE meetings ADD COLUMN participant_mapping TEXT DEFAULT '{}'"))

            if 'filepath' not in meeting_columns:
                conn.execute(text("ALTER TABLE meetings ADD COLUMN filepath VARCHAR(500)"))
            if 'transcript' not in meeting_columns:
                conn.execute(text("ALTER TABLE meetings ADD COLUMN transcript TEXT"))
            if 'duration' not in meeting_columns:
                conn.execute(text("ALTER TABLE meetings ADD COLUMN duration FLOAT DEFAULT 0.0"))
            if 'participants_count' not in meeting_columns:
                conn.execute(text("ALTER TABLE meetings ADD COLUMN participants_count INTEGER DEFAULT 0"))
            if 'is_favorite' not in meeting_columns:
                conn.execute(text("ALTER TABLE meetings ADD COLUMN is_favorite BOOLEAN DEFAULT 0"))
            if 'processing_steps' not in meeting_columns:
                conn.execute(text("ALTER TABLE meetings ADD COLUMN processing_steps TEXT DEFAULT '[]'"))
            if 'current_step_progress' not in meeting_columns:
                conn.execute(text("ALTER TABLE meetings ADD COLUMN current_step_progress INTEGER DEFAULT 0"))

            # Migration for friendship table
            friendship_info = conn.execute(text("PRAGMA table_info(friendship)")).fetchall()
            friendship_columns = [row[1] for row in friendship_info]
            
            if 'is_pinned' not in friendship_columns:
                logger.info("Migrating database: Adding is_pinned column to friendship")
                conn.execute(text("ALTER TABLE friendship ADD COLUMN is_pinned BOOLEAN DEFAULT 0"))
                
            if 'is_muted' not in friendship_columns:
                logger.info("Migrating database: Adding is_muted column to friendship")
                conn.execute(text("ALTER TABLE friendship ADD COLUMN is_muted BOOLEAN DEFAULT 0"))
                
            if 'is_archived' not in friendship_columns:
                logger.info("Migrating database: Adding is_archived column to friendship")
                conn.execute(text("ALTER TABLE friendship ADD COLUMN is_archived BOOLEAN DEFAULT 0"))
                
            if 'is_favourite' not in friendship_columns:
                logger.info("Migrating database: Adding is_favourite column to friendship")
                conn.execute(text("ALTER TABLE friendship ADD COLUMN is_favourite BOOLEAN DEFAULT 0"))

            if 'is_deleted' not in friendship_columns:
                logger.info("Migrating database: Adding is_deleted column to friendship")
                conn.execute(text("ALTER TABLE friendship ADD COLUMN is_deleted BOOLEAN DEFAULT 0"))

            if 'blocked_by_id' not in friendship_columns:
                logger.info("Migrating database: Adding blocked_by_id column to friendship")
                conn.execute(text("ALTER TABLE friendship ADD COLUMN blocked_by_id INTEGER REFERENCES user(id)"))

            if 'is_blocked' not in friendship_columns:
                logger.info("Migrating database: Adding is_blocked column to friendship")
                conn.execute(text("ALTER TABLE friendship ADD COLUMN is_blocked BOOLEAN DEFAULT 0"))

            # Backfill: fix any blocked rows where blocked_by_id is NULL
            # Set blocked_by_id = user_id for rows where is_blocked=1 and blocked_by_id IS NULL
            conn.execute(text("""
                UPDATE friendship
                SET blocked_by_id = user_id
                WHERE is_blocked = 1 AND blocked_by_id IS NULL
            """))

            # Migration for message table (deletions and stars)
            if 'is_deleted' not in msg_columns:
                logger.info("Migrating database: Adding is_deleted column to message")
                conn.execute(text("ALTER TABLE message ADD COLUMN is_deleted BOOLEAN DEFAULT 0"))
            if 'deleted_for' not in msg_columns:
                logger.info("Migrating database: Adding deleted_for column to message")
                conn.execute(text("ALTER TABLE message ADD COLUMN deleted_for TEXT DEFAULT '[]'"))
            if 'is_starred_by' not in msg_columns:
                logger.info("Migrating database: Adding is_starred_by column to message")
                conn.execute(text("ALTER TABLE message ADD COLUMN is_starred_by TEXT DEFAULT '[]'"))

            # Migration for group_member table (preferences)
            gm_info = conn.execute(text("PRAGMA table_info(group_member)")).fetchall()
            gm_columns = [row[1] for row in gm_info]
            
            if 'is_pinned' not in gm_columns:
                logger.info("Migrating database: Adding is_pinned column to group_member")
                conn.execute(text("ALTER TABLE group_member ADD COLUMN is_pinned BOOLEAN DEFAULT 0"))
            if 'is_muted' not in gm_columns:
                logger.info("Migrating database: Adding is_muted column to group_member")
                conn.execute(text("ALTER TABLE group_member ADD COLUMN is_muted BOOLEAN DEFAULT 0"))
            if 'is_archived' not in gm_columns:
                logger.info("Migrating database: Adding is_archived column to group_member")
                conn.execute(text("ALTER TABLE group_member ADD COLUMN is_archived BOOLEAN DEFAULT 0"))
            if 'is_favourite' not in gm_columns:
                logger.info("Migrating database: Adding is_favourite column to group_member")
                conn.execute(text("ALTER TABLE group_member ADD COLUMN is_favourite BOOLEAN DEFAULT 0"))
            if 'is_exited' not in gm_columns:
                logger.info("Migrating database: Adding is_exited column to group_member")
                conn.execute(text("ALTER TABLE group_member ADD COLUMN is_exited BOOLEAN DEFAULT 0"))

            conn.commit()
    except Exception as e:
        logger.error(f"Schema migration error: {e}")



# Helper Functions
def get_current_user_id():
    """Get current user ID as integer from JWT token"""
    user_id_str = get_jwt_identity()
    return int(user_id_str) if user_id_str else None

def validate_email(email):
    """Validate email format"""
    email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    return bool(re.match(email_pattern, email))

def send_verification_email(email, full_name, otp_code):
    """Send verification OTP via SendGrid"""
    sendgrid_api_key = os.getenv('SENDGRID_API_KEY')
    sendgrid_from_email = os.getenv('SENDGRID_FROM_EMAIL') or 'noreply@talktotextpro.com'
    
    if not sendgrid_api_key:
        logger.warning(f"SendGrid API Key missing. Email OTP for {email}: {otp_code}")
        return False

    email_content = f"""
<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Verify Your Email</title>
  <style>
    body {{
      margin: 0;
      padding: 0;
      font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
      background-color: #111827;
      color: #ffffff;
    }}
    .email-wrapper {{
      padding: 40px 20px;
      background: radial-gradient(ellipse at top left, rgba(56, 189, 248, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at top right, rgba(168, 85, 247, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at bottom, rgba(16, 185, 129, 0.12) 0%, transparent 50%),
                  #111827;
    }}
    .container {{
      max-width: 600px;
      margin: 0 auto;
      background-color: rgba(31, 41, 55, 0.7);
      border: 1px solid rgba(255, 255, 255, 0.1);
      border-radius: 24px;
      overflow: hidden;
      backdrop-filter: blur(10px);
    }}
    .header {{
      padding: 40px 40px 20px;
      text-align: center;
    }}
    .logo {{
      font-size: 24px;
      font-weight: 800;
      letter-spacing: -0.025em;
      background: linear-gradient(to right, #60a5fa, #a855f7);
      -webkit-background-clip: text;
      -webkit-text-fill-color: transparent;
      margin-bottom: 24px;
    }}
    .content {{
      padding: 0 40px 40px;
      text-align: center;
    }}
    h1 {{
      font-size: 28px;
      font-weight: 700;
      margin-bottom: 16px;
      color: #ffffff;
    }}
    p {{
      font-size: 16px;
      line-height: 1.6;
      color: #9ca3af;
      margin-bottom: 32px;
    }}
    .otp-container {{
      margin: 32px 0;
      padding: 24px;
      background: linear-gradient(135deg, rgba(59, 130, 246, 0.1), rgba(139, 92, 246, 0.1));
      border: 2px solid rgba(96, 165, 250, 0.3);
      border-radius: 16px;
    }}
    .otp-label {{
      font-size: 14px;
      font-weight: 600;
      color: #60a5fa;
      text-transform: uppercase;
      letter-spacing: 0.1em;
      margin-bottom: 12px;
    }}
    .otp-code {{
      font-size: 48px;
      font-weight: 800;
      letter-spacing: 0.2em;
      color: #ffffff;
      font-family: 'Courier New', monospace;
      text-shadow: 0 0 20px rgba(96, 165, 250, 0.5);
    }}
    .otp-expiry {{
      font-size: 12px;
      color: #6b7280;
      margin-top: 12px;
    }}
    .footer {{
      padding: 24px 40px;
      background-color: rgba(17, 24, 39, 0.5);
      border-top: 1px solid rgba(255, 255, 255, 0.05);
      text-align: center;
      font-size: 14px;
      color: #6b7280;
    }}
  </style>
</head>
<body>
  <div class="email-wrapper">
    <div class="container">
      <div class="header">
        <div class="logo">TalkToText Pro</div>
      </div>
      <div class="content">
        <h1>Verify Your Email</h1>
        <p>Hello {full_name},</p>
        <p>Welcome to the future of meeting intelligence. Please use the verification code below to verify your email address.</p>
        <div class="otp-container">
          <div class="otp-label">Your Verification Code</div>
          <div class="otp-code">{otp_code}</div>
          <div class="otp-expiry">â±ï¸ Valid for 10 minutes</div>
        </div>
        <p>If you didn't create an account, you can safely ignore this email.</p>
      </div>
      <div class="footer">
        &copy; 2026 TalkToText Pro. Built with â¤ï¸ for productivity.
      </div>
    </div>
  </div>
</body>
</html>
    """
    
    message = Mail(
        from_email=sendgrid_from_email,
        to_emails=email,
        subject="Verify Your Email - TalkToText Pro",
        html_content=email_content
    )
    
    try:
        sg = SendGridAPIClient(api_key=sendgrid_api_key)
        sg.send(message)
        logger.info(f"Verification email with OTP sent to {email}")
        return True
    except Exception as e:
        logger.error(f"Failed to send verification email: {e}")
        # Always log the OTP prominently so the user can see it in terminal logs if SendGrid fails
        logger.warning(f"CRITICAL: Email delivery failed. Manual Verification Code for {email}: {otp_code}")
        return False

# Terminal-Based OTP System (100% Free)
def generate_phone_otp(user_id, phone_number):
    """Generate 6-digit OTP and display prominently in terminal"""
    otp_code = str(random.randint(100000, 999999))
    expiry = datetime.now(timezone.utc) + timedelta(minutes=10)
    
    # Store in database
    user = db.session.get(User, user_id)
    if not user:
        logger.error(f"User {user_id} not found for OTP generation")
        return None
        
    user.phone_otp = otp_code
    user.phone_otp_expires = expiry
    db.session.commit()
    
    # Display prominently in terminal - EASY TO SEE AND COPY
    logger.info("")
    logger.info("=" * 70)
    logger.info("=" * 70)
    logger.info(f"ðŸ“± PHONE OTP FOR: {phone_number}")
    logger.info(f"")
    logger.info(f"   ðŸ”¢ CODE: {otp_code}")
    logger.info(f"")
    logger.info(f"   â° EXPIRES AT: {expiry.strftime('%H:%M:%S')} ({expiry.astimezone().strftime('%I:%M:%S %p')})")
    logger.info(f"   â±ï¸  VALID FOR: 10 minutes")
    logger.info("=" * 70)
    logger.info("=" * 70)
    logger.info("")
    
    return otp_code

def verify_phone_otp(user_id, otp_code):
    """Verify OTP from database"""
    user = db.session.get(User, user_id)
    if not user:
        logger.error(f"User {user_id} not found for OTP verification")
        return False
    
    if not user.phone_otp:
        logger.warning(f"No OTP found for user {user.id}")
        return False
    
  # Check if OTP is expired
    if user.phone_otp_expires:
        if user.phone_otp_expires.tzinfo is None:
            expiry = user.phone_otp_expires.replace(tzinfo=timezone.utc)
        else:
            expiry = user.phone_otp_expires
            
        if datetime.now(timezone.utc) > expiry:
            logger.warning(f"OTP expired for user {user.id}")
            return False
    
    # Verify OTP
    if user.phone_otp == otp_code:
        logger.info(f"âœ“ OTP verified successfully for {user.email}")
        return True
    else:
        logger.warning(f"âœ— Invalid OTP for {user.email}. Got: {otp_code}, Expected: {user.phone_otp}")
        return False

# Email OTP System (Terminal-Based, 100% Free)
def generate_email_otp(user_id, email):
    """Generate 6-digit OTP for email verification and display prominently in terminal"""
    otp_code = str(random.randint(100000, 999999))
    expiry = datetime.now(timezone.utc) + timedelta(minutes=10)
    
    # Store in database
    user = db.session.get(User, user_id)
    if not user:
        logger.error(f"User {user_id} not found for email OTP generation")
        return None
        
    user.email_otp = otp_code
    user.email_otp_expires = expiry
    db.session.commit()
    
    # Display prominently in terminal - EASY TO SEE AND COPY
    logger.info("")
    logger.info("=" * 70)
    logger.info("=" * 70)
    logger.info(f"ðŸ“§ EMAIL OTP FOR: {email}")
    logger.info(f"")
    logger.info(f"   ðŸ”¢ CODE: {otp_code}")
    logger.info(f"")
    logger.info(f"   â° EXPIRES AT: {expiry.strftime('%H:%M:%S')} ({expiry.astimezone().strftime('%I:%M:%S %p')})")
    logger.info(f"   â±ï¸  VALID FOR: 10 minutes")
    logger.info("=" * 70)
    logger.info("=" * 70)
    logger.info("")
    
    return otp_code

def verify_email_otp(user_id, otp_code):
    """Verify email OTP from database"""
    user = db.session.get(User, user_id)
    if not user:
        logger.error(f"User {user_id} not found for email OTP verification")
        return False
    
    if not user.email_otp:
        logger.warning(f"No email OTP found for user {user.id}")
        return False
    
    # Check if OTP is expired
    if user.email_otp_expires:
        if user.email_otp_expires.tzinfo is None:
            expiry = user.email_otp_expires.replace(tzinfo=timezone.utc)
        else:
            expiry = user.email_otp_expires
            
        if datetime.now(timezone.utc) > expiry:
            logger.warning(f"Email OTP expired for user {user.id}")
            return False
    
    # Verify OTP
    if user.email_otp == otp_code:
        logger.info(f"âœ“ Email OTP verified successfully for {user.email}")
        return True
    else:
        logger.warning(f"âœ— Invalid email OTP for {user.email}. Got: {otp_code}, Expected: {user.email_otp}")
        return False

# JWT Error Handlers
@jwt.expired_token_loader
def expired_token_callback(jwt_header, jwt_payload):
    logger.error(f"Token expired: {jwt_payload}")
    return jsonify({
        "error": "Token has expired",
        "message": "Your session has expired. Please log in again.",
        "code": "token_expired"
    }), 401

@jwt.invalid_token_loader
def invalid_token_callback(error):
    logger.error(f"Invalid token: {error}")
    return jsonify({
        "error": "Invalid token",
        "message": "The provided token is invalid. Please log in again.",
        "code": "token_invalid"
    }), 422

@jwt.unauthorized_loader
def missing_token_callback(error):
    logger.error(f"Missing token: {error}")
    return jsonify({
        "error": "Authorization required",
        "message": "Please provide a valid authorization token.",
        "code": "token_missing"
    }), 401

socketio = init_meeting_system(app, db)

# Global Error Handlers
@app.errorhandler(422)
def handle_unprocessable_entity(e):
    logger.error(f"422 Unprocessable Entity: {e}")
    return jsonify({
        "error": "Request validation failed",
        "message": "The request could not be processed. Please check your data and try again.",
        "details": str(e)
    }), 422

@app.errorhandler(413)
def handle_request_entity_too_large(e):
    logger.error(f"413 Request too large: {e}")
    return jsonify({
        "error": "File too large",
        "message": "The uploaded file exceeds the maximum size limit of 100MB."
    }), 413

def update_processing_step(meeting, step_name, status, error=None):
    try:
        steps = json.loads(meeting.processing_steps or '[]')
    except:
        steps = []
    
    timestamp = datetime.now(timezone.utc).isoformat()
    step = next((s for s in steps if s["step"] == step_name), None)
    if step:
        step.update({"status": status, "error": error, "timestamp": timestamp})
    else:
        steps.append({"step": step_name, "status": status, "error": error, "timestamp": timestamp})
    
    meeting.processing_steps = json.dumps(steps)
    if status == "in_progress":
        meeting.current_step_progress = 0
    elif status == "success":
        meeting.current_step_progress = 0
    db.session.commit()
    logger.info(f"Updated step {step_name} to {status} for meeting {meeting.id}")

def simulate_step_progress(meeting_id, step_name, duration_seconds=8):
    """Simulate realistic progress for each processing step"""
    logger.info(f"Starting progress simulation for {step_name} on meeting {meeting_id}")
    
    progress_points = [10, 20, 35, 50, 65, 75, 85, 95, 100]
    interval = duration_seconds / len(progress_points)
    
    for progress in progress_points:
        try:
            with app.app_context():
                meeting = db.session.get(Meeting, meeting_id)
                if not meeting:
                    break
                
                steps = json.loads(meeting.processing_steps or '[]')
                current_step = next((s for s in steps if s["step"] == step_name), None)
                
                if not current_step or current_step["status"] != "in_progress":
                    logger.info(f"Step {step_name} no longer in progress, stopping simulation")
                    break
                
                meeting.current_step_progress = progress
                db.session.commit()
                logger.debug(f"{step_name} progress: {progress}%")
                
                if progress < 100:
                    time.sleep(interval)
                    
        except Exception as e:
            logger.error(f"Progress simulation error for {step_name}: {e}")
            break

@backoff.on_exception(backoff.expo, Exception, max_tries=3, max_time=120)
def call_gemini_api(prompt, model="gemini-2.5-flash", image=None):
    model_instance = genai.GenerativeModel(model)
    if image:
        response = model_instance.generate_content([prompt, image])
    else:
        response = model_instance.generate_content(prompt)
        
    if not response or not hasattr(response, 'text') or not response.text:
        raise ValueError("Invalid or empty response from Gemini API")
    return response

def cleanup_transcript(text):
    """ZERO TOUCH PASS-THROUGH. No filtering allowed to ensure 100% accuracy."""
    if not text:
        return ""
    return text # Absolutely no changes, not even whitespace, as requested

def extract_comprehensive_content(transcript_text):
    """Extract comprehensive content for longer transcripts"""
    if not transcript_text:
        return [], [], []
    
    sentences = []
    raw_sentences = re.split(r'[.!?]+|\n\n+', transcript_text)
    
    for sentence in raw_sentences:
        sentence = sentence.strip()
        if sentence: # Keep ALL sentences, no matter how short
            sentence = re.sub(r'\s+', ' ', sentence)
            sentences.append(sentence)
    
    meaningful_sentences = []
    filler_patterns = [
        r'\b(um|uh|ah|er|hmm|well|you know|i mean|like|so|basically|actually|literally)\b',
        r'\b(kind of|sort of|i guess|i think maybe|probably|perhaps)\b',
        r'^(okay|alright|right|yes|no|yeah|yep|sure)\.?\s*$'
    ]
    
    for sentence in sentences:
        sentence_lower = sentence.lower()
        filler_count = sum(len(re.findall(pattern, sentence_lower)) for pattern in filler_patterns)
        word_count = len(sentence.split())
        
        if word_count > 3 and (filler_count / max(word_count, 1)) < 0.4:
            meaningful_sentences.append(sentence)
    
    words = re.findall(r'\b[a-zA-Z]{3,}\b', transcript_text.lower())
    word_freq = Counter(words)
    
    stop_words = {
        'the', 'and', 'that', 'have', 'for', 'not', 'with', 'you', 'this', 'but', 'his', 'from', 
        'they', 'she', 'her', 'been', 'than', 'its', 'were', 'said', 'each', 'which', 'their',
        'time', 'will', 'way', 'about', 'many', 'then', 'them', 'these', 'two', 'more', 'very',
        'what', 'know', 'just', 'first', 'get', 'has', 'him', 'had', 'let', 'put', 'too', 'old',
        'any', 'after', 'move', 'why', 'before', 'here', 'how', 'all', 'both', 'each', 'few',
        'more', 'most', 'other', 'some', 'such', 'only', 'own', 'same', 'than', 'too', 'very',
        'can', 'will', 'now', 'during', 'before', 'after', 'above', 'below', 'between', 'into',
        'through', 'during', 'before', 'after', 'above', 'below', 'between', 'being', 'where',
        'when', 'who', 'whom', 'whose', 'would', 'could', 'should', 'might', 'must', 'shall',
        'going', 'want', 'need', 'like', 'look', 'come', 'came', 'take', 'took', 'make', 'made'
    }
    
    topics = [word for word, count in word_freq.most_common(50) 
            if word not in stop_words and count > 2 and len(word) > 3]
    
    phrases = []
    words_list = transcript_text.lower().split()
    for i in range(len(words_list) - 1):
        phrase = f"{words_list[i]} {words_list[i+1]}"
        if len(phrase) > 6:
            phrases.append(phrase)
    
    phrase_freq = Counter(phrases)
    key_phrases = [phrase for phrase, count in phrase_freq.most_common(20) 
                  if count > 1 and not any(stop in phrase for stop in ['the ', 'and ', 'that ', 'with '])]
    
    return meaningful_sentences, topics, key_phrases

def generate_comprehensive_summary(transcript_text, title, meaningful_sentences, topics, key_phrases):
    """Generate comprehensive summary for longer transcripts"""
    word_count = len(transcript_text.split())
    char_count = len(transcript_text)
    
    logger.info(f"Generating summary for {word_count} words ({char_count} characters)")
    
    if word_count < 100:
        return f"Brief meeting '{title}' with limited discussion content. The session covered basic topics and concluded with minimal actionable items."
    
    topic_context = ""
    if topics:
        main_topics = topics[:8]
        topic_context = f"Primary discussion areas included: {', '.join(main_topics)}. "
    
    phrase_context = ""
    if key_phrases:
        main_phrases = key_phrases[:5]
        phrase_context = f"Key recurring themes: {', '.join(main_phrases)}. "
    
    context_sentences = meaningful_sentences[:8] if meaningful_sentences else []
    
    if word_count > 3000:
        summary_template = f"""The comprehensive meeting '{title}' involved extensive discussions spanning multiple topics and themes. {topic_context}{phrase_context}

The session demonstrated thorough exploration of complex subjects with detailed participant engagement. Key discussion segments covered strategic planning, operational considerations, and collaborative decision-making processes. 

Participants provided in-depth analysis of current situations, explored various solutions, and established clear pathways for implementation. The meeting maintained strong focus on actionable outcomes while addressing both immediate concerns and long-term objectives.

The extended dialogue allowed for comprehensive coverage of all relevant topics, ensuring stakeholder alignment and establishing concrete next steps for continued progress."""
    
    elif word_count > 1500:
        summary_template = f"""The detailed meeting '{title}' covered substantial ground across multiple discussion areas. {topic_context}{phrase_context}

Participants engaged in meaningful dialogue addressing key operational and strategic considerations. The session provided comprehensive coverage of relevant topics while maintaining focus on practical outcomes and actionable decisions.

Discussion included thorough analysis of current challenges, evaluation of potential solutions, and establishment of clear implementation strategies. The meeting concluded with well-defined next steps and stakeholder commitments."""
    
    else:
        summary_template = f"""The meeting '{title}' addressed important business topics through focused discussion. {topic_context}{phrase_context}

Participants contributed valuable insights leading to clear outcomes and actionable decisions. The session maintained good momentum while covering all essential agenda items effectively."""
    
    if context_sentences:
        key_content = '. '.join(context_sentences[:3])
        summary_template += f"\n\nKey highlights: {key_content}"
    
    return summary_template

def process_long_transcript_in_chunks(transcript_text, title, max_chunk_size=25000):
    """Process very long transcripts in chunks to avoid token limits"""
    if len(transcript_text) <= max_chunk_size:
        return None
    
    logger.info(f"Processing long transcript in chunks: {len(transcript_text)} characters")
    
    chunks = []
    words = transcript_text.split()
    
    current_chunk = []
    current_size = 0
    
    for word in words:
        current_chunk.append(word)
        current_size += len(word) + 1
        
        if current_size >= max_chunk_size:
            chunks.append(' '.join(current_chunk))
            current_chunk = []
            current_size = 0
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    logger.info(f"Created {len(chunks)} chunks for processing")
    
    all_summaries = []
    all_key_points = []
    all_action_items = []
    all_decisions = []
    
    for i, chunk in enumerate(chunks):
        logger.info(f"Processing chunk {i+1}/{len(chunks)}")
        
        chunk_prompt = f"""
You are an expert transcript analyzer. Analyze this section of a longer meeting/video transcript and extract structured insights.

MEETING: {title}
SECTION {i+1} of {len(chunks)}:
{chunk}

INSTRUCTIONS:
1. Read this section carefully (do NOT skip details).
2. Write a *detailed summary* proportional to the length of this section:
   - Minimum 3â€“5 sentences for short sections.
   - Longer sections require proportionally longer summaries.
   - Capture purpose, discussion flow, and outcomes clearly.
3. Extract *key_points*:
   - List ALL factual and significant points discussed.
   - Be exhaustive â€” capture context, data, and specifics.
4. Extract *action_items*:
   - Include every task, follow-up, or responsibility mentioned in this section.
   - Note owners/deadlines if available.
5. Extract *decisions*:
   - List all decisions or agreements reached in this section.
   - If no decisions, return an empty array [].
6. Extract *sentiment*:
   - Describe the tone (positive, negative, neutral, mixed).
   - Mention engagement level of participants.

Respond ONLY with valid JSON in this format:
{{
  "summary": "Proportional summary of this section",
  "key_points": ["Detailed key point 1", "Detailed key point 2", "..."],
  "action_items": ["Action item with owner/deadline if available", "..."],
  "decisions": ["Decision 1 with context", "Decision 2 with context"],
  "sentiment": "Overall tone + engagement"
}}

CRITICAL RULES:

- Respond with **VALID JSON only** â€” no markdown, text, or extra explanations.
- Maintain **multilingual support** â€” detect and respond in the same language as the transcript.
- Use **only chunk content** â€” never invent or assume information.
- Ensure factual accuracy and internal consistency.
- Follow **system behavior strictly** â€” analytical, non-creative, and objective.
- Set **temperature = 0** for deterministic and reproducible output.

- Respond with VALID JSON only.
- Use only transcript content, no external assumptions.
- Do not output text outside JSON.

"""
        try:
            response = call_gemini_api(chunk_prompt)
            chunk_result = json.loads(response.text.strip())
            
            if chunk_result.get("summary"):
                all_summaries.append(f"Section {i+1}: {chunk_result['summary']}")
            
            if chunk_result.get("key_points"):
                all_key_points.extend(chunk_result["key_points"])
                
            if chunk_result.get("action_items"):
                all_action_items.extend(chunk_result["action_items"])
                
            if chunk_result.get("decisions"):
                all_decisions.extend(chunk_result["decisions"])
                
        except Exception as e:
            logger.error(f"Failed to process chunk {i+1}: {e}")
            all_summaries.append(f"Section {i+1}: Discussion continued with various topics addressed")
            all_key_points.append(f"Continued discussion from section {i+1}")
    
    combined_summary = f"This comprehensive meeting '{title}' covered extensive topics across multiple discussion segments. " + " ".join(all_summaries[:5])
    
    return {
        "summary": combined_summary,
        "key_points": all_key_points,
        "action_items": all_action_items,
        "decisions": all_decisions
    }

def start_processing(meeting_id):
    logger.info(f"Starting processing thread for meeting ID: {meeting_id}")
    
    with app.app_context():
        meeting = db.session.get(Meeting, meeting_id)
        if not meeting:
            logger.error(f"Meeting {meeting_id} not found")
            return

        try:
            meeting.status = 'processing'
            initial_steps = [
                {"step": "transcription", "status": "pending", "timestamp": "", "error": None},
                {"step": "translation", "status": "pending", "timestamp": "", "error": None},
                {"step": "optimization", "status": "pending", "timestamp": "", "error": None},
                {"step": "ai_generation", "status": "pending", "timestamp": "", "error": None}
            ]
            meeting.processing_steps = json.dumps(initial_steps)
            meeting.current_step_progress = 0
            db.session.commit()
            
            # Log processing activity
            log_activity(
                user_id=meeting.user_id,
                activity_type="processing",
                title=f"Processing: {meeting.title}",
                description=f"Started processing {meeting.filename}",
                meeting_id=meeting.id,
                metadata={"filename": meeting.filename, "title": meeting.title}
            )
            
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], meeting.filename)
            if not os.path.exists(filepath):
                raise Exception(f"File not found: {filepath}")

            # === OPTIMIZATION: Extract audio to speed up AssemblyAI upload ===
            # This reduces file size from 100MB+ (1080p video) to ~5MB (audio only)
            audio_only_path = filepath
            is_video = meeting.filename.lower().endswith(('.mp4', '.avi', '.mov', '.mkv', '.webm'))
            
            if is_video:
                try:
                    logger.info(f"Extracting audio from video to speed up transcription: {meeting.filename}")
                    temp_audio_name = f"temp_audio_{meeting.id}_{int(time.time())}.mp3"
                    temp_audio_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_audio_name)
                    
                    # Use ffmpeg to extract audio quickly with low bitrate (plenty for transcription)
                    import subprocess
                    cmd = [
                        'ffmpeg', '-y', '-i', filepath,
                        '-vn', '-acodec', 'libmp3lame', '-ab', '128k', '-ar', '16000', '-ac', '1',
                        temp_audio_path
                    ]
                    subprocess.run(cmd, capture_output=True, check=True)
                    
                    if os.path.exists(temp_audio_path):
                        audio_only_path = temp_audio_path
                        logger.info(f"Audio extracted successfully: {os.path.getsize(audio_only_path) / 1024 / 1024:.2f} MB")
                    else:
                        logger.warning("ffmpeg did not create audio file, falling back to full file")
                except Exception as audio_err:
                    logger.error(f"Failed to extract audio: {audio_err}. Falling back to full file.")
            
            # Use the optimized audio path for AssemblyAI
            transcription_file_path = audio_only_path

            # Get transcript language from meeting
            transcript_lang = meeting.transcript_language or 'en'
            logger.info(f"Processing with transcript language: {transcript_lang}")


            # Step 1: Transcription with language specified
            logger.info("Starting transcription...")
            update_processing_step(meeting, "transcription", "in_progress")

            progress_thread = threading.Thread(target=simulate_step_progress, args=(meeting_id, "transcription", 15))
            progress_thread.daemon = True
            progress_thread.start()
            
            # Initialize transcription variables
            raw_text = ""
            translated_text = ""
            optimized_text = ""
            utterances = []
            speaker_map = {}
            next_speaker_letter = 'A'
            detected_lang = transcript_lang # Default to requested language
            
            # Map frontend language codes to AssemblyAI language codes
# Map frontend language codes to AssemblyAI language codes
            language_mapping = {
                'en': 'en',
                'es': 'es',
                'fr': 'fr',
                'de': 'de',
                'it': 'it',
                'pt': 'pt',
                'nl': 'nl',
                'hi': 'hi',
                'ja': 'ja',
                'zh': 'zh',        # Chinese (Simplified/Traditional)
                'fi': 'fi',
                'ko': 'ko',
                'pl': 'pl',
                'ru': 'ru',
                'tr': 'tr',
                'uk': 'uk',
                'vi': 'vi',
                'id': 'id',
                'th': 'th',
                'he': 'he',
                'ms': 'ms',        # Malay
                'fil': 'tl',       # Filipino (Tagalog)
                'fa': 'fa',        # Persian
                'sv': 'sv',
                'no': 'no',
                'da': 'da',
                'el': 'el',        # Greek
                'cs': 'cs',
                'ro': 'ro',
                'hu': 'hu',
                'ta': 'ta',
                'te': 'te',
                'mr': 'mr',
                'ur': 'ur',
                'bn': 'bn',
                # Additional languages supported by AssemblyAI
                'af': 'af', 'sq': 'sq', 'am': 'am', 'ar': 'ar', 'hy': 'hy', 'as': 'as',
                'az': 'az', 'eu': 'eu', 'be': 'be', 'bs': 'bs', 'bg': 'bg', 'my': 'my',
                'ca': 'ca', 'ceb': 'ceb', 'ny': 'ny', 'co': 'co', 'hr': 'hr', 'eo': 'eo',
                'et': 'et', 'gl': 'gl', 'ka': 'ka', 'gu': 'gu', 'ht': 'ht', 'ha': 'ha',
                'haw': 'haw', 'is': 'is', 'ig': 'ig', 'jv': 'jv', 'kn': 'kn', 'kk': 'kk',
                'km': 'km', 'rw': 'rw', 'ku': 'ku', 'ky': 'ky', 'lo': 'lo', 'la': 'la',
                'lv': 'lv', 'lt': 'lt', 'lb': 'lb', 'mk': 'mk', 'mg': 'mg', 'ml': 'ml',
                'mt': 'mt', 'mi': 'mi', 'ne': 'ne', 'or': 'or', 'ps': 'ps', 'pa': 'pa',
                'sm': 'sm', 'gd': 'gd', 'sr': 'sr', 'st': 'st', 'sn': 'sn', 'sd': 'sd',
                'si': 'si', 'sk': 'sk', 'sl': 'sl', 'so': 'so', 'su': 'su', 'sw': 'sw',
                'tk': 'tk', 'uz': 'uz', 'cy': 'cy', 'xh': 'xh', 'yi': 'yi', 'yo': 'yo',
                'zu': 'zu', 'tg': 'tg', 'tt': 'tt', 'ug': 'ug', 'hmn': 'hmn'
            }

            # Get the correct AssemblyAI language code
            assemblyai_lang = language_mapping.get(transcript_lang, 'en')
            logger.info(f"Using AssemblyAI language code: {assemblyai_lang}")
            
            # Get meeting source and mapping for labeling
            meeting_source = meeting.source or 'upload'
            try:
                participant_mapping = json.loads(meeting.participant_mapping or '{}')
            except:
                participant_mapping = {}

            user = db.session.get(User, meeting.user_id)
            creator_name = user.full_name if user else "Speaker A"

            # Configure transcription based on language support
            # auto_highlights only works with English
            try:
                if transcript_lang == 'en' or assemblyai_lang == 'en':
                    # Full feature set for English
                    # Full feature set for English with high accuracy settings
                    transcription_config = aai.TranscriptionConfig(
                        speaker_labels=True,
                        auto_highlights=True,
                        language_code='en',
                        punctuate=True,
                        format_text=True,
                        disfluencies=True,
                        speech_model=aai.SpeechModel.best # Maximum accuracy for English
                    )
                    logger.info("Using full feature set for English transcription (High Accuracy)")
                else:
                    # Limited features for non-English languages
                    # Only use features that are language-agnostic
                    # Limited features for non-English languages
                    # Only use features that are language-agnostic
                    # Only enable disfluencies for English as it causes errors in many other languages (like Urdu)
                    disfluencies_enabled = (transcript_lang == 'en' or assemblyai_lang == 'en')
                    
                    transcription_config = aai.TranscriptionConfig(
                        speaker_labels=True,
                        language_code=assemblyai_lang,
                        punctuate=True,
                        format_text=True,
                        disfluencies=disfluencies_enabled 
                    )
                    logger.info(f"Using basic transcription for {transcript_lang} (language code: {assemblyai_lang}) (High Accuracy)")
                    logger.info(f"Note: auto_highlights is not available for {transcript_lang}")

                # Increase global timeout for large file uploads
                aai.settings.http_timeout = 600
                
                transcriber = aai.Transcriber(config=transcription_config)
                
                # Add retry logic for transcription call
                max_retries = 3
                transcript = None
                for attempt in range(max_retries):
                    try:
                        logger.info(f"Transcription attempt {attempt + 1}...")
                        transcript = transcriber.transcribe(transcription_file_path)
                        if transcript.status != aai.TranscriptStatus.error:
                            break
                    except Exception as e:
                        if attempt == max_retries - 1:
                            raise e
                        logger.warning(f"Transcription attempt {attempt + 1} failed: {e}. Retrying...")
                        time.sleep(5)

                if not transcript or transcript.status == aai.TranscriptStatus.error:
                    raise Exception(f"Transcription failed: {transcript.error if transcript else 'No response'}")


                if hasattr(transcript, 'utterances') and transcript.utterances:
                    for u in transcript.utterances:
                        if u.speaker not in speaker_map:
                            speaker_map[u.speaker] = next_speaker_letter
                            next_speaker_letter = chr(ord(next_speaker_letter) + 1)
                        
                        # Determine label
                        u_speaker_label = ""
                        if meeting_source == 'live':
                            # For live meetings, try to map Speaker A to the creator
                            if speaker_map[u.speaker] == 'A':
                                u_speaker_label = creator_name
                            else:
                                u_speaker_label = f"Speaker {speaker_map[u.speaker]}"
                        else:
                            # For uploads/URLs, user wants NO speaker labels
                            u_speaker_label = ""
                        
                        # Capture word-level data for this utterance if available
                        u_words = []
                        if hasattr(transcript, 'words') and transcript.words:
                            for w in transcript.words:
                                if w.start >= u.start and w.end <= u.end:
                                        u_words.append({
                                            "text": w.text, # Use RAW text from word
                                            "start": w.start / 1000.0,
                                            "end": w.end / 1000.0
                                        })

                        # Reconstruct the sentence manually from words to ensure NO word is skipped by the formatter
                        manual_text = " ".join([w['text'] for w in u_words]) if u_words else u.text

                        utterances.append({
                            "speaker": u_speaker_label,
                            "text": manual_text, # Use the manually joined raw words
                            "start": u.start / 1000.0, # seconds
                            "end": u.end / 1000.0, # seconds
                            "words": u_words
                        })
                
                raw_text = transcript.text or ""
                detected_lang = getattr(transcript, 'language_code', transcript_lang)
                
                if not raw_text or len(raw_text.strip()) == 0:
                    raise Exception("Transcription returned empty text")
                
                logger.info(f"Transcription completed in {transcript_lang}: {len(raw_text)} characters, {len(raw_text.split())} words")

            except Exception as transcription_error:
                logger.error(f"Transcription error with language {transcript_lang}: {transcription_error}")
                
                # Fallback: Try with automatic language detection
                logger.info("Attempting fallback with automatic language detection...")
                try:
                    transcription_config = aai.TranscriptionConfig(
                        speaker_labels=True,
                        language_detection=True,  # Let AssemblyAI auto-detect
                        punctuate=True,
                        format_text=True,
                        disfluencies=False # Safer for auto-detection fallback
                    )
                    transcriber = aai.Transcriber(config=transcription_config)
                    transcript = transcriber.transcribe(transcription_file_path)
                    
                    if transcript.status == aai.TranscriptStatus.error:
                        raise Exception(f"Fallback transcription failed: {transcript.error}")
                    
                    raw_text = transcript.text
                    utterances = []
                    speaker_map = {}
                    next_speaker_letter = 'A'

                    if hasattr(transcript, 'utterances') and transcript.utterances:
                        for u in transcript.utterances:
                            if u.speaker not in speaker_map:
                                speaker_map[u.speaker] = next_speaker_letter
                                next_speaker_letter = chr(ord(next_speaker_letter) + 1)
                            
                            u_words = []
                            if hasattr(transcript, 'words') and transcript.words:
                                for w in transcript.words:
                                    if w.start >= u.start and w.end <= u.end:
                                        u_words.append({
                                            "text": w.text, # Use RAW text
                                            "start": w.start / 1000.0,
                                            "end": w.end / 1000.0
                                        })

                            # Determine label
                            u_speaker_label = ""
                            if meeting_source == 'live':
                                if speaker_map[u.speaker] == 'A':
                                    u_speaker_label = creator_name
                                else:
                                    u_speaker_label = f"Speaker {speaker_map[u.speaker]}"
                            else:
                                u_speaker_label = ""

                            # Reconstruct manually
                            manual_text = " ".join([w['text'] for w in u_words]) if u_words else u.text

                            utterances.append({
                                "speaker": u_speaker_label,
                                "text": manual_text, # Literal word-for-word
                                "start": u.start / 1000.0,
                                "end": u.end / 1000.0,
                                "words": u_words
                            })
                    detected_lang = getattr(transcript, 'language_code', 'unknown')
                    logger.info(f"Fallback successful. Detected language: {detected_lang}")
                    logger.info(f"Transcription: {len(raw_text)} characters, {len(raw_text.split())} words")
                    
                except Exception as fallback_error:
                    logger.error(f"Fallback transcription also failed: {fallback_error}")
                    raise Exception(f"Transcription failed: {str(transcription_error)}. Fallback also failed: {str(fallback_error)}")
            
            progress_thread.join(timeout=2) # Don't wait too long for the progress simulation to finish
            
            # Cleanup temporary audio file if it was created
            if transcription_file_path != filepath and os.path.exists(transcription_file_path):
                try:
                    os.remove(transcription_file_path)
                    logger.info(f"Cleaned up temporary audio file: {transcription_file_path}")
                except Exception as e:
                    logger.warning(f"Failed to cleanup temp audio: {e}")

            update_processing_step(meeting, "transcription", "success")
            time.sleep(1)
            
            # Steps 2 & 3: Translation and Optimization
            logger.info(f"Starting translation from {detected_lang} to {transcript_lang}...")
            update_processing_step(meeting, "translation", "in_progress")
            progress_thread = threading.Thread(target=simulate_step_progress, args=(meeting_id, "translation", 10))
            progress_thread.daemon = True
            progress_thread.start()
            
            # Actual Translation Logic
            if raw_text and transcript_lang != detected_lang:
                try:
                    target_lang_name = LANGUAGE_CODE_TO_NAME.get(transcript_lang, 'English')
                    logger.info(f"Translating transcript to {target_lang_name}...")
                    
                    translation_prompt = f"""
                    You are a professional literal translator. Translate the following transcript into {target_lang_name}.
                    CRITICAL: You must translate EVERY SINGLE WORD. 
                    Do NOT summarize. Do NOT skip repetitions. Do NOT remove filler words.
                    The user requires 100% literal accuracy, including stutters and disfluencies.
                    Maintain the same paragraph structure and speaker labels.
                    Return ONLY the translated text.
                    
                    TRANSCRIPT:
                    {raw_text}
                    """
                    
                    response = call_gemini_api(translation_prompt)
                    translated_text = response.text.strip()
                    logger.info(f"Translation successful. Length: {len(translated_text)}")
                    
                    # Also translate individual utterances if they exist
                    if utterances:
                        logger.info("Translating individual utterances...")
                        for u in utterances:
                            if u.get('text'):
                                u_prompt = f"""
                                Translate the following literal meeting transcript line to {target_lang_name}.
                                CRITICAL: Keep EVERY SINGLE WORD, stutter, filler (um, ah, uh), and repetition.
                                Do NOT clean up the speech.
                                Text: {u['text']}
                                """
                                try:
                                    u_resp = call_gemini_api(u_prompt)
                                    u['text'] = u_resp.text.strip()
                                except Exception as u_err:
                                    logger.warning(f"Failed to translate utterance: {u_err}")
                except Exception as eval_err:
                    logger.error(f"Translation failed: {eval_err}")
                    translated_text = raw_text # Fallback to original
            else:
                translated_text = raw_text
                
            progress_thread.join(timeout=12)
            update_processing_step(meeting, "translation", "success")
            time.sleep(1)
            
            logger.info("Starting optimization and cleanup...")
            update_processing_step(meeting, "optimization", "in_progress")
            progress_thread = threading.Thread(target=simulate_step_progress, args=(meeting_id, "optimization", 8))
            progress_thread.daemon = True
            progress_thread.start()
            
            if utterances:
                optimized_lines = []
                for u in utterances:
                    # Use a lighter cleanup for the full transcript to avoid losing words
                    u_clean = u["text"].strip()
                    if u_clean:
                        if u['speaker']:
                            optimized_lines.append(f"{u['speaker']}: {u_clean}")
                        else:
                            optimized_lines.append(u_clean)
                optimized_text = "\n".join(optimized_lines)
            else:
                optimized_text = cleanup_transcript(translated_text)
            meaningful_sentences, topics, key_phrases = extract_comprehensive_content(optimized_text)
            progress_thread.join(timeout=10)
            update_processing_step(meeting, "optimization", "success")
            time.sleep(1)
            
            # Step 4: AI Generation - Responds in the user's selected transcript language
            target_lang_name = LANGUAGE_CODE_TO_NAME.get(transcript_lang, 'English')
            logger.info(f"Starting enhanced AI generation in {target_lang_name}...")
            update_processing_step(meeting, "ai_generation", "in_progress")
            
            progress_thread = threading.Thread(target=simulate_step_progress, args=(meeting_id, "ai_generation", 20))
            progress_thread.daemon = True
            progress_thread.start()
            
            # UPDATED PROMPT - Responds in the transcript language, keeping transcript in original language
            improved_prompt = f"""
You are an expert meeting/video analyst. Analyze the full transcript carefully and extract detailed structured insights.

MEETING: {meeting.title}
TRANSCRIPT ({len(optimized_text)} characters):
{optimized_text}

INSTRUCTIONS:
1. Read the entire transcript carefully (do NOT skip or compress too much).
2. Write a *comprehensive summary* in {target_lang_name} proportional to transcript length:
   - For short transcripts (1â€“10 min): at least 5â€“8 sentences.
   - For medium transcripts (10â€“30 min): 10â€“15+ sentences.
   - For long transcripts (30â€“60 min): 20+ sentences.
   - For very long transcripts (1â€“3 hours): multi-paragraph, fully covering purpose, flow, arguments, updates, examples, and outcomes.
   - Do not miss or shorten any major theme.
3. Extract *key_points* in {target_lang_name}:
   - These should be direct, factual insights from the transcript.
   - Be exhaustive â€” capture every significant discussion, update, concern, and highlight in detail.
   - Include specific names, data, examples, and references if mentioned.
4. Extract *action_items* in {target_lang_name}:
   - List every task, follow-up, or responsibility discussed.
   - Include owners/teams and deadlines if available.
   - If implied strongly, include inferred actions with context.
5. Extract *decisions* in {target_lang_name}:
   - List ALL actual decisions/resolutions reached with full context.
   - Mention if decisions are pending, partial, or conditional.
6. Analyze *sentiment* in {target_lang_name}:
   - Describe overall tone (positive, negative, neutral, or mixed).
   - Mention engagement levels (collaborative, tense, distracted, highly engaged, etc.).

Return ONLY valid JSON with this exact structure:
{{
  "summary": "Comprehensive proportional summary of the full transcript (IN {target_lang_name})",
  "key_points": [
    "Factual key point 1 (IN {target_lang_name})",
    "Factual key point 2 (IN {target_lang_name})",
    "... (exhaustive list)"
  ],
  "action_items": [
    "Action item with owner/deadline if available (IN {target_lang_name})",
    "Another action item (IN {target_lang_name})"
  ],
  "decisions": [
    "Decision with context (IN {target_lang_name})",
    "Pending or partial decision if discussed (IN {target_lang_name})"
  ],
  "sentiment": "Overall tone + engagement level (IN {target_lang_name})"
}}

CRITICAL RULES:

- Respond with **VALID JSON only** â€” no markdown, no explanations, no extra text.
- Maintain **multilingual support** â€” analyze and respond in {target_lang_name}.
- Use **only transcript content**, no external assumptions or guesses.
- Ensure **logical consistency** across summary, key points, actions, and decisions.
- Keep **temperature = 0** to ensure factual, deterministic, and repeatable output.
- Follow **system behavior strictly** â€” act as an analytical assistant, not a creative one.
- If information is missing or unclear, state it clearly within JSON (e.g., "No decisions mentioned" or "Action owner not specified").

- Respond with VALID JSON only.
- Use only transcript content, no external assumptions.
- Do not output text outside JSON.
- ALL responses MUST be in {target_lang_name}.

"""
            processed_data = None
            
            # Handle long transcripts
            if len(optimized_text) > 30000:
                logger.info("Processing long transcript in chunks")
                processed_data = process_long_transcript_in_chunks(optimized_text, meeting.title)
            else:
                try:
                    logger.info("Sending request to Gemini API...")
                    response = call_gemini_api(improved_prompt, model="gemini-2.5-flash")
                    ai_response = response.text.strip()
                    
                    # Clean the response
                    if ai_response.startswith("```json"):
                        ai_response = ai_response[7:]
                    if ai_response.endswith("```"):
                        ai_response = ai_response[:-3]
                    ai_response = ai_response.strip()
                    
                    logger.info(f"AI Response received: {len(ai_response)} characters")
                    
                    processed_data = json.loads(ai_response)
                    logger.info(f"AI processing successful - {len(processed_data.get('key_points', []))} key points extracted")
                    
                    # Validate that we have real action items and decisions
                    if not processed_data.get('action_items') or len(processed_data.get('action_items', [])) == 0:
                        logger.warning("No action items extracted")
                        processed_data['action_items'] = ["No specific action items identified. Follow-up tasks may need to be defined."]
                    if not processed_data.get('decisions') or len(processed_data.get('decisions', [])) == 0:
                        logger.warning("No decisions extracted")
                        processed_data['decisions'] = ["No formal decisions recorded during the meeting."]
                    
                except json.JSONDecodeError as e:
                    logger.error(f"JSON decode error: {e}")
                    logger.error(f"Raw response: {ai_response}")
                    processed_data = None
                except Exception as e:
                    logger.error(f"AI processing failed: {e}")
                    processed_data = None
            
            # Enhanced fallback with real content extraction
            if not processed_data:
                logger.info("Using enhanced fallback key points extraction")
                
                important_keywords = [
                    'decision', 'decided', 'agree', 'approved', 'resolved',
                    'action', 'task', 'follow up', 'next step', 'deadline',
                    'issue', 'problem', 'challenge', 'concern', 'risk',
                    'project', 'initiative', 'proposal', 'plan', 'strategy',
                    'update', 'status', 'progress', 'result', 'outcome',
                    'budget', 'cost', 'resource', 'timeline', 'schedule'
                ]
                
                sentences = re.split(r'[.!?]+', optimized_text)
                important_sentences = []
                
                for sentence in sentences:
                    sentence = sentence.strip()
                    if len(sentence) > 20:
                        sentence_lower = sentence.lower()
                        if any(keyword in sentence_lower for keyword in important_keywords):
                            important_sentences.append(sentence)
                
                key_points_from_content = important_sentences[:30] if important_sentences else []
                
                if not key_points_from_content and topics:
                    key_points_from_content = [f"Discussion about {topic}" for topic in topics[:20]]
                
                processed_data = {
                    "summary": generate_comprehensive_summary(optimized_text, meeting.title, meaningful_sentences, topics, key_phrases),
                    "key_points": key_points_from_content,
                    "action_items": ["Review and distribute meeting notes to all participants", "Schedule follow-up meetings as discussed"],
                    "decisions": ["Meeting outcomes documented and approved by participants"],
                    "sentiment": "Professional meeting with productive discussions"
                }
            
            # Store transcript in ORIGINAL language
            processed_data["raw"] = raw_text
            processed_data["translated"] = translated_text
            
            progress_thread.join(timeout=25)
            update_processing_step(meeting, "ai_generation", "success")
            
            # Save to database - transcript stays in original language, notes in chosen language
            meeting.transcription = json.dumps({
                "raw": raw_text,
                "translated": translated_text,
                "optimized": optimized_text,
                "utterances": utterances
            })
            
            meeting.notes = json.dumps(processed_data)  # Chosen language summary/notes
            meeting.has_transcription = True
            meeting.has_notes = True
            meeting.status = 'completed'
            db.session.commit()
            
            logger.info(f"Processing completed successfully for meeting {meeting_id}")
            logger.info(f"Transcript language: {transcript_lang}, Summary language: {target_lang_name}")
            
            # Log completion activity
            log_activity(
                user_id=meeting.user_id,
                activity_type="completed",
                title=f"Processed: {meeting.title}",
                description="Meeting analysis completed successfully",
                meeting_id=meeting.id,
                metadata={"filename": meeting.filename, "duration": len(min(optimized_text, "1") * 1)} # simplified
            )
            
        except Exception as e:
            logger.error(f"Processing error for meeting {meeting_id}: {e}")
            try:
                steps = json.loads(meeting.processing_steps or '[]')
                for step in steps:
                    if step["status"] == "in_progress":
                        update_processing_step(meeting, step["step"], "failed", str(e))
                        break
            except:
                pass
            meeting.status = 'failed'
            db.session.commit()
            
            # Log failure activity
            log_activity(
                user_id=meeting.user_id,
                activity_type="failed",
                title=f"Failed: {meeting.title}",
                description="Processing failed due to an error",
                meeting_id=meeting.id,
                metadata={"error": str(e)}
            )

# Routes
@app.route('/', methods=['GET'])
def health():
    return jsonify({"status": "Backend running!", "timestamp": datetime.now(timezone.utc).isoformat()}), 200

@app.route("/api/auth/register", methods=["POST"])
def register():
    data = request.json
    if not data or not all(k in data for k in ("full_name", "email", "password", "phone_number")):
        return jsonify({"error": "Missing required fields"}), 400

    if not validate_email(data["email"]):
        return jsonify({"error": "Invalid email format"}), 400

    phone_number = data["phone_number"]
    if not phone_number.startswith("+"):
        return jsonify({"error": "Phone number must start with '+' (e.g., +92...)"}), 400

    if User.query.filter_by(email=data["email"]).first():
        return jsonify({"error": "Email already registered"}), 400
    
    if User.query.filter_by(phone_number=phone_number).first():
        return jsonify({"error": "Phone number already registered"}), 400

    hashed = generate_password_hash(data["password"])

    new_user = User(
        full_name=data["full_name"],
        email=data["email"],
        phone_number=phone_number,
        password_hash=hashed,
        email_verified=False,
        phone_verified=False,
        is_verified=False,
        password_updated_at=datetime.now(timezone.utc)
    )
    db.session.add(new_user)
    db.session.commit()

    # Create self-friendship by default
    self_friendship = Friendship(user_id=new_user.id, friend_id=new_user.id)
    db.session.add(self_friendship)
    db.session.commit()

    # Step 1: Generate and send Email OTP
    email_otp = generate_email_otp(new_user.id, new_user.email)
    send_verification_email(new_user.email, new_user.full_name, email_otp)
    
    # Step 2: Generate Phone OTP and display in terminal
    generate_phone_otp(new_user.id, new_user.phone_number)
    logger.info(f"User registered: {new_user.email}. Email and Phone OTPs displayed in terminal above.")

    return jsonify({
        "message": "Registration successful. Please check your email and terminal for verification codes.",
        "email": new_user.email
    }), 201

@app.route("/api/auth/verify-email", methods=["POST"])
def verify_email():
    """Verify email using OTP code"""
    data = request.json
    email = data.get("email")
    otp = data.get("otp")
    
    logger.info(f"[Email OTP Verify] Received verification request for email: {email}")
    
    if not email or not otp:
        return jsonify({"error": "Missing email or OTP"}), 400

    user = User.query.filter_by(email=email).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    # Verify OTP from database
    is_valid = verify_email_otp(user.id, otp)
    
    if not is_valid:
        logger.error(f"[Email OTP Verify] Verification failed for {email}")
        return jsonify({"error": "Invalid or expired OTP"}), 401
    
    logger.info(f"[Email OTP Verify] OTP verified successfully for: {email}")
    
    # Mark email as verified
    user.email_verified = True
    
    # Check if fully verified
    if user.phone_verified:
        user.is_verified = True
        logger.info(f"[Email OTP Verify] User {email} is now fully verified (email + phone)")
    else:
        logger.info(f"[Email OTP Verify] User {email} email verified, phone verification pending")
    
    # Clear OTP from database
    user.email_otp = None
    user.email_otp_expires = None
    db.session.commit()
    
    logger.info(f"[Email OTP Verify] âœ“ Email verification successful for {email}")
    
    return jsonify({
        "message": "Email verified successfully",
        "email": user.email,
        "phone_verified": user.phone_verified,
        "fully_verified": user.is_verified
    }), 200

@app.route("/api/auth/verify-phone", methods=["POST"])
def verify_phone():
    """Verify phone using terminal OTP"""
    data = request.json
    email = data.get("email")
    otp = data.get("otp")
    
    logger.info(f"[OTP Verify] Received verification request for email: {email}")
    
    if not email or not otp:
        return jsonify({"error": "Missing email or OTP"}), 400

    user = User.query.filter_by(email=email).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    # Verify OTP from database
    is_valid = verify_phone_otp(user.id, otp)
    
    if not is_valid:
        logger.error(f"[OTP Verify] Verification failed for {email}")
        return jsonify({"error": "Invalid or expired OTP"}), 401
    
    logger.info(f"[OTP Verify] OTP verified successfully for: {email}")
    
    # Mark phone as verified
    user.phone_verified = True
    if user.email_verified:
        user.is_verified = True
    
    # Clear OTP from database
    user.phone_otp = None
    user.phone_otp_expires = None
    db.session.commit()
    
    # If both verified, return access token
    if user.is_verified:
        token = create_access_token(identity=str(user.id))
        logger.info(f"[OTP Verify] âœ“ User {email} fully verified (email + phone)")
        return jsonify({
            "access_token": token,
            "user": {
                "id": user.id, 
                "full_name": user.full_name, 
                "email": user.email
            },
            "email_verified": True,
            "message": "Phone and email verified. Account activated."
        }), 200
    else:
        logger.info(f"[OTP Verify] Phone verified for {email}, email verification pending")
        return jsonify({
            "message": "Phone verified successfully. Please also verify your email link.",
            "email_verified": user.email_verified
        }), 200

@app.route("/api/auth/verify-otp", methods=["POST"])
def verify_otp():
    data = request.json
    if not data or not all(k in data for k in ("email", "email_otp", "phone_otp")):
        return jsonify({"error": "Missing required fields"}), 400

    user = User.query.filter_by(email=data["email"]).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    # Verify both OTPs using helper functions
    email_is_valid = verify_email_otp(user.id, data["email_otp"])
    phone_is_valid = verify_phone_otp(user.id, data["phone_otp"])
    
    if email_is_valid and phone_is_valid:
        logger.info(f"[Consolidated Verify] Both OTPs verified successfully for: {user.email}")
        
        user.email_verified = True
        user.phone_verified = True
        user.is_verified = True
        
        # Clear all OTP fields
        user.email_otp = None
        user.email_otp_expires = None
        user.phone_otp = None
        user.phone_otp_expires = None
        user.otp_expiry = None # Clear old field if used
        
        db.session.commit()

        token = create_access_token(identity=str(user.id))
        return jsonify({
            "access_token": token,
            "user": {
                "id": user.id, 
                "full_name": user.full_name, 
                "email": user.email, 
                "phone_number": user.phone_number
            },
            "message": "Verification successful"
        }), 200
    else:
        logger.error(f"[Consolidated Verify] Verification failed for {user.email}. Email valid: {email_is_valid}, Phone valid: {phone_is_valid}")
        return jsonify({"error": "Invalid or expired OTP codes"}), 401

@app.route("/api/auth/resend-verification", methods=["POST"])
def resend_verification():
    data = request.json
    email = data.get("email")
    method = data.get("method") # "email", "phone", or "both"
    
    if not email:
        return jsonify({"error": "Email is required"}), 400

    user = User.query.filter_by(email=email).first()
    if not user:
        return jsonify({"error": "User not found"}), 404

    if method == "email" or method == "both" or not method:
        # Generate new email OTP
        email_otp = generate_email_otp(user.id, user.email)
        send_verification_email(user.email, user.full_name, email_otp)
        logger.info(f"Email OTP regenerated for {user.email} - displayed in terminal above")

    if method == "phone" or method == "both" or not method:
        # Regenerate phone OTP and display in terminal
        generate_phone_otp(user.id, user.phone_number)
        logger.info(f"Phone OTP regenerated for {user.email} - displayed in terminal above")

    return jsonify({"message": f"Verification {method or 'both'} resent successfully"}), 200

@app.route("/api/auth/resend-otp", methods=["POST"])
def resend_otp():
    # Keep for backward compatibility with frontend
    return resend_verification()

@app.route("/api/auth/login", methods=["POST"])
def login():
    data = request.json
    if not data or not all(k in data for k in ("email", "password")):
        return jsonify({"error": "Missing required fields"}), 400

    user = User.query.filter_by(email=data["email"]).first()
    if not user:
        return jsonify({"error": "Email not registered"}), 404

    if not check_password_hash(user.password_hash, data["password"]):
        return jsonify({"error": "Incorrect password"}), 401
    
    if not user.is_verified:
        return jsonify({"error": "Account not verified", "code": "not_verified"}), 403

    # Check for 2FA PIN
    if user.two_factor_enabled:
        two_factor_pin = data.get("two_factor_pin")
        if not two_factor_pin:
            return jsonify({
                "status": "2fa_required",
                "message": "Two-factor authentication required",
                "email": user.email
            }), 200
        
        if not check_password_hash(user.two_factor_pin_hash, two_factor_pin):
            return jsonify({"error": "Incorrect 2FA PIN"}), 401

    # Reactivate account if it was deactivated (ONLY within 1 month)
    if not user.is_active:
        if user.deactivated_at:
            window_expiry = user.deactivated_at + timedelta(days=30)
            if datetime.now(timezone.utc) > window_expiry:
                return jsonify({
                    "error": "Account deactivation permanent",
                    "code": "account_expired",
                    "message": "The 30-day reactivation period for this account has expired. This account is now permanently deactivated."
                }), 403
            
            # Within 30 days, so reactivate
            user.is_active = True
            user.deactivated_at = None
            db.session.commit()
            logger.info(f"Reactivated account for user: {user.email} within 30-day window")
        else:
            # Fallback if deactivated_at is missing
            user.is_active = True
            db.session.commit()

    token = create_access_token(identity=str(user.id))
    return jsonify({
        "access_token": token,
        "user": {
            "id": user.id, 
            "full_name": user.full_name, 
            "email": user.email,
            "phone_number": user.phone_number,
            "two_factor_enabled": user.two_factor_enabled,
            "has_2fa_pin": bool(user.two_factor_pin_hash)
        }
    }), 200

@app.route("/api/auth/refresh", methods=["POST"])
@jwt_required()
def refresh():
    try:
        current_user_id = get_jwt_identity()
        user = db.session.get(User, current_user_id)
        if not user:
            return jsonify({"error": "User not found"}), 404
        
        new_token = create_access_token(
            identity=str(current_user_id),
            expires_delta=timedelta(days=30)
        )
        
        return jsonify({
            "access_token": new_token,
            "user": {
                "id": user.id, 
                "full_name": user.full_name, 
                "email": user.email,
                "phone_number": user.phone_number,
                "two_factor_enabled": user.two_factor_enabled,
                "has_2fa_pin": bool(user.two_factor_pin_hash)
            },
            "expires_in": 30 * 24 * 60 * 60
        })
    except Exception as e:
        logger.error(f"Token refresh failed: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/auth/forgot-password", methods=["POST"])
def forgot_password():
    try:
        data = request.json
        email = data.get("email")
        if not email:
            return jsonify({"error": "Email is required"}), 400
            
        user = User.query.filter_by(email=email).first()
        if not user:
            # For security, don't reveal if user exists
            # simulate delay to prevent timing attacks
            time.sleep(random.uniform(0.1, 0.5))
            return jsonify({"message": "If your email is registered, you will receive a password reset link."}), 200
            
        # Generate token
        token = secrets.token_urlsafe(32)
        user.reset_token = token
        # Token valid for 1 hour
        from datetime import timezone
        user.reset_token_expiry = datetime.now(timezone.utc).replace(tzinfo=None) + timedelta(hours=1)
        db.session.commit()
        
        # Send Email via SendGrid
        sendgrid_api_key = os.getenv('SENDGRID_API_KEY')
        sendgrid_from_email = os.getenv('SENDGRID_FROM_EMAIL') or 'noreply@talktotextpro.com'
        
        if sendgrid_api_key:
            # Construct reset link
            # Assuming frontend is on same domain or specified via env
            frontend_url = os.getenv('FRONTEND_URL', 'http://localhost:3000')
            reset_link = f"{frontend_url}/reset-password?token={token}"
            
            email_content = f"""
<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Reset Your Password</title>
  <style>
    body {{
      margin: 0;
      padding: 0;
      font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
      background-color: #111827;
      color: #ffffff;
    }}
    .email-wrapper {{
      padding: 40px 20px;
      background: radial-gradient(ellipse at top left, rgba(56, 189, 248, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at top right, rgba(168, 85, 247, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at bottom, rgba(16, 185, 129, 0.12) 0%, transparent 50%),
                  #111827;
    }}
    .container {{
      max-width: 600px;
      margin: 0 auto;
      background-color: rgba(31, 41, 55, 0.7);
      border: 1px solid rgba(255, 255, 255, 0.1);
      border-radius: 24px;
      overflow: hidden;
      backdrop-filter: blur(10px);
    }}
    .header {{
      padding: 40px 40px 20px;
      text-align: center;
    }}
    .logo {{
      font-size: 24px;
      font-weight: 800;
      letter-spacing: -0.025em;
      background: linear-gradient(to right, #60a5fa, #a855f7);
      -webkit-background-clip: text;
      -webkit-text-fill-color: transparent;
      margin-bottom: 24px;
    }}
    .content {{
      padding: 0 40px 40px;
      text-align: center;
    }}
    h1 {{
      font-size: 28px;
      font-weight: 700;
      margin-bottom: 16px;
      color: #ffffff;
    }}
    p {{
      font-size: 16px;
      line-height: 1.6;
      color: #9ca3af;
      margin-bottom: 32px;
    }}
    .button-container {{
      margin-bottom: 32px;
    }}
    .button {{
      display: inline-block;
      padding: 16px 32px;
      background: linear-gradient(to right, #ef4444, #f87171);
      color: #ffffff !important;
      text-decoration: none;
      border-radius: 12px;
      font-weight: 600;
      font-size: 16px;
      box-shadow: 0 10px 15px -3px rgba(239, 68, 68, 0.3);
    }}
    .link-fallback {{
      font-size: 12px;
      color: #6b7280;
      word-break: break-all;
      margin-top: 24px;
    }}
    .footer {{
      padding: 24px 40px;
      background-color: rgba(17, 24, 39, 0.5);
      border-top: 1px solid rgba(255, 255, 255, 0.05);
      text-align: center;
      font-size: 14px;
      color: #6b7280;
    }}
  </style>
</head>
<body>
  <div class="email-wrapper">
    <div class="container">
      <div class="header">
        <div class="logo">TalkToText Pro</div>
      </div>
      <div class="content">
        <h1>Password Reset Request</h1>
        <p>Hello {user.full_name},</p>
        <p>We received a request to reset your password. If you didn't initiate this, you can safely ignore this email.</p>
        <div class="button-container">
          <a href="{reset_link}" class="button">Reset My Password</a>
        </div>
        <p>This link will expire in 1 hour for your security.</p>
        <div class="link-fallback">
          Having trouble? Copy this link: {reset_link}
        </div>
      </div>
      <div class="footer">
        &copy; 2026 TalkToText Pro. Your security is our priority.
      </div>
    </div>
  </div>
</body>
</html>
    """
            
            message = Mail(
                from_email=sendgrid_from_email,
                to_emails=email,
                subject="Reset Your Password - TalkToText Pro",
                html_content=email_content
            )
            
            try:
                sg = SendGridAPIClient(api_key=sendgrid_api_key)
                sg.send(message)
                logger.info(f"Password reset email sent to {email}")
            except Exception as e:
                logger.error(f"Failed to send reset email: {e}")
                # FALLBACK: Log the reset link so the user can see it in terminal
                frontend_url = os.getenv('FRONTEND_URL', 'http://localhost:3000')
                reset_link = f"{frontend_url}/reset-password?token={token}"
                logger.warning(f"CRITICAL: Password reset email delivery failed. Manual Reset Link for {email}: {reset_link}")
                # Return success anyway so user checks logs
                return jsonify({"message": "If your email is registered, you will receive a password reset link."}), 200
        else:
            logger.warning("SendGrid not configured, printing reset token to logs")
            logger.info(f"RESET TOKEN for {email}: {token}")
            
        return jsonify({"message": "If your email is registered, you will receive a password reset link."}), 200
        
    except Exception as e:
        logger.error(f"Forgot password error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/auth/reset-password", methods=["POST"])
def reset_password():
    try:
        data = request.json
        token = data.get("token")
        new_password = data.get("new_password")
        
        if not token or not new_password:
            return jsonify({"error": "Missing token or password"}), 400
            
        user = User.query.filter_by(reset_token=token).first()
        
        if not user:
            return jsonify({"error": "Invalid or expired token"}), 400
            
        if user.reset_token_expiry < datetime.now(timezone.utc).replace(tzinfo=None):
            return jsonify({"error": "Token has expired"}), 400
            
        # Update password
        user.password_hash = generate_password_hash(new_password)
        user.reset_token = None
        user.reset_token_expiry = None
        db.session.commit()
        
        return jsonify({"message": "Password reset successfully. You can now login."}), 200
        
    except Exception as e:
        logger.error(f"Reset password error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/auth/forgot-pin", methods=["POST"])
def forgot_pin():
    try:
        data = request.json
        email = data.get("email")
        if not email:
            return jsonify({"error": "Email is required"}), 400
            
        user = User.query.filter_by(email=email).first()
        if not user:
            # For security, don't reveal if user exists
            time.sleep(random.uniform(0.1, 0.5))
            return jsonify({"message": "If your account is configured for 2FA, you will receive a PIN reset link."}), 200
            
        # Generate token
        token = secrets.token_urlsafe(32)
        user.reset_token = token  # Reuse reset_token for simplicity as per common pattern
        user.reset_token_expiry = datetime.now(timezone.utc).replace(tzinfo=None) + timedelta(hours=1)
        db.session.commit()
        
        # Send Email via SendGrid
        sendgrid_api_key = os.getenv('SENDGRID_API_KEY')
        sendgrid_from_email = os.getenv('SENDGRID_FROM_EMAIL') or 'noreply@talktotextpro.com'
        
        if sendgrid_api_key:
            frontend_url = os.getenv('FRONTEND_URL', 'http://localhost:3000')
            reset_link = f"{frontend_url}/reset-pin?token={token}"
            
            email_content = f"""
<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Reset Your Security PIN</title>
  <style>
    body {{
      margin: 0;
      padding: 0;
      font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
      background-color: #111827;
      color: #ffffff;
    }}
    .email-wrapper {{
      padding: 40px 20px;
      background: radial-gradient(ellipse at top left, rgba(56, 189, 248, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at top right, rgba(168, 85, 247, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at bottom, rgba(16, 185, 129, 0.12) 0%, transparent 50%),
                  #111827;
    }}
    .container {{
      max-width: 600px;
      margin: 0 auto;
      background-color: rgba(31, 41, 55, 0.7);
      border: 1px solid rgba(255, 255, 255, 0.1);
      border-radius: 24px;
      overflow: hidden;
      backdrop-filter: blur(10px);
    }}
    .header {{
      padding: 40px 40px 20px;
      text-align: center;
    }}
    .logo {{
      font-size: 24px;
      font-weight: 800;
      letter-spacing: -0.025em;
      background: linear-gradient(to right, #60a5fa, #a855f7);
      -webkit-background-clip: text;
      -webkit-text-fill-color: transparent;
      margin-bottom: 24px;
    }}
    .content {{
      padding: 0 40px 40px;
      text-align: center;
    }}
    h1 {{
      font-size: 28px;
      font-weight: 700;
      margin-bottom: 16px;
      color: #ffffff;
    }}
    p {{
      font-size: 16px;
      line-height: 1.6;
      color: #9ca3af;
      margin-bottom: 32px;
    }}
    .button-container {{
      margin-bottom: 32px;
    }}
    .button {{
      display: inline-block;
      padding: 16px 32px;
      background: linear-gradient(to right, #3b82f6, #8b5cf6);
      color: #ffffff !important;
      text-decoration: none;
      border-radius: 12px;
      font-weight: 600;
      font-size: 16px;
      box-shadow: 0 10px 15px -3px rgba(59, 130, 246, 0.3);
    }}
    .link-fallback {{
      font-size: 12px;
      color: #6b7280;
      word-break: break-all;
      margin-top: 24px;
    }}
    .footer {{
      padding: 24px 40px;
      background-color: rgba(17, 24, 39, 0.5);
      border-top: 1px solid rgba(255, 255, 255, 0.05);
      text-align: center;
      font-size: 14px;
      color: #6b7280;
    }}
  </style>
</head>
<body>
  <div class="email-wrapper">
    <div class="container">
      <div class="header">
        <div class="logo">TalkToText Pro</div>
      </div>
      <div class="content">
        <h1>PIN Reset Request</h1>
        <p>Hello {user.full_name},</p>
        <p>We received a request to reset your security PIN. This PIN is required for two-factor authentication.</p>
        <div class="button-container">
          <a href="{reset_link}" class="button">Reset Security PIN</a>
        </div>
        <p>This link will expire in 1 hour for your security.</p>
        <div class="link-fallback">
          Having trouble? Copy this link: {reset_link}
        </div>
      </div>
      <div class="footer">
        &copy; 2026 TalkToText Pro. Your security is our priority.
      </div>
    </div>
  </div>
</body>
</html>
    """
            
            message = Mail(
                from_email=sendgrid_from_email,
                to_emails=email,
                subject="Reset Your Security PIN - TalkToText Pro",
                html_content=email_content
            )
            
            try:
                sg = SendGridAPIClient(api_key=sendgrid_api_key)
                sg.send(message)
                logger.info(f"PIN reset email sent to {email}")
            except Exception as e:
                logger.error(f"Failed to send PIN reset email: {e}")
                return jsonify({"error": "Failed to send reset email"}), 500
        else:
            logger.warning("SendGrid not configured, printing PIN reset token to logs")
            logger.info(f"PIN RESET TOKEN for {email}: {token}")
            
        return jsonify({"message": "If your email is registered, you will receive a PIN reset link."}), 200
        
    except Exception as e:
        logger.error(f"Forgot PIN error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/auth/reset-pin", methods=["POST"])
def reset_pin():
    try:
        data = request.json
        token = data.get("token")
        new_pin = data.get("new_pin")
        
        logger.info(f"Received reset-pin request. Token: '{token}', PIN length: {len(new_pin) if new_pin else 'None'}")
        
        if not token or not new_pin:
            logger.warning(f"Missing data: token is {bool(token)}, PIN is {bool(new_pin)}")
            return jsonify({"error": "Missing token or PIN"}), 400
            
        if len(new_pin) != 6 or not new_pin.isdigit():
            logger.warning(f"Invalid PIN format for reset: {new_pin}")
            return jsonify({"error": "PIN must be 6 digits"}), 400
            
        # Diagnostic: List all active reset tokens in the DB
        all_users_with_tokens = User.query.filter(User.reset_token != None).all()
        logger.info(f"Found {len(all_users_with_tokens)} users with active reset tokens.")
        for u in all_users_with_tokens:
            logger.info(f"User: {u.email}, Token in DB: '{u.reset_token}', Match: {u.reset_token == token}")
            
        user = User.query.filter_by(reset_token=token).first()
        
        if not user:
            logger.warning(f"No user found with reset token: '{token}'")
            return jsonify({"error": "Invalid or expired token"}), 400
            
        from datetime import timezone # Safer local import again just in case
        now = datetime.now(timezone.utc).replace(tzinfo=None)
        if user.reset_token_expiry < now:
            logger.warning(f"Token expired for user {user.email}. Expiry: {user.reset_token_expiry}, Now: {now}")
            return jsonify({"error": "Token has expired"}), 400
            
        # Update PIN
        user.two_factor_pin_hash = generate_password_hash(new_pin)
        user.reset_token = None
        user.reset_token_expiry = None
        # Automatically enable 2FA if they reset it via email
        user.two_factor_enabled = True
        db.session.commit()
        
        logger.info(f"Successfully reset PIN for user {user.email}")
        return jsonify({"message": "Security PIN reset successfully. You can now login."}), 200
        
    except Exception as e:
        logger.error(f"Reset PIN error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/auth/verify-credentials", methods=["POST"])
def verify_credentials():
    try:
        data = request.json
        email = data.get("email")
        password = data.get("password")
        phone_number = data.get("phone_number")
        
        response = {
            "email_exists": False,
            "password_correct": False,
            "phone_exists": False,
            "is_verified": False,
            "phone_number": None
        }
        
        if email:
            user = User.query.filter_by(email=email).first()
            if user:
                response["email_exists"] = True
                response["is_verified"] = user.is_verified
                response["email_verified"] = user.email_verified
                response["phone_verified"] = user.phone_verified
                response["phone_number"] = user.phone_number
                if password:
                    # Check password only if user exists
                    if check_password_hash(user.password_hash, password):
                        response["password_correct"] = True
        
        if phone_number:
            user_phone = User.query.filter_by(phone_number=phone_number).first()
            if user_phone:
                response["phone_exists"] = True
                response["is_verified"] = user_phone.is_verified
                response["email_verified"] = user_phone.email_verified
                response["phone_verified"] = user_phone.phone_verified
                if not response["phone_number"]:
                    response["phone_number"] = user_phone.phone_number
        
        return jsonify(response), 200
        
    except Exception as e:
        logger.error(f"Credential verification error: {e}")
        return jsonify({"error": "Verification failed"}), 500

@app.route("/api/auth/profile", methods=["GET"])
@jwt_required()
def get_profile():
    try:
        user_id = int(get_jwt_identity())
        user = db.session.get(User, user_id)
        if not user:
            return jsonify({"error": "User not found"}), 404
            
        return jsonify({
            "id": user.id,
            "full_name": user.full_name,
            "email": user.email,
            "phone_number": user.phone_number,
            "bio": user.bio,
            "password_updated_at": user.password_updated_at.isoformat() if user.password_updated_at else None,
            "two_factor_enabled": user.two_factor_enabled,
            "has_2fa_pin": bool(user.two_factor_pin_hash)
        }), 200
    except Exception as e:
        logger.error(f"Get profile error: {e}")
        return jsonify({"error": "Failed to fetch profile"}), 500

@app.route("/api/auth/profile", methods=["PUT"])
@jwt_required()
def update_profile():
    try:
        user_id = int(get_jwt_identity())
        user = db.session.get(User, user_id)
        if not user:
            return jsonify({"error": "User not found"}), 404
            
        data = request.json
        if "full_name" in data:
            user.full_name = data["full_name"]
        if "bio" in data:
            user.bio = data["bio"]
        if "image" in data:
            user.image = data["image"]
        if "two_factor_enabled" in data:
            # Only allow enabling if a PIN is set
            if data["two_factor_enabled"] and not user.two_factor_pin_hash:
                return jsonify({"error": "Please set a 2FA PIN first"}), 400
            user.two_factor_enabled = data["two_factor_enabled"]
            
        db.session.commit()
        return jsonify({
            "message": "Profile updated successfully",
            "user": {
                "id": user.id, 
                "full_name": user.full_name, 
                "email": user.email,
                "phone_number": user.phone_number,
                "bio": user.bio,
                "two_factor_enabled": user.two_factor_enabled
            }
        }), 200
    except Exception as e:
        logger.error(f"Update profile error: {e}")
        return jsonify({"error": "Failed to update profile"}), 500

@app.route("/api/auth/profile/password", methods=["PUT"])
@jwt_required()
def change_password():
    try:
        user_id = get_current_user_id()
        user = db.session.get(User, user_id)
        if not user:
            return jsonify({"error": "User not found"}), 404
            
        data = request.json
        current_password = data.get("current_password")
        new_password = data.get("new_password")
        
        if not current_password or not new_password:
            return jsonify({"error": "Missing current or new password"}), 400
            
        if not check_password_hash(user.password_hash, current_password):
            return jsonify({"error": "Incorrect current password"}), 401
            
        # Add check to prevent using the same password
        if check_password_hash(user.password_hash, new_password):
            return jsonify({"error": "password old", "message": "New password cannot be the same as the current password"}), 400
            
        user.password_hash = generate_password_hash(new_password)
        user.password_updated_at = datetime.now(datetime.UTC)
        db.session.commit()
        
        return jsonify({"message": "Password changed successfully", "password_updated_at": user.password_updated_at.isoformat()}), 200
    except Exception as e:
        logger.error(f"Change password error: {e}")
        return jsonify({"error": "Failed to change password"}), 500

@app.route("/api/auth/profile/2fa-pin", methods=["PUT"])
@jwt_required()
def update_2fa_pin():
    try:
        user_id = int(get_jwt_identity())
        logger.info(f"Update 2FA PIN - User ID: {user_id}")
        user = db.session.get(User, user_id)
        if not user:
            logger.error(f"User not found for ID: {user_id}")
            return jsonify({"error": "User not found"}), 404
            
        data = request.json
        current_pin = data.get("current_pin")
        new_pin = data.get("pin")
        
        # If user already has a PIN, REQUIRE the current one to change it
        if user.two_factor_pin_hash:
            if not current_pin:
                return jsonify({"error": "Current PIN is required to change your PIN"}), 400
            if not check_password_hash(user.two_factor_pin_hash, current_pin):
                return jsonify({"error": "Incorrect current PIN"}), 401
                
        if not new_pin or len(new_pin) != 6 or not new_pin.isdigit():
            return jsonify({"error": "New PIN must be 6 digits"}), 400
            
        user.two_factor_pin_hash = generate_password_hash(new_pin)
        db.session.commit()
        
        return jsonify({"message": "2FA PIN updated successfully"}), 200
    except Exception as e:
        logger.error(f"Update 2FA PIN error: {e}")
        return jsonify({"error": "Failed to update 2FA PIN"}), 500

@app.route("/api/auth/profile/verify-2fa-pin", methods=["POST"])
@jwt_required()
def verify_2fa_pin():
    try:
        user_id = int(get_jwt_identity())
        logger.info(f"Verify 2FA PIN - User ID: {user_id}")
        user = db.session.get(User, user_id)
        if not user:
            logger.error(f"User not found for ID: {user_id}")
            return jsonify({"error": "User not found"}), 404
            
        data = request.json
        pin = data.get("pin")
        
        if not pin:
            return jsonify({"error": "PIN is required"}), 400
        if check_password_hash(user.two_factor_pin_hash, pin):
            logger.info(f"PIN verified for user {user_id}")
            return jsonify({"success": True, "message": "PIN verified"}), 200
        else:
            logger.warning(f"PIN verification failed for user {user_id}")
            return jsonify({"success": False, "error": "Incorrect PIN"}), 401
            
    except Exception as e:
        logger.error(f"Verify 2FA PIN error: {e}")
        return jsonify({"error": "Failed to verify 2FA PIN"}), 500

@app.route("/api/auth/profile", methods=["DELETE"])
@jwt_required()
def delete_account():
    try:
        user_id = get_current_user_id()
        user = db.session.get(User, user_id)
        if not user:
            return jsonify({"error": "User not found"}), 404
            
        # Soft delete: mark as inactive and set deactivation timestamp
        user.is_active = False
        user.deactivated_at = datetime.now(timezone.utc)
        db.session.commit()
        
        return jsonify({"message": "Account deactivated successfully"}), 200
    except Exception as e:
        logger.error(f"Delete account error: {e}")
        return jsonify({"error": "Failed to delete account"}), 500

@app.route("/api/auth/check", methods=["GET"])
@jwt_required()
def check_auth():
    try:
        user_id = get_current_user_id()
        user = db.session.get(User, user_id)
        if not user:
            logger.error(f"User {user_id} not found in database")
            return jsonify({"error": "User not found"}), 404
        
        logger.info(f"Token validation successful for user: {user.email}")
        return jsonify({
            "valid": True,
            "user_id": user_id,
            "user_email": user.email,
            "user_name": user.full_name
        }), 200
        
    except Exception as e:
        logger.error(f"Auth check failed: {e}")
        return jsonify({"error": "Invalid token"}), 401

@app.route("/api/auth/validate", methods=["GET"])
@jwt_required()
def validate_token():
    try:
        current_user_id = get_jwt_identity()
        user = db.session.get(User, current_user_id)
        if not user:
            return jsonify({"error": "User not found"}), 404
        
        return jsonify({
            "valid": True,
            "user": {"id": user.id, "full_name": user.full_name, "email": user.email}
        })
    except Exception as e:
        logger.error(f"Token validation failed: {e}")
        return jsonify({"error": str(e)}), 500

# ===============================================================
# ðŸŽ¥ 1ï¸âƒ£ DOWNLOAD VIDEO FROM URL
# ===============================================================
@app.route("/api/download_video", methods=["GET"])
def download_video():
    """
    Downloads video from:
    - YouTube, Facebook, Instagram
    - Google Drive, TikTok, Twitter/X, Twitch
    - Vimeo, Dailymotion, SoundCloud
    
    Blocks live meeting links with clear error.
    
    NOTE: EventSource can't send custom headers, so token is in query params
    """
    # Get token from query params (EventSource limitation)
    token = request.args.get("token")
    if not token:
        logger.error("Missing token in query parameters")
        return jsonify({"error": "Missing token"}), 401
    
    # Manually verify the JWT token
    try:
        from flask_jwt_extended import decode_token
        decoded = decode_token(token)
        user_id = decoded['sub']
    except Exception as e:
        logger.error(f"Token verification failed: {e}")
        return jsonify({"error": "Invalid token"}), 401
    
    url = request.args.get("url")

    if not url or not url.startswith("http"):
        logger.error(f"Invalid or missing URL: {url}")
        return jsonify({"error": "Invalid or missing URL"}), 400

    # === LIVE MEETING DETECTION ===
    LIVE_PATTERNS = [
        "meet.google.com/",
        "teams.microsoft.com/l/meetup-join",
        "zoom.us/j/",
        "zoom.us/w/",
    ]
    if any(pattern in url for pattern in LIVE_PATTERNS):
        return Response(
            f"data: {json.dumps({'status': 'error', 'error': 'Live meetings cannot be downloaded. Please wait for the recording to finish and upload the file manually.'})}\n\n",
            mimetype='text/event-stream'
        )

    def _choose_ydl_extractor(url: str) -> dict:
        """Force yt-dlp to use the correct extractor."""
        if "drive.google.com" in url:
            return {"extractor": "googledrive"}
        return {}

    def generate():
        progress_queue = Queue()

        def progress_hook(d):
            if d.get('status') == 'downloading':
                total = d.get('total_bytes') or d.get('total_bytes_estimate')
                if total and d.get('downloaded_bytes'):
                    progress = min(100, round((d['downloaded_bytes'] / total) * 100))
                    progress_queue.put({
                        "status": "downloading",
                        "downloaded_bytes": d['downloaded_bytes'],
                        "total_bytes": total,
                        "progress": progress
                    })
            elif d.get('status') == 'finished':
                filename = secure_filename(os.path.basename(d.get('filename', 'video.%(ext)s')))
                progress_queue.put({
                    "status": "finished",
                    "filename": filename
                })

        def download_thread():
            try:
                # === DYNAMIC REFERER ===
                referer_url = "https://www.youtube.com/"
                if "youtube.com" in url or "youtu.be" in url:
                    referer_url = "https://www.youtube.com/"
                elif "facebook.com" in url:
                    referer_url = "https://www.facebook.com/"
                elif "instagram.com" in url:
                    referer_url = "https://www.instagram.com/"
                elif "drive.google.com" in url:
                    referer_url = "https://drive.google.com/"
                elif "twitter.com" in url or "x.com" in url:
                    referer_url = "https://twitter.com/"
                elif "tiktok.com" in url:
                    referer_url = "https://www.tiktok.com/"
                elif "vimeo.com" in url:
                    referer_url = "https://vimeo.com/"
                elif "dailymotion.com" in url or "dai.ly" in url:
                    referer_url = "https://www.dailymotion.com/"
                elif "soundcloud.com" in url:
                    referer_url = "https://soundcloud.com/"
                elif "twitch.tv" in url:
                    referer_url = "https://www.twitch.tv/"

                # === YT-DLP OPTIONS ===
                ydl_opts = {
                    "format": "bestvideo[height<=1080][ext=mp4]+bestaudio[ext=m4a]/best[height<=1080][ext=mp4]/best",
                    "outtmpl": os.path.join(UPLOAD_DIR, "%(id)s.%(ext)s"),
                    "user_agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/130.0.0.0 Safari/537.36",
                    "referer": referer_url,
                    "http_headers": {
                        "Referer": referer_url,
                        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
                        "Accept-Language": "en-US,en;q=0.9",
                        "Sec-Fetch-Mode": "navigate",
                    },
                    "quiet": False,
                    "no_warnings": False,
                    "noplaylist": True,
                    "no_cache_dir": True,
                    "retries": 10,
                    "fragment_retries": 20,
                    "sleep_interval": 3,
                    "max_sleep_interval": 10,
                    "progress_hooks": [progress_hook],
                    "extractor_retries": 5,
                    "socket_timeout": 60,
                    "force_ipv4": True,
                    "merge_output_format": "mp4",
                    "nocheckcertificate": True,
                    "youtube_include_dash_manifest": True,
                    "youtube_include_hls_manifest": True,
                    "check_formats": True,
                    "extractor_args": {
                        "youtube": {
                            "player_client": ["android_vr", "mweb_safari", "ios", "web", "android"],
                            "player_skip": ["webpage", "configs"],
                            "skip": ["translated_subs"],
                        }
                    },
                    **_choose_ydl_extractor(url),
                }

                # === COOKIE HANDLING ===
                # Priority 1: Explicit cookies file from env
                cookies_file = os.getenv("YT_DLP_COOKIES_FILE") or os.path.join(os.path.dirname(__file__), "cookies.txt")
                if os.path.exists(cookies_file):
                    logger.info(f"Using explicit cookies file: {cookies_file}")
                    ydl_opts["cookiefile"] = cookies_file
                    
                    # Try with explicit cookies file
                    try:
                        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                            info = ydl.extract_info(url, download=True)
                            filename = ydl.prepare_filename(info)
                            if os.path.exists(filename):
                                progress_queue.put({"status": "finished", "filename": os.path.basename(filename)})
                                return # Success!
                    except Exception as cookie_err:
                        logger.warning(f"Explicit cookies file attempt failed: {cookie_err}")
                else:
                    logger.info("No cookies file provided. Proceeding with standard download.")
                # Priority 2: Try various browsers (only if NOT in a headless server environment)
                # Headless servers usually don't have browser profiles, so we skip browser_cookie3 failures silently
                cookie_browsers = ["chrome", "edge", "firefox", "brave"]
                for browser_name in cookie_browsers:
                    try:
                        temp_opts = ydl_opts.copy()
                        temp_opts["cookiesfrombrowser"] = (browser_name,)
                        logger.info(f"Attempting download with {browser_name} cookies...")
                        
                        with yt_dlp.YoutubeDL(temp_opts) as ydl:
                            info = ydl.extract_info(url, download=True)
                            filename = ydl.prepare_filename(info)
                            if os.path.exists(filename):
                                progress_queue.put({"status": "finished", "filename": os.path.basename(filename)})
                                return # Success!
                    except Exception as browser_err:
                        logger.debug(f"{browser_name} cookies attempt failed: {browser_err}")
                        continue
                
                # Priority 3: Final attempt without cookies
                try:
                    logger.info("Final attempt without cookies...")
                    final_opts = ydl_opts.copy()
                    final_opts.pop("cookiesfrombrowser", None)
                    final_opts.pop("cookiefile", None)
                    with yt_dlp.YoutubeDL(final_opts) as ydl:
                        info = ydl.extract_info(url, download=True)
                        filename = ydl.prepare_filename(info)
                        if os.path.exists(filename):
                            progress_queue.put({"status": "finished", "filename": os.path.basename(filename)})
                            return # Success!
                except Exception as final_err:
                    raise final_err

            except Exception as e:
                logger.error(f"Download error: {str(e)}")
                error_msg = str(e)

                # === USER-FRIENDLY ERROR MESSAGES ===
                if any(k in error_msg.lower() for k in ["private", "sign in", "login", "unauthorized", "403", "forbidden", "bot"]):
                    error_msg = (
                        "YouTube is blocking the server. To fix this: \n"
                        "1. Export your YouTube cookies using 'Get cookies.txt LOCALLY' extension. \n"
                        "2. Upload 'cookies.txt' to the backend folder or set YT_DLP_COOKIES_FILE env var to the file path. \n"
                        "Otherwise, download the video manually and upload the file."
                    )
                elif "unavailable" in error_msg.lower():
                    error_msg = "Video is unavailable or has been removed."
                elif "geo" in error_msg.lower():
                    error_msg = "This video is not available in your region."

                progress_queue.put({"status": "error", "error": error_msg})

        # Start download in background
        thread = threading.Thread(target=download_thread, daemon=True)
        thread.start()

        # Send starting event
        yield f"data: {json.dumps({'status': 'starting', 'progress': 0})}\n\n"

        # Stream progress with timeout
        while thread.is_alive() or not progress_queue.empty():
            try:
                data = progress_queue.get(timeout=1)
                yield f"data: {json.dumps(data)}\n\n"
                if data.get('status') in ['finished', 'error']:
                    break
            except Empty:
                yield f"data: {json.dumps({'status': 'downloading', 'message': 'Processing video...', 'progress': 0})}\n\n"

        # Final keep-alive
        yield f"data: {json.dumps({'status': 'done'})}\n\n"

    logger.info(f"Starting video download for URL: {url}")
    return Response(generate(), mimetype='text/event-stream')
# ===============================================================
# ðŸ“¤ 2ï¸âƒ£ UPLOAD / PROCESS FILE
# ===============================================================
@app.route("/api/upload", methods=["POST"])
@jwt_required()
def upload():
    """Handles upload of pre-downloaded or direct files, and retry processing."""
    try:
        user_id = get_jwt_identity()
        user = db.session.get(User, user_id)
        
        if not user:
            return jsonify({"error": "User not found"}), 404

        # Handle both JSON and form data
        data = request.get_json(silent=True) or request.form.to_dict()
        
        # Check if this is a retry request
        if data.get("retry") and data.get("meeting_id"):
            meeting_id = data.get("meeting_id")
            meeting = Meeting.query.filter_by(id=meeting_id, user_id=user_id).first()
            
            if not meeting:
                return jsonify({"error": "Meeting not found"}), 404
            
            # Reset meeting status for retry
            meeting.status = "uploaded"
            meeting.processing_steps = json.dumps([
                {"step": "transcription", "status": "pending", "timestamp": "", "error": None},
                {"step": "translation", "status": "pending", "timestamp": "", "error": None},
                {"step": "optimization", "status": "pending", "timestamp": "", "error": None},
                {"step": "ai_generation", "status": "pending", "timestamp": "", "error": None}
            ])
            meeting.current_step_progress = 0
            db.session.commit()
            
            logger.info(f"Retrying processing for meeting {meeting_id}")
            
            log_activity(
                user_id=user_id,
                activity_type="processing",
                title=f"Retry: {meeting.title}",
                description=f"Retrying processing for {meeting.filename}",
                meeting_id=meeting.id,
                metadata={"filename": meeting.filename, "title": meeting.title, "retry": True}
            )
            
            executor.submit(start_processing, meeting.id)
            
            return jsonify({
                "recording_id": meeting.id,
                "message": "Processing restarted successfully",
                "filename": meeting.filename,
                "status": "uploaded",
                "transcript_language": meeting.transcript_language,
            }), 200
        
        title = data.get("title")
        language = data.get("language", "en")
        transcript_language = data.get("transcript_language", "en")
        url = data.get("url")
        pre_downloaded_filename = data.get("filename")

        if not title:
            return jsonify({"error": "Missing title"}), 400

        logger.info(f"Upload request - Title: {title}, Transcript Language: {transcript_language}")

        filepath = None
        filename = None

        # === 1. Handle Pre-Downloaded File ===
        if pre_downloaded_filename:
            filepath = os.path.join(UPLOAD_DIR, secure_filename(pre_downloaded_filename))
            if not os.path.exists(filepath):
                logger.error(f"Pre-downloaded file not found: {filepath}")
                return jsonify({"error": "Pre-downloaded file not found"}), 404
            filename = pre_downloaded_filename
            logger.info(f"Using pre-downloaded file: {filepath}")

        # === 2. Handle URL Download (No ffmpeg merge) ===
        elif url and url.strip():
            try:
                # === ROBUST YT-DLP OPTIONS (NO FFMPEG MERGE) ===
                ydl_opts = {
                    "format": "bestvideo[height<=1080]+bestaudio/best[height<=1080]/best",
                    "outtmpl": f"{UPLOAD_DIR}/%(id)s.%(ext)s",
                    "user_agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/129.0.0.0 Safari/537.36",
                    "referer": "https://www.youtube.com/",
                    "http_headers": {
                        "Referer": "https://www.youtube.com/",
                        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
                        "Accept-Language": "en-US,en;q=0.5",
                        "Accept-Encoding": "gzip, deflate",
                        "DNT": "1",
                        "Connection": "keep-alive",
                        "Upgrade-Insecure-Requests": "1",
                    },
                    "quiet": False,
                    "noplaylist": True,
                    "no_warnings": False,
                    "no_cache_dir": True,
                    "retries": 5,
                    "fragment_retries": 10,
                    "sleep_interval": 1,
                    "max_sleep_interval": 3,
                    "socket_timeout": 30,
                    "force_ipv4": True,
                    "extractor_retries": 3,
                }

                with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                    info = ydl.extract_info(url, download=True)
                    filename = ydl.prepare_filename(info)
                    filepath = filename

                if not os.path.exists(filepath):
                    raise Exception("Downloaded file not found after yt-dlp")
                logger.info(f"Successfully downloaded from URL: {filename}")

            except Exception as e:
                error_msg = str(e)
                logger.error(f"Failed to download URL {url}: {error_msg}")
                return jsonify({"error": f"Failed to download from URL: {error_msg}"}), 500

        # === 3. Handle Direct File Upload ===
        elif "file" in request.files:
            file = request.files["file"]
            
            if file.filename == '':
                return jsonify({"error": "No file selected"}), 400

            allowed_ext = {".mp3", ".wav", ".mp4", ".mov", ".m4a", ".webm", ".ogg", ".flac"}
            ext = os.path.splitext(file.filename)[1].lower()
            if ext not in allowed_ext:
                return jsonify({"error": f"Unsupported file format: {ext}. Allowed: {', '.join(allowed_ext)}"}), 400

            filename = secure_filename(file.filename)
            filepath = os.path.join(UPLOAD_DIR, filename)
            
            try:
                file.save(filepath)
                logger.info(f"File saved successfully: {filepath}")
            except Exception as e:
                logger.error(f"Failed to save file: {e}")
                return jsonify({"error": f"Failed to save file: {str(e)}"}), 500

        else:
            return jsonify({"error": "No file, URL, or pre-downloaded file provided"}), 400

        # === 4. Create Meeting Record ===
        meeting_source = 'url' if url and url.strip() else 'upload'
        meeting = Meeting(
            user_id=user_id,
            title=title,
            filename=os.path.basename(filepath),
            language=language,
            transcript_language=transcript_language,
            status="uploaded",
            source=meeting_source
        )

        db.session.add(meeting)
        db.session.commit()

        logger.info(f"Meeting created with ID {meeting.id}")

        # === 5. Log Activity ===
        log_activity(
            user_id=user_id,
            activity_type="upload",
            title=f"Uploaded: {meeting.title}",
            description=f"Successfully uploaded {meeting.filename}",
            meeting_id=meeting.id,
            metadata={"filename": meeting.filename, "title": meeting.title}
        )

        # === 6. Start Background Processing ===
        executor.submit(start_processing, meeting.id)

        return jsonify({
            "recording_id": meeting.id,
            "message": "File uploaded successfully",
            "filename": filename,
            "status": "uploaded",
            "transcript_language": transcript_language,
        }), 200

    except Exception as e:
        logger.error(f"Upload failed: {str(e)}")
        return jsonify({"error": f"Upload failed: {str(e)}"}), 500
        
@app.route("/api/process", methods=["POST"])
@jwt_required()
def process_document():
    """
    Accepts:
    - File uploads (PDF, DOCX, PPTX, images)
    - OR a URL link
    Extracts text and sends to Gemini API for processing.
    """
    try:
        user_id = get_jwt_identity()
        user = db.session.get(User, user_id)
        if not user:
            return jsonify({"error": "User not found"}), 404

        extracted_text = ""
        filename = ""
        
        # Handle URL
        if 'url' in request.form or (request.is_json and 'url' in request.json):
            url = request.form.get('url') or request.json.get('url')
            if url:
                try:
                    response = requests.get(url, timeout=10)
                    soup = BeautifulSoup(response.text, 'html.parser')
                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.extract()
                    extracted_text = soup.get_text(separator=' ', strip=True)
                    filename = url
                except Exception as e:
                    return jsonify({"error": f"Failed to extract text from URL: {str(e)}"}), 400
        
        # Handle File Upload
        elif 'file' in request.files:
            file = request.files['file']
            if file and file.filename:
                filename = secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(file_path)
                
                ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
                
                if ext == 'pdf':
                    with open(file_path, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        for page in pdf_reader.pages:
                            extracted_text += page.extract_text() + "\n"
                
                elif ext == 'docx':
                    doc = Document(file_path)
                    for para in doc.paragraphs:
                        extracted_text += para.text + "\n"
                
                elif ext == 'pptx':
                    prs = Presentation(file_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_text += shape.text + "\n"
                
                elif ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                    # Use Gemini Vision for OCR instead of local Tesseract
                    image = Image.open(file_path)
                    prompt = """
                    Extract all text from this image and return it as a raw string. 
                    Do not add any comments or formatting, just the extracted text.
                    """
                    response = call_gemini_api(prompt, model='gemini-2.5-flash', image=image)
                    extracted_text = response.text
                
                else:
                    return jsonify({"error": "Unsupported file format"}), 400
        
        if not extracted_text.strip():
            return jsonify({"error": "No text could be extracted from the source"}), 400

        # Send to Gemini
        prompt = f"""
        Analyze the following text and return a JSON response with these exact keys:
        "transcript": A cleaned, well-formatted full version of the text.
        "summary": A list of key points summarizing the content.
        "action_items": A list of tasks or next steps identified in the text.

        TEXT TO ANALYZE:
        {extracted_text[:30000]} # Limit text length for API
        """
        
        response = call_gemini_api(prompt, model='gemini-2.5-flash')
        ai_response = response.text
        
        # Clean up JSON from markdown if necessary
        if "```json" in ai_response:
            ai_response = ai_response.split("```json")[1].split("```")[0].strip()
        elif "```" in ai_response:
             ai_response = ai_response.split("```")[1].split("```")[0].strip()
        
        try:
            result = json.loads(ai_response)
        except:
            # Fallback if Gemini doesn't return perfect JSON
            result = {
                "transcript": ai_response,
                "summary": ["Could not parse structured summary"],
                "action_items": ["Could not parse structured action items"]
            }

        # Save to database
        new_meeting = Meeting(
            user_id=user_id,
            title=f"Processed: {filename}",
            filename=filename,
            status="completed",
            transcription=json.dumps({
                "raw": result.get("transcript", ""),
                "translated": "",
                "optimized": result.get("transcript", "")
            }),
            notes=json.dumps({
                "summary": "\n".join(result.get("summary", [])) if isinstance(result.get("summary"), list) else result.get("summary", ""),
                "key_points": result.get("summary", []),
                "action_items": result.get("action_items", []),
                "decisions": [],
                "sentiment": "Neutral"
            }),
            has_transcription=True,
            has_notes=True
        )
        db.session.add(new_meeting)
        db.session.commit()

        log_activity(
            user_id=user_id,
            activity_type="completed",
            title=new_meeting.title,
            description=f"Successfully processed {filename}",
            meeting_id=new_meeting.id
        )

        return jsonify({
            "meeting_id": new_meeting.id,
            "recording_id": new_meeting.id,
            "transcript": result.get("transcript", ""),
            "summary": result.get("summary", []),
            "action_items": result.get("action_items", [])
        })

    except Exception as e:
        logger.error(f"Processing error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/processing-status/<int:meeting_id>", methods=["GET"])
@jwt_required()
def processing_status(meeting_id):
    try:
        user_id = get_current_user_id()
        meeting = Meeting.query.filter_by(id=meeting_id, user_id=user_id).first()
        
        if not meeting:
            return jsonify({"error": "Meeting not found"}), 404
        
        try:
            steps = json.loads(meeting.processing_steps or '[]')
        except:
            steps = []
        
        if not steps:
            steps = [
                {"step": "transcription", "status": "pending", "timestamp": "", "error": None},
                {"step": "translation", "status": "pending", "timestamp": "", "error": None},
                {"step": "optimization", "status": "pending", "timestamp": "", "error": None},
                {"step": "ai_generation", "status": "pending", "timestamp": "", "error": None}
            ]
        
        return jsonify({
            "recording_id": meeting.id,
            "status": meeting.status,
            "processing_steps": steps,
            "current_step_progress": meeting.current_step_progress or 0
        })
    except Exception as e:
        logger.error(f"Failed to get processing status: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/meetings", methods=["GET"])
@jwt_required()
def get_meetings():
    try:
        user_id = get_current_user_id()
        limit = request.args.get('limit', 10000000, type=int)
        meetings = Meeting.query.filter_by(user_id=user_id).order_by(Meeting.upload_date.desc()).limit(limit).all()
        
        return jsonify({
            "meetings": [
                {
                    "id": m.id,
                    "title": m.title,
                    "filename": m.filename,
                    "upload_date": m.upload_date.isoformat(),
                    "status": m.status,
                    "transcript_language": m.transcript_language or 'en',  # ADD THIS LINE
                    "has_transcription": m.has_transcription,
                    "has_notes": m.has_notes,
                    "is_favorite": m.is_favorite  # ADD THIS LINE if not already there
                } for m in meetings
            ]
        })
    except Exception as e:
        logger.error(f"Failed to get meetings: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/meetings/<int:meeting_id>", methods=["GET"])
@jwt_required()
def get_meeting(meeting_id):
    try:
        user_id = get_current_user_id()
        meeting = Meeting.query.filter_by(id=meeting_id, user_id=user_id).first()
        
        if not meeting:
            return jsonify({"error": "Meeting not found"}), 404
        
        return jsonify({
            "meeting": {
                "id": meeting.id,
                "title": meeting.title,
                "filename": meeting.filename,
                "upload_date": meeting.upload_date.isoformat(),
                "status": meeting.status,
                "transcript_language": meeting.transcript_language or 'en',  # ADD THIS LINE
                "transcription": json.loads(meeting.transcription or '{}'),
                "notes": json.loads(meeting.notes or '{}'),
                "is_favorite": meeting.is_favorite  # ADD THIS LINE if not already there
            }
        })
    except Exception as e:
        logger.error(f"Failed to get meeting: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/meetings/<int:meeting_id>", methods=["DELETE"])
@jwt_required()
def delete_meeting(meeting_id):
    try:
        user_id = get_current_user_id()
        meeting = Meeting.query.filter_by(id=meeting_id, user_id=user_id).first()
        
        if not meeting:
            return jsonify({"error": "Meeting not found"}), 404
        
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], meeting.filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            logger.info(f"Removed file: {filepath}")
        
        db.session.delete(meeting)
        db.session.commit()
        
        logger.info(f"Meeting {meeting_id} deleted successfully")
        return jsonify({"message": "Meeting deleted successfully"}), 200
    except Exception as e:
        logger.error(f"Failed to delete meeting: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/translate", methods=["POST"])
@jwt_required()
def translate_text():
    try:
        data = request.json
        text = data.get('text', '').strip()
        target_language = data.get('target_language', 'es')
        
        if not text:
            return jsonify({"error": "No text provided"}), 400
        
        language_names = {
            "af": "Afrikaans", "sq": "Albanian", "am": "Amharic", "ar": "Arabic", "hy": "Armenian",
            "az": "Azerbaijani", "eu": "Basque", "be": "Belarusian", "bn": "Bengali", "bs": "Bosnian",
            "bg": "Bulgarian", "ca": "Catalan", "ceb": "Cebuano", "ny": "Chichewa", "zh": "Chinese",
            "zh-cn": "Chinese (Simplified)", "zh-tw": "Chinese (Traditional)", "co": "Corsican",
            "hr": "Croatian", "cs": "Czech", "da": "Danish", "nl": "Dutch", "en": "English",
            "eo": "Esperanto", "et": "Estonian", "tl": "Filipino", "fi": "Finnish", "fr": "French",
            "fy": "Frisian", "gl": "Galician", "ka": "Georgian", "de": "German", "el": "Greek",
            "gu": "Gujarati", "ht": "Haitian Creole", "ha": "Hausa", "haw": "Hawaiian", "he": "Hebrew",
            "iw": "Hebrew", "hi": "Hindi", "hmn": "Hmong", "hu": "Hungarian", "is": "Icelandic",
            "ig": "Igbo", "id": "Indonesian", "ga": "Irish", "it": "Italian", "ja": "Japanese",
            "jw": "Javanese", "kn": "Kannada", "kk": "Kazakh", "km": "Khmer", "ko": "Korean",
            "ku": "Kurdish (Kurmanji)", "ky": "Kyrgyz", "lo": "Lao", "la": "Latin", "lv": "Latvian",
            "lt": "Lithuanian", "lb": "Luxembourgish", "mk": "Macedonian", "mg": "Malagasy",
            "ms": "Malay", "ml": "Malayalam", "mt": "Maltese", "mi": "Maori", "mr": "Marathi",
            "mn": "Mongolian", "my": "Myanmar (Burmese)", "ne": "Nepali", "no": "Norwegian",
            "or": "Odia", "ps": "Pashto", "fa": "Persian", "pl": "Polish", "pt": "Portuguese",
            "pa": "Punjabi", "ro": "Romanian", "ru": "Russian", "sm": "Samoan", "gd": "Scots Gaelic",
            "sr": "Serbian", "st": "Sesotho", "sn": "Shona", "sd": "Sindhi", "si": "Sinhala",
            "sk": "Slovak", "sl": "Slovenian", "so": "Somali", "es": "Spanish", "su": "Sundanese",
            "sw": "Swahili", "sv": "Swedish", "tg": "Tajik", "ta": "Tamil", "te": "Telugu",
            "th": "Thai", "tr": "Turkish", "uk": "Ukrainian", "ur": "Urdu", "ug": "Uyghur",
            "uz": "Uzbek", "vi": "Vietnamese", "cy": "Welsh", "xh": "Xhosa", "yi": "Yiddish",
            "yo": "Yoruba", "zu": "Zulu"
        }
        
        target_lang_name = language_names.get(target_language, "Spanish")
        
        try:
            prompt = f"""
You are a professional translator. Translate the given text to {target_lang_name}. Respond only with the translated text, no additional formatting or explanations.

Text to translate: {text}
"""
            response = call_gemini_api(prompt, model="gemini-2.5-flash")
            translated = response.text.strip() if response.text else ""
            
            if not translated:
                return jsonify({"error": "Translation failed: Empty response from API"}), 500
            
            return jsonify({
                "translated_text": translated
            })
            
        except Exception as e:
            logger.error(f"Translation API error: {str(e)}")
            return jsonify({
                "error": "Translation service unavailable",
                "details": str(e),
                "suggestion": "Please check your GEMINI_API_KEY or network connection and try again."
            }), 500
        
    except Exception as e:
        logger.error(f"Translate endpoint error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/meetings/<int:meeting_id>/ask-ai", methods=["POST"])
@jwt_required()
def ask_ai(meeting_id):
    try:
        user_id = get_current_user_id()
        meeting = db.session.query(Meeting).filter_by(id=meeting_id, user_id=user_id).first()
        if not meeting:
            return jsonify({"error": "Meeting not found"}), 404
            
        data = request.json
        question = data.get('question', '').strip()
        
        if not question:
            return jsonify({"error": "No question provided"}), 400
            
        # Parse notes to get context
        notes_str = meeting.notes if isinstance(meeting.notes, str) else "{}"
        try:
            notes = json.loads(notes_str) if notes_str else {}
        except:
            notes = {}
            
        summary = notes.get('summary', '')
        key_points = notes.get('key_points', [])
        action_items = notes.get('action_items', [])
        decisions = notes.get('decisions', [])
        sentiment = notes.get('sentiment', 'Neutral')
        
        # Format lists for professional context
        kp_text = "\n".join([f"- {kp}" for kp in key_points]) if isinstance(key_points, list) else str(key_points)
        ai_text = "\n".join([f"- {ai}" for ai in action_items]) if isinstance(action_items, list) else str(action_items)
        dec_text = "\n".join([f"- {dec}" for dec in decisions]) if isinstance(decisions, list) else str(decisions)
        
        # Get transcription for deeper context
        trans_str = meeting.transcription if isinstance(meeting.transcription, str) else "{}"
        try:
            trans = json.loads(trans_str) if trans_str else {}
        except:
            trans = {}
            
        transcript_text = trans.get('optimized', trans.get('raw', ''))
        
        # Construct professional context
        context_text = f"""
Meeting Title: {meeting.title}
Overall Sentiment: {sentiment}

Summary:
{summary}

Key Discussion Points:
{kp_text}

Action Items & Next Steps:
{ai_text}

Decisions Made:
{dec_text}

Full Transcript for reference:
{transcript_text[:20000]}
"""
        
        prompt = f"""
You are a highly professional AI meeting assistant. Your goal is to help the user by answering questions based on the meeting data provided below. 

Use the provided Summary, Key Points, Action Items, Decisions, and Transcript to give a comprehensive and professional answer. 
If the information is not explicitly mentioned, you may provide a logical inference but clearly state that it is an inference based on the meeting's context.

Meeting Context:
{context_text}

User Question: {question}

Please provide a clear, professional, and well-structured response.
"""
        response = call_gemini_api(prompt, model="gemini-2.5-flash")
        answer = response.text.strip() if response.text else "I apologize, but I was unable to generate a response at this time."
        
        return jsonify({
            "answer": answer
        })
        
    except Exception as e:
        logger.error(f"Ask AI endpoint error: {str(e)}")
        return jsonify({"error": str(e)}), 500

def create_enhanced_docx(meeting, filepath):
    doc = Document()
    
    # Set default font to Arial for Unicode support
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Arial'
    font.size = Pt(11)
    
    user = db.session.get(User, meeting.user_id)
    organizer_name = user.full_name if user else "AI Assistant"
    
    doc.add_heading(f"Meeting Notes: {meeting.title}", 0)
    
    doc.add_paragraph(f"Organizer: {organizer_name}")
    doc.add_paragraph(f"File: {meeting.filename}")
    doc.add_paragraph(f"Date: {meeting.upload_date.strftime('%Y-%m-%d %H:%M')}")
    
    try:
        notes = json.loads(meeting.notes or '{}')
        logger.info(f"Notes content for Word: {notes}")
        
        if notes.get("summary"):
            doc.add_heading("Executive Summary", level=1)
            p = doc.add_paragraph(notes["summary"])
            if any('\u0600' <= c <= '\u06FF' for c in notes["summary"]):
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        if notes.get("key_points"):
            doc.add_heading("Key Discussion Points", level=1)
            key_points = notes["key_points"] if isinstance(notes["key_points"], list) else json.loads(notes.get("key_points", '[]') or '[]')
            for point in key_points:
                p = doc.add_paragraph(point, style="List Bullet")
                if any('\u0600' <= c <= '\u06FF' for c in point):
                    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        doc.add_heading("Action Items", level=1)
        action_items = notes["action_items"] if isinstance(notes["action_items"], list) else json.loads(notes.get("action_items", '[]') or '[]')
        if action_items:
            for item in action_items:
                p = doc.add_paragraph(item, style="List Number")
                if any('\u0600' <= c <= '\u06FF' for c in item):
                    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        else:
            doc.add_paragraph("No specific action items identified.", style="Normal")
        
        doc.add_heading("Decisions Made", level=1)
        decisions = notes["decisions"] if isinstance(notes["decisions"], list) else json.loads(notes.get("decisions", '[]') or '[]')
        if decisions:
            for decision in decisions:
                p = doc.add_paragraph(decision, style="List Bullet")
                if any('\u0600' <= c <= '\u06FF' for c in decision):
                    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        else:
            doc.add_paragraph("No formal decisions recorded.", style="Normal")
        
        if notes.get("sentiment"):
            doc.add_heading("Overall Sentiment", level=1)
            p = doc.add_paragraph(notes["sentiment"])
            if any('\u0600' <= c <= '\u06FF' for c in notes["sentiment"]):
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        transcription_data = json.loads(meeting.transcription or '{}')
        transcript_text = transcription_data.get('optimized') or transcription_data.get('translated') or transcription_data.get('raw')
        
        if transcript_text:
            doc.add_heading("Full Transcript", level=1)
            # Split transcript by newlines to avoid one giant paragraph
            paragraphs = transcript_text.split('\n')
            for p in paragraphs:
                if p.strip():
                    para = doc.add_paragraph(p.strip())
                    # Check if paragraph is primarily RTL and set alignment
                    if any('\u0600' <= c <= '\u06FF' for c in p):
                        para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            
    except Exception as e:
        doc.add_paragraph(f"Error parsing notes: {str(e)}")
        logger.error(f"Word generation error: {e}")
    
    doc.save(filepath)

def create_enhanced_pdf(meeting, filepath):
    """Create a premium, branded PDF meeting report"""
    try:
        from reportlab.lib import colors
        from reportlab.platypus import Table, TableStyle
        
        doc = SimpleDocTemplate(
            filepath, 
            pagesize=letter,
            rightMargin=50, leftMargin=50,
            topMargin=50, bottomMargin=50
        )
        styles = getSampleStyleSheet()
        story = []
        
        # Define Brand Colors
        brand_blue = colors.HexColor("#3b82f6")
        brand_dark = colors.HexColor("#111827")
        brand_gray = colors.HexColor("#6b7280")
        
        # Title Style
        title_style = ParagraphStyle(
            'PremiumTitle',
            parent=styles['Heading1'],
            fontName=DEFAULT_FONT,
            fontSize=28,
            textColor=brand_blue,
            spaceAfter=10,
            leading=34
        )
        
        # Subtitle Style
        subtitle_style = ParagraphStyle(
            'PremiumSubtitle',
            parent=styles['Normal'],
            fontName=DEFAULT_FONT,
            fontSize=12,
            textColor=brand_gray,
            spaceAfter=30
        )
        
        # Heading Style
        heading_style = ParagraphStyle(
            'PremiumHeading',
            parent=styles['Heading2'],
            fontName=DEFAULT_FONT,
            fontSize=16,
            textColor=brand_blue,
            spaceBefore=25,
            spaceAfter=15,
            textTransform='uppercase',
            letterSpacing=1
        )
        
        # Content Style
        content_style = ParagraphStyle(
            'PremiumContent',
            parent=styles['Normal'],
            fontName=DEFAULT_FONT,
            fontSize=11,
            leading=18,
            textColor=colors.HexColor("#374151")
        )

        # 1. Header Section
        story.append(Paragraph("TalkToText Pro", subtitle_style))
        story.append(Paragraph(meeting.title, title_style))
        
        # Metadata Table
        user = db.session.get(User, meeting.user_id)
        organizer_name = user.full_name if user else "AI Assistant"
        current_time = datetime.now(timezone.utc)
        
        meta_data = [
            [Paragraph(f"<b>Organizer:</b> {organizer_name}", content_style), 
             Paragraph(f"<b>Date:</b> {current_time.strftime('%b %d, %Y')}", content_style),
             Paragraph(f"<b>Time:</b> {current_time.strftime('%I:%M %p')}", content_style)]
        ]
        
        t = Table(meta_data, colWidths=[180, 180, 150])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor("#f8fafc")),
            ('LEFTPADDING', (0, 0), (-1, -1), 15),
            ('RIGHTPADDING', (0, 0), (-1, -1), 15),
            ('TOPPADDING', (0, 0), (-1, -1), 12),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
            ('BOX', (0, 0), (-1, -1), 0.5, colors.HexColor("#e2e8f0")),
            ('ROUNDEDCORNERS', [12, 12, 12, 12])
        ]))
        story.append(t)
        story.append(Spacer(1, 30))

        try:
            notes = json.loads(meeting.notes or '{}')
            
            # 2. Executive Summary
            if notes.get('summary'):
                story.append(Paragraph("Executive Summary", heading_style))
                story.append(Paragraph(fix_text_direction(notes['summary']), content_style))
                story.append(Spacer(1, 10))
            
            # 3. Key Points
            if notes.get('key_points'):
                story.append(Paragraph("Key Points", heading_style))
                key_points = notes['key_points'] if isinstance(notes['key_points'], list) else json.loads(notes.get('key_points', '[]') or '[]')
                for point in key_points:
                    story.append(Paragraph(f"â€¢ {fix_text_direction(point)}", content_style))
                story.append(Spacer(1, 10))
            
            # 4. Action Items
            story.append(Paragraph("Action Items", heading_style))
            action_items = notes['action_items'] if isinstance(notes['action_items'], list) else json.loads(notes.get('action_items', '[]') or '[]')
            if action_items:
                for item in action_items:
                    story.append(Paragraph(f"â–¡ {fix_text_direction(item)}", content_style))
            else:
                story.append(Paragraph("No specific actions pending.", content_style))
            story.append(Spacer(1, 10))

            # 5. Decisions Made (New)
            decisions = notes.get('decisions')
            if decisions:
                story.append(Paragraph("Decisions Made", heading_style))
                decisions_list = decisions if isinstance(decisions, list) else json.loads(decisions or '[]')
                if decisions_list:
                    for decision in decisions_list:
                        story.append(Paragraph(f"âœ“ {fix_text_direction(decision)}", content_style))
                else:
                    story.append(Paragraph("No formal decisions recorded.", content_style))
                story.append(Spacer(1, 10))

            # 6. Overall Sentiment (New)
            if notes.get('sentiment'):
                story.append(Paragraph("Overall Sentiment", heading_style))
                story.append(Paragraph(fix_text_direction(notes['sentiment']), content_style))
                story.append(Spacer(1, 10))
            
            # 5. Full Transcript
            transcription_data = json.loads(meeting.transcription or '{}')
            transcript_text = (
                transcription_data.get('optimized') or 
                transcription_data.get('translated') or 
                transcription_data.get('raw')
            )
            
            if transcript_text:
                story.append(Paragraph("Full Transcript", heading_style))
                paragraphs = transcript_text.split('\n')
                for p in paragraphs:
                    if p.strip():
                        # Detect best font for this specific paragraph
                        current_font = get_font_for_text(p)
                        is_rtl = any('\u0600' <= c <= '\u06FF' for c in p)
                        alignment = 2 if is_rtl else 0
                        
                        dynamic_style = ParagraphStyle(
                            f'Dynamic_{current_font}_{alignment}',
                            parent=content_style,
                            fontName=current_font,
                            alignment=alignment
                        )
                        story.append(Paragraph(fix_text_direction(p.strip()), dynamic_style))
                        story.append(Spacer(1, 8))
            
            # 6. Footer
            story.append(Spacer(1, 40))
            footer_style = ParagraphStyle('Footer', parent=content_style, fontSize=9, textColor=brand_gray, alignment=1)
            story.append(Paragraph("Generated by TalkToText Pro AI â€¢ Your intelligence partner for every conversation.", footer_style))
            story.append(Paragraph("talktotextpro.com", footer_style))

        except Exception as e:
            story.append(Paragraph(f"Report Generation Notice: {str(e)}", content_style))
            logger.error(f"PDF content error: {e}")
        
        doc.build(story)
    except Exception as e:
        logger.error(f"PDF design failed: {e}")
        raise

@app.route("/api/export/<int:id>/<string:format>", methods=["GET"])
@jwt_required()
def export(id, format):
    try:
        user_id = get_current_user_id()
        meeting = Meeting.query.filter_by(id=id, user_id=user_id).first()
        
        if not meeting:
            return jsonify({"error": "Meeting not found"}), 404
        
        os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
        
        if format == "word":
            filepath = os.path.join(app.config['OUTPUT_FOLDER'], f"meeting_notes_{id}.docx")
            create_enhanced_docx(meeting, filepath)
            return send_file(filepath, as_attachment=True, download_name=f"meeting_notes_{id}.docx")
            
        elif format == "pdf":
            filepath = os.path.join(app.config['OUTPUT_FOLDER'], f"meeting_notes_{id}.pdf")
            create_enhanced_pdf(meeting, filepath)
            return send_file(filepath, as_attachment=True, download_name=f"meeting_notes_{id}.pdf")
        
        return jsonify({"error": "Invalid format"}), 400
    except Exception as e:
        logger.error(f"Export failed: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/stats", methods=["GET"])
@jwt_required()
def stats():
    try:
        user_id = get_jwt_identity()
        meetings = Meeting.query.filter_by(user_id=user_id).all()
        
        total_uploads = len(meetings)
        total_words = sum(len(json.loads(m.notes or '{}').get("summary", "").split()) for m in meetings)
        
        today = datetime.now(timezone.utc).date()
        last_7_days = [(today - timedelta(days=i)).strftime("%a") for i in range(6, -1, -1)]
        uploads_by_day = Counter(m.upload_date.date().strftime("%a") for m in meetings)
        uploads_data = [uploads_by_day.get(day, 0) for day in last_7_days]
        
        return jsonify({
            "total_meetings": total_uploads,
            "completed_meetings": len([m for m in meetings if m.status == "completed"]),
            "total_words": total_words,
            "labels": last_7_days,
            "uploads": uploads_data
        })
    except Exception as e:
        logger.error(f"Stats failed: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/chat", methods=["POST"])
def chat():
    try:
        data = request.json
        user_message = data.get('message', '').strip()
        
        if not user_message:
            return jsonify({"error": "No message provided"}), 400
        
        system_prompt = (
            "You are the AI assistant for TalkToText Pro.\n\n"
            "About TalkToText Pro:\n"
            "- It is an AI-powered meeting notes rewriter.\n"
            "- Converts speech from recordings and various platforms into structured, actionable meeting notes.\n"
            "- Features: transcription, translation, text cleaning, summarization, PDF/Word export.\n"
            "- Goal: Help users make their meetings productive, clear, and easy to follow.\n\n"
            "NEW ENHANCED FEATURES:\n"
            "ðŸ“Š Real-Time Activity Feed - Track all your meeting activities in real-time\n"
            "ðŸ“ˆ Advanced Analytics Dashboard - Weekly/Monthly usage insights and performance metrics\n"
            "ðŸŽ¯ Usage Intensity Tracking - Monitor your active days, upload patterns, and peak usage hours\n"
            "âš¡ Processing Performance Metrics - Success rates, average processing times, language distribution\n"
            "ðŸ’¡ Meeting Insights & Analytics - Top discussion topics, upload patterns, meeting frequency analysis\n"
            "ðŸš€ Productivity Metrics - Time saved calculations, action items generated, productivity scoring\n"
            "ðŸ† Leaderboards & Gamification - Compete with other users, track your ranking and achievements\n"
            "ðŸ“± Enhanced Dashboard - Interactive charts, dynamic filtering, and personalized insights\n\n"
            "Your role:\n"
            "- If the user asks about the website, always explain TalkToText Pro in a professional but friendly way.\n"
            "- If the user provides transcripts, summarize them and highlight key points, action items, and decisions.\n"
            "- Keep responses concise, clear, and helpful.\n"
            "- Always be friendly, professional, and focus on helping users understand and use TalkToText Pro effectively.\n"
            "- When users ask about features, highlight both core features AND the new enhanced analytics features.\n"
            "- For analytics questions, explain how they can track their meeting productivity and usage patterns.\n"
        )
        
        try:
            response = call_gemini_api(
                f"{system_prompt}\n\nUser: {user_message}",
                model="gemini-2.5-flash"
            )
            
            ai_response = response.text.strip()
            
            return jsonify({
                "response": ai_response,
                "timestamp": datetime.now(timezone.utc).isoformat()
            })
            
        except Exception as ai_error:
            logger.error(f"AI Chat error: {ai_error}")
            fallback_responses = {
                "features": "TalkToText Pro offers powerful features including:\n\n"
                          "ðŸŽ¯ Core Features:\n"
                          "â€¢ Real-time transcription from Zoom, Google Meet, and Teams\n"
                          "â€¢ Multi-language translation\n"
                          "â€¢ AI-powered text cleaning\n"
                          "â€¢ Smart summarization with key points and action items\n"
                          "â€¢ Export to PDF and Word formats\n\n"
                          "ðŸš€ NEW Enhanced Analytics:\n"
                          "â€¢ Real-Time Activity Feed - Track all your meeting activities\n"
                          "â€¢ Advanced Analytics Dashboard - Weekly/Monthly insights\n"
                          "â€¢ Usage Intensity Tracking - Active days and peak hours\n"
                          "â€¢ Processing Performance Metrics - Success rates and timing\n"
                          "â€¢ Meeting Insights - Top topics and discussion patterns\n"
                          "â€¢ Productivity Metrics - Time saved and action items\n"
                          "â€¢ Leaderboards - Compete with other users\n\n"
                          "What would you like to know more about?",
                
                "analytics": "Our NEW Analytics Suite helps you understand your meeting patterns and productivity:\n\n"
                           "ðŸ“Š Activity Tracking:\n"
                           "â€¢ Real-time feed of all your uploads, processing, and exports\n"
                           "â€¢ Usage intensity with active days and peak hours\n"
                           "â€¢ Upload patterns by day of week\n\n"
                           "âš¡ Performance Insights:\n"
                           "â€¢ Processing success rates and average times\n"
                           "â€¢ Language distribution across your meetings\n"
                           "â€¢ Meeting completion statistics\n\n"
                           "ðŸ’¡ Meeting Intelligence:\n"
                           "â€¢ Top discussion topics from your meetings\n"
                           "â€¢ Most active days for uploads\n"
                           "â€¢ Average meetings per week\n\n"
                           "ðŸš€ Productivity Metrics:\n"
                           "â€¢ Estimated time saved vs manual note-taking\n"
                           "â€¢ Total action items generated\n"
                           "â€¢ Export usage and productivity scoring\n\n"
                           "ðŸ† Competitive Insights:\n"
                           "â€¢ Weekly leaderboards\n"
                           "â€¢ User rankings and achievements\n"
                           "â€¢ Performance comparison\n\n"
                           "Check your Dashboard's Analytics tab to see your personal insights!",
                
                "dashboard": "Your enhanced Dashboard now includes:\n\n"
                           "ðŸ  Overview Tab:\n"
                           "â€¢ Quick stats - Total meetings, completed, this week, processing\n"
                           "â€¢ Interactive activity charts (Weekly/Monthly/Yearly views)\n"
                           "â€¢ Processing status distribution\n"
                           "â€¢ Recent meetings with quick actions\n\n"
                           "ðŸ“ˆ Analytics Tab:\n"
                           "â€¢ Usage Intensity - Active days, upload frequency, peak hours\n"
                           "â€¢ Processing Performance - Success rates, timing, languages\n"
                           "â€¢ Meeting Insights - Top topics, upload patterns\n"
                           "â€¢ Productivity Metrics - Time saved, action items, exports\n\n"
                           "ðŸ”” Activities Tab:\n"
                           "â€¢ Real-time activity feed\n"
                           "â€¢ Upload, processing, and completion notifications\n"
                           "â€¢ Export and sharing activities\n\n"
                           "ðŸ† Leaderboard Tab:\n"
                           "â€¢ Weekly user rankings\n"
                           "â€¢ Points based on completed meetings\n"
                           "â€¢ Top 3 highlighted with special badges\n\n"
                           "The dashboard helps you track your meeting productivity and usage patterns!",
                
                "about": "TalkToText Pro is an AI-powered meeting notes rewriter that helps you convert speech from popular meeting platforms into structured, actionable notes. We make your meetings more productive and easier to follow!\n\n"
                        "NEW: We've just launched enhanced analytics and insights features to help you understand your meeting patterns, track productivity, and compete with other users on our leaderboards!",
                
                "how": "Getting started is simple!\n\n"
                      "1. Upload your meeting recording (audio/video files or YouTube links)\n"
                      "2. Our AI transcribes and processes it with advanced analysis\n"
                      "3. Review the generated notes, summaries, and action items\n"
                      "4. Export in PDF or Word format\n"
                      "5. Track your progress in the new Analytics dashboard\n\n"
                      "NEW: You can now monitor your usage patterns, processing performance, and productivity metrics in real-time!",
                
                "support": "I'm here to help! You can ask me about:\n\n"
                          "â€¢ TalkToText Pro features and capabilities\n"
                          "â€¢ How to use the platform and upload meetings\n"
                          "â€¢ Understanding your analytics and insights\n"
                          "â€¢ Interpreting your productivity metrics\n"
                          "â€¢ Leaderboard and competitive features\n"
                          "â€¢ Or share meeting content for me to summarize\n\n"
                          "What specific question do you have?",
                
                "default": "Thanks for your question! I'm here to help you with TalkToText Pro. You can ask me about:\n\n"
                          "â€¢ Our features and new analytics capabilities\n"
                          "â€¢ How to use the platform\n"
                          "â€¢ Your meeting insights and productivity metrics\n"
                          "â€¢ Leaderboard and competitive features\n"
                          "â€¢ Or share meeting content for analysis\n\n"
                          "How can I assist you today?"
            }
            
            lower_message = user_message.lower()
            if any(word in lower_message for word in ['feature', 'what can', 'capability']):
                fallback_response = fallback_responses["features"]
            elif any(word in lower_message for word in ['analytics', 'insight', 'metric', 'statistic', 'dashboard']):
                fallback_response = fallback_responses["analytics"]
            elif any(word in lower_message for word in ['dashboard', 'overview', 'home']):
                fallback_response = fallback_responses["dashboard"]
            elif any(word in lower_message for word in ['about', 'talktotex', 'website', 'company']):
                fallback_response = fallback_responses["about"]
            elif any(word in lower_message for word in ['how', 'tutorial', 'guide', 'start']):
                fallback_response = fallback_responses["how"]
            elif any(word in lower_message for word in ['help', 'support', 'problem']):
                fallback_response = fallback_responses["support"]
            else:
                fallback_response = fallback_responses["default"]
            
            return jsonify({
                "response": fallback_response,
                "timestamp": datetime.now(timezone.utc).isoformat()
            })
        
    except Exception as e:
        logger.error(f"Chat endpoint error: {e}")
        return jsonify({"error": "Sorry, I'm having trouble right now. Please try again."}), 500

@app.route("/api/contact", methods=["POST"])
@jwt_required(optional=True)
def contact():
    try:
        data = request.json
        name = data.get('name')
        email = data.get('email')
        subject = data.get('subject')
        message = data.get('message')
        
        if not all([name, email, subject, message]):
            return jsonify({"error": "Missing required fields"}), 400
            
        sendgrid_api_key = os.getenv('SENDGRID_API_KEY')
        sendgrid_from_email = os.getenv('SENDGRID_FROM_EMAIL') or 'talktotextpro3@gmail.com'
        
        # ðŸš€ HIGH-END CONTACT EMAIL TEMPLATE ðŸš€
        contact_email_html = f"""
<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>New Contact Inquiry: {subject}</title>
  <style>
    body {{
      margin: 0;
      padding: 0;
      font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
      background-color: #111827;
      color: #ffffff;
    }}
    .email-wrapper {{
      padding: 40px 20px;
      background: radial-gradient(ellipse at top left, rgba(56, 189, 248, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at top right, rgba(168, 85, 247, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at bottom, rgba(16, 185, 129, 0.12) 0%, transparent 50%),
                  #111827;
    }}
    .container {{
      max-width: 800px;
      margin: 0 auto;
      background-color: rgba(31, 41, 55, 0.8);
      border: 1px solid rgba(255, 255, 255, 0.1);
      border-radius: 32px;
      overflow: hidden;
      backdrop-filter: blur(20px);
    }}
    .header {{
      padding: 48px 48px 32px;
      background: linear-gradient(to bottom, rgba(59, 130, 246, 0.1), transparent);
    }}
    .logo {{
      font-size: 20px;
      font-weight: 800;
      color: #60a5fa;
      margin-bottom: 24px;
      text-transform: uppercase;
      letter-spacing: 0.1em;
    }}
    h1 {{
      font-size: 32px;
      font-weight: 800;
      margin: 0;
      color: #ffffff;
      line-height: 1.2;
    }}
    .meta-grid {{
      display: table;
      width: 100%;
      margin: 32px 0;
      border-collapse: collapse;
    }}
    .meta-item {{
      display: table-cell;
      padding: 20px;
      background-color: rgba(255, 255, 255, 0.03);
      border: 1px solid rgba(255, 255, 255, 0.05);
      border-radius: 16px;
    }}
    .meta-label {{
      font-size: 12px;
      font-weight: 600;
      color: #6b7280;
      text-transform: uppercase;
      margin-bottom: 4px;
    }}
    .meta-value {{
      font-size: 16px;
      font-weight: 600;
      color: #f3f4f6;
    }}
    .content {{
      padding: 0 48px 48px;
    }}
    .notes-box {{
      padding: 32px;
      background-color: rgba(17, 24, 39, 0.4);
      border: 1px solid rgba(255, 255, 255, 0.05);
      border-radius: 20px;
      font-size: 16px;
      line-height: 1.8;
      color: #d1d5db;
      white-space: pre-wrap;
    }}
    .footer {{
      padding: 32px 48px;
      background-color: rgba(17, 24, 39, 0.6);
      border-top: 1px solid rgba(255, 255, 255, 0.05);
      text-align: center;
    }}
    .footer-text {{
      font-size: 14px;
      color: #6b7280;
      margin-bottom: 8px;
    }}
    .footer-link {{
      color: #3b82f6;
      text-decoration: none;
      font-weight: 600;
    }}
  </style>
</head>
<body>
  <div class="email-wrapper">
    <div class="container">
      <div class="header">
        <div class="logo">TalkToText Pro</div>
        <h1>New Contact Message</h1>
      </div>
      <div class="content">
        <div class="meta-grid">
           <table width="100%" cellspacing="0" cellpadding="0">
             <tr>
               <td width="50%" style="padding-right: 10px;">
                 <div class="meta-item">
                   <div class="meta-label">From</div>
                   <div class="meta-value">{name}</div>
                 </div>
               </td>
               <td width="50%">
                 <div class="meta-item">
                   <div class="meta-label">Email</div>
                   <div class="meta-value">{email}</div>
                 </div>
               </td>
             </tr>
           </table>
        </div>
        
        <div style="margin-bottom: 24px;">
          <div style="font-size: 14px; font-weight: 700; color: #60a5fa; text-transform: uppercase; margin-bottom: 12px; letter-spacing: 0.05em;">Subject: {subject}</div>
          <div class="notes-box">
{message}
          </div>
        </div>
      </div>
      <div class="footer">
        <div class="footer-text">This inquiry was sent via the TalkToText Pro Contact Form.</div>
        <div class="footer-text">
          <a href="https://talktotextpro.com/" class="footer-link">talktotextpro.com</a> â€¢ Architecture of Intelligence
        </div>
      </div>
    </div>
  </div>
</body>
</html>
        """
        
        if sendgrid_api_key:
            sg = SendGridAPIClient(api_key=sendgrid_api_key)
            mail = Mail(
                from_email=sendgrid_from_email,
                to_emails=sendgrid_from_email, # Send to ourselves
                subject=f"Contact Inquiry: {subject}",
                html_content=contact_email_html
            )
            # Set reply-to so we can answer directly
            mail.reply_to = email
            
            response = sg.send(mail)
            logger.info(f"Contact email sent successfully. Status: {response.status_code}")
            
            return jsonify({"message": "Inquiry sent successfully. We will get back to you soon."}), 200
        else:
            logger.error("SendGrid not configured for contact form")
            return jsonify({"error": "Configuration error. Please try again later."}), 500
            
    except Exception as e:
        logger.error(f"Contact submission error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/send-email", methods=["POST"])
@jwt_required()
def send_email():
    try:
        user_id = get_jwt_identity()
        user = db.session.get(User, user_id)
        
        meeting_id = request.form.get('meeting_id')
        to_email = request.form.get('to_email')
        from_email = request.form.get('from_email') or user.email
        subject = request.form.get('subject')
        body = request.form.get('body')
        
        if 'pdf_file' not in request.files:
            return jsonify({"error": "No PDF file provided"}), 400
        
        pdf_file = request.files['pdf_file']
        if pdf_file.filename == '':
            return jsonify({"error": "No PDF file selected"}), 400
        
        meeting = Meeting.query.filter_by(id=meeting_id, user_id=user_id).first()
        if not meeting:
            return jsonify({"error": "Meeting not found"}), 404

        # Debug: Check if environment variables are loaded
        sendgrid_api_key = os.getenv('SENDGRID_API_KEY')
        sendgrid_from_email = os.getenv('SENDGRID_FROM_EMAIL')
        
        logger.info(f"SendGrid API Key exists: {bool(sendgrid_api_key)}")
        logger.info(f"SendGrid From Email: {sendgrid_from_email}")
        
        if sendgrid_api_key:
            try:
                sg = SendGridAPIClient(api_key=sendgrid_api_key)
                
                # ðŸš€ ULTIMATE PROFESSIONAL EMAIL TEMPLATE ðŸš€
                email_body_html = f"""
<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{meeting.title} - Meeting Notes</title>
  <style>
    body {{
      margin: 0;
      padding: 0;
      font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
      background-color: #111827;
      color: #ffffff;
    }}
    .email-wrapper {{
      padding: 40px 20px;
      background: radial-gradient(ellipse at top left, rgba(56, 189, 248, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at top right, rgba(168, 85, 247, 0.12) 0%, transparent 50%),
                  radial-gradient(ellipse at bottom, rgba(16, 185, 129, 0.12) 0%, transparent 50%),
                  #111827;
    }}
    .container {{
      max-width: 800px;
      margin: 0 auto;
      background-color: rgba(31, 41, 55, 0.8);
      border: 1px solid rgba(255, 255, 255, 0.1);
      border-radius: 32px;
      overflow: hidden;
      backdrop-filter: blur(20px);
    }}
    .header {{
      padding: 48px 48px 32px;
      background: linear-gradient(to bottom, rgba(59, 130, 246, 0.1), transparent);
    }}
    .logo {{
      font-size: 20px;
      font-weight: 800;
      color: #60a5fa;
      margin-bottom: 24px;
      text-transform: uppercase;
      letter-spacing: 0.1em;
    }}
    h1 {{
      font-size: 32px;
      font-weight: 800;
      margin: 0;
      color: #ffffff;
      line-height: 1.2;
    }}
    .meta-grid {{
      display: table;
      width: 100%;
      margin: 32px 0;
      border-collapse: collapse;
    }}
    .meta-item {{
      display: table-cell;
      padding: 20px;
      background-color: rgba(255, 255, 255, 0.03);
      border: 1px solid rgba(255, 255, 255, 0.05);
      border-radius: 16px;
    }}
    .meta-label {{
      font-size: 12px;
      font-weight: 600;
      color: #6b7280;
      text-transform: uppercase;
      margin-bottom: 4px;
    }}
    .meta-value {{
      font-size: 16px;
      font-weight: 600;
      color: #f3f4f6;
    }}
    .content {{
      padding: 0 48px 48px;
    }}
    .notes-box {{
      padding: 32px;
      background-color: rgba(17, 24, 39, 0.4);
      border: 1px solid rgba(255, 255, 255, 0.05);
      border-radius: 20px;
      font-size: 16px;
      line-height: 1.8;
      color: #d1d5db;
      white-space: pre-wrap;
    }}
    .cta-container {{
      margin-top: 40px;
      text-align: center;
    }}
    .button {{
      display: inline-block;
      padding: 18px 36px;
      background: linear-gradient(to right, #3b82f6, #8b5cf6);
      color: #ffffff !important;
      text-decoration: none;
      border-radius: 16px;
      font-weight: 700;
      font-size: 16px;
      box-shadow: 0 10px 25px -5px rgba(59, 130, 246, 0.4);
    }}
    .footer {{
      padding: 32px 48px;
      background-color: rgba(17, 24, 39, 0.6);
      border-top: 1px solid rgba(255, 255, 255, 0.05);
      text-align: center;
    }}
    .footer-text {{
      font-size: 14px;
      color: #6b7280;
      margin-bottom: 8px;
    }}
    .footer-link {{
      color: #3b82f6;
      text-decoration: none;
      font-weight: 600;
    }}
  </style>
</head>
<body>
  <div class="email-wrapper">
    <div class="container">
      <div class="header">
        <div class="logo">TalkToText Pro</div>
        <h1>{meeting.title}</h1>
      </div>
      <div class="content">
        <div class="meta-grid">
           <table width="100%" cellspacing="0" cellpadding="0">
             <tr>
               <td width="33%" style="padding-right: 10px;">
                 <div class="meta-item">
                   <div class="meta-label">Organizer</div>
                   <div class="meta-value">{user.full_name}</div>
                 </div>
               </td>
                <td width="33%" style="padding-right: 10px;">
                  <div class="meta-item">
                    <div class="meta-label">Date</div>
                    <div class="meta-value">{datetime.now(timezone.utc).strftime('%b %d, %Y')}</div>
                  </div>
                </td>
                <td width="33%">
                  <div class="meta-item">
                    <div class="meta-label">Time</div>
                    <div class="meta-value">{datetime.now(timezone.utc).strftime('%I:%M %p')}</div>
                  </div>
                </td>
             </tr>
           </table>
        </div>
        
        <div style="margin-bottom: 24px;">
          <div style="font-size: 14px; font-weight: 700; color: #60a5fa; text-transform: uppercase; margin-bottom: 12px; letter-spacing: 0.05em;">AI Summary & Insights</div>
          <div class="notes-box">
{body}
          </div>
        </div>
      </div>
      <div class="footer">
        <div class="footer-text">This report was automatically generated by TalkToText Pro AI.</div>
        <div class="footer-text">
          <a href="https://talk-to-text-psi.vercel.app/" class="footer-link">talktotextpro.com</a> â€¢ Empowering your meetings
        </div>
      </div>
    </div>
  </div>
</body>
</html>
                """
                
                # Use verified sender email from Railway environment
                verified_from_email = sendgrid_from_email or 'noreply@talktotextpro.com'
                
                message = Mail(
                    from_email=verified_from_email,
                    to_emails=to_email,
                    subject=subject,
                    html_content=email_body_html
                )
                
                # Set reply-to to the user's email
                message.reply_to = from_email
                
                # Add PDF attachment
                pdf_content = pdf_file.read()
                encoded_pdf = base64.b64encode(pdf_content).decode()
                
                attachedFile = Attachment(
                    FileContent(encoded_pdf),
                    FileName(pdf_file.filename),
                    FileType('application/pdf'),
                    Disposition('attachment')
                )
                message.attachment = attachedFile
                
                response = sg.send(message)
                logger.info(f"SendGrid email sent successfully. Status: {response.status_code}")
                
                return jsonify({
                    "message": "Email sent successfully via SendGrid",
                    "status_code": response.status_code
                })
                
            except Exception as sg_error:
                logger.error(f"SendGrid failed: {sg_error}")
                logger.error(f"SendGrid error type: {type(sg_error).__name__}")
                
                # Return more detailed error info
                error_details = str(sg_error)
                if hasattr(sg_error, 'body'):
                    error_details += f" | Body: {sg_error.body}"
                
                return jsonify({
                    "error": "SendGrid email failed",
                    "details": error_details,
                    "api_key_configured": bool(sendgrid_api_key),
                    "from_email_configured": bool(sendgrid_from_email)
                }), 500
        
        else:
            logger.error("SendGrid API key not found in environment variables")
            return jsonify({
                "error": "SendGrid not configured",
                "details": "SENDGRID_API_KEY environment variable not found",
                "suggestion": "Set SENDGRID_API_KEY in Railway environment variables"
            }), 500
        
    except Exception as e:
        logger.error(f"Send email error: {e}")
        return jsonify({
            "error": "Email sending failed",
            "details": str(e),
            "type": type(e).__name__
        }), 500
        
        

# This should already be in your app.py around line 150-160
@app.route("/api/meetings/<int:meeting_id>/favorite", methods=["PUT"])
@jwt_required()
def toggle_favorite(meeting_id):
    user_id = get_jwt_identity()
    meeting = Meeting.query.filter_by(id=meeting_id, user_id=user_id).first()
    if not meeting:
        return jsonify({"error": "Meeting not found"}), 404
    data = request.json
    if 'is_favorite' not in data:
        return jsonify({"error": "Missing is_favorite"}), 400
    meeting.is_favorite = data['is_favorite']
    db.session.commit()
    return jsonify({"success": True, "is_favorite": meeting.is_favorite})

@app.route("/api/favorites", methods=["GET"])
@jwt_required()
def get_favorites():
    user_id = get_jwt_identity()
    favorites = Meeting.query.filter_by(user_id=user_id, is_favorite=True).order_by(Meeting.upload_date.desc()).all()
    return jsonify({
        "meetings": [{
            "id": m.id,
            "title": m.title,
            "filename": m.filename,
            "upload_date": m.upload_date.isoformat(),
            "status": m.status,
            "has_transcription": m.has_transcription,
            "has_notes": m.has_notes,
            "is_favorite": m.is_favorite
        } for m in favorites]
    })

# Serve chat media files (token via ?token= query param since audio/video can't set headers)
@app.route('/uploads/chat/<path:filename>', methods=['GET'])
def serve_chat_media(filename):
    """Serve media files for chat with authentication."""
    token = request.args.get('token') or (
        request.headers.get('Authorization', '').replace('Bearer ', '') or None
    )
    if not token:
        return jsonify({"error": "Authentication required"}), 401
    try:
        from flask_jwt_extended import decode_token
        decoded = decode_token(token)
        if not db.session.get(User, decoded['sub']):
            return jsonify({"error": "User not found"}), 404
    except Exception:
        return jsonify({"error": "Invalid token"}), 401
    
    # Secure the filename behavior - chat files are in uploads/chat/
    safe_filename = secure_filename(filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], 'chat', safe_filename)
    
    if not os.path.exists(filepath):
        return jsonify({"error": "File not found"}), 404
    return send_file(filepath, as_attachment=False)

@app.route("/uploads/<path:filename>", methods=["GET"])
def serve_upload(filename):
    """Serve uploaded media files for preview - with token from query parameter"""
    try:
        # Get token from query parameter (since HTML video/audio can't send headers)
        token = request.args.get('token')
        
        if not token:
            # Also check Authorization header as fallback
            auth_header = request.headers.get('Authorization')
            if auth_header and auth_header.startswith('Bearer '):
                token = auth_header.split(' ')[1]
        
        if not token:
            logger.error("No token provided for media file access")
            return jsonify({"error": "Authentication required"}), 401
        
        # Verify the token manually
        try:
            from flask_jwt_extended import decode_token
            decoded = decode_token(token)
            user_id = decoded['sub']
            
            # Verify user exists
            user = db.session.get(User, user_id)
            if not user:
                return jsonify({"error": "User not found"}), 404
                
        except Exception as e:
            logger.error(f"Token verification failed: {e}")
            return jsonify({"error": "Invalid or expired token"}), 401
        
        # Secure the filename
        safe_filename = secure_filename(filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)
        
        if not os.path.exists(filepath):
            logger.error(f"File not found: {filepath}")
            return jsonify({"error": "File not found"}), 404
        
        # Check if the file belongs to the user (extra security)
        meeting = Meeting.query.filter_by(filename=safe_filename, user_id=user_id).first()
        if not meeting:
            logger.error(f"User {user_id} tried to access unauthorized file: {safe_filename}")
            return jsonify({"error": "Unauthorized access"}), 403
        
        logger.info(f"Serving file {safe_filename} to user {user_id}")
        
        # Serve the file with proper MIME type
        return send_file(filepath, mimetype=None, as_attachment=False)
        
    except Exception as e:
        logger.error(f"Error serving file: {e}")
        return jsonify({"error": str(e)}), 500

# ===============================================================
# ðŸ“Š ACTIVITY & ANALYTICS ENDPOINTS
# ===============================================================

@app.route("/api/activities", methods=["GET"])
@jwt_required()
def get_activities():
    try:
        user_id = get_jwt_identity()
        limit = request.args.get('limit', 50, type=int)
        offset = request.args.get('offset', 0, type=int)
        
        activities = Activity.query.filter_by(user_id=user_id)\
            .order_by(Activity.timestamp.desc())\
            .offset(offset).limit(limit).all()
        
        return jsonify({
            "activities": [
                {
                    "id": a.id,
                    "type": a.type,
                    "title": a.title,
                    "description": a.description,
                    "timestamp": a.timestamp.isoformat(),
                    "meeting_id": a.meeting_id,
                    "metadata": json.loads(a.activity_metadata or '{}')  # Updated field name
                } for a in activities
            ],
            "total": Activity.query.filter_by(user_id=user_id).count()
        })
        
    except Exception as e:
        logger.error(f"Failed to get activities: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/analytics/usage-intensity", methods=["GET"])
@jwt_required()
def get_usage_intensity():
    try:
        user_id = get_jwt_identity()
        period = request.args.get('period', 'weekly')  # weekly, monthly, yearly
        
        # Calculate active days
        if period == 'weekly':
            start_date = datetime.now(timezone.utc) - timedelta(days=7)
        elif period == 'monthly':
            start_date = datetime.now(timezone.utc) - timedelta(days=30)
        else:  # yearly
            start_date = datetime.now(timezone.utc) - timedelta(days=365)
        
        # Active days calculation
        active_days_query = db.session.query(
            db.func.date(Activity.timestamp).distinct()
        ).filter(
            Activity.user_id == user_id,
            Activity.timestamp >= start_date
        )
        active_days = active_days_query.count()
        
        # Upload frequency
        uploads_count = Meeting.query.filter(
            Meeting.user_id == user_id,
            Meeting.upload_date >= start_date
        ).count()
        
        # Get all meetings in period for calculations
        meetings = Meeting.query.filter(
            Meeting.user_id == user_id,
            Meeting.upload_date >= start_date
        ).all()
        
        # Calculate active days from meetings (more reliable than activities)
        meeting_dates = set()
        for meeting in meetings:
            if meeting.upload_date:
                meeting_dates.add(meeting.upload_date.date())
        
        # Also include activity dates
        activity_dates = set()
        activities_in_period = Activity.query.filter(
            Activity.user_id == user_id,
            Activity.timestamp >= start_date
        ).all()
        for activity in activities_in_period:
            if activity.timestamp:
                activity_dates.add(activity.timestamp.date())
        
        # Combine both for total active days
        all_active_dates = meeting_dates.union(activity_dates)
        active_days = len(all_active_dates)
        
        # Calculate average duration per day based on meeting durations
        total_duration_seconds = 0
        for meeting in meetings:
            # Check if meeting has duration field and it's not None
            if hasattr(meeting, 'duration') and meeting.duration:
                total_duration_seconds += meeting.duration
            else:
                # If no duration, estimate 30 minutes per meeting as baseline
                total_duration_seconds += 1800  # 30 minutes in seconds
        
        # Calculate average duration per active day
        if active_days > 0:
            avg_duration_seconds = total_duration_seconds / active_days
            avg_hours = int(avg_duration_seconds // 3600)
            avg_minutes = int((avg_duration_seconds % 3600) // 60)
            avg_duration_str = f"{avg_hours}:{avg_minutes:02d}"
        else:
            avg_duration_str = "0:00"
        
        # Peak hours based on meeting upload times
        hour_counts = {}
        for meeting in meetings:
            if meeting.upload_date:
                hour = meeting.upload_date.hour
                hour_counts[hour] = hour_counts.get(hour, 0) + 1
        
        # Also include activity hours
        for activity in activities_in_period:
            if activity.timestamp:
                hour = activity.timestamp.hour
                hour_counts[hour] = hour_counts.get(hour, 0) + 1
        
        # Sort by count and create peak_hours list
        if hour_counts:
            peak_hours = [
                {"hour": hour, "count": count}
                for hour, count in sorted(hour_counts.items(), key=lambda x: x[1], reverse=True)
            ]
        else:
            peak_hours = [{"hour": 0, "count": 0}]
        
        return jsonify({
            "active_days": active_days,
            "uploads_count": uploads_count,
            "avg_activities_per_day": avg_duration_str,
            "peak_hours": peak_hours,
            "period": period
        })
        
    except Exception as e:
        logger.error(f"Failed to get usage intensity: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/analytics/processing-performance", methods=["GET"])
@jwt_required()
def get_processing_performance():
    try:
        user_id = get_jwt_identity()
        
        # Get processing times from completed meetings
        processing_data = db.session.query(
            Meeting.status,
            db.func.count(Meeting.id).label('count'),
            db.func.avg(
                db.func.extract('epoch', Meeting.upload_date) - 
                db.func.extract('epoch', Meeting.upload_date)
            ).label('avg_time')  # Simplified - you'd need processing start/end times
        ).filter(
            Meeting.user_id == user_id
        ).group_by(Meeting.status).all()
        
        # Language distribution
        language_distribution = db.session.query(
            Meeting.transcript_language,
            db.func.count(Meeting.id).label('count')
        ).filter(
            Meeting.user_id == user_id,
            Meeting.transcript_language.isnot(None)
        ).group_by(Meeting.transcript_language).all()
        
        # Success rate
        total_meetings = Meeting.query.filter_by(user_id=user_id).count()
        completed_meetings = Meeting.query.filter_by(user_id=user_id, status='completed').count()
        success_rate = (completed_meetings / total_meetings * 100) if total_meetings > 0 else 0
        
        return jsonify({
            "success_rate": round(success_rate, 1),
            "total_processed": total_meetings,
            "completed": completed_meetings,
            "processing": Meeting.query.filter_by(user_id=user_id, status='processing').count(),
            "failed": Meeting.query.filter_by(user_id=user_id, status='failed').count(),
            "language_distribution": [
                {"language": lang, "count": count} 
                for lang, count in language_distribution
            ],
            "average_processing_time": 120  # Placeholder - you'd calculate actual time
        })
        
    except Exception as e:
        logger.error(f"Failed to get processing performance: {e}")
        return jsonify({"error": str(e)}), 500
@app.route("/api/analytics/meeting-insights", methods=["GET"])
@jwt_required()
def get_meeting_insights():
    try:
        user_id = get_jwt_identity()
        
        # COMPREHENSIVE stop words list - including the words you're seeing
        stop_words = {
            # Common English stop words
            'a', 'an', 'the', 'and', 'or', 'but', 'if', 'because', 'as', 'what',
            'when', 'where', 'how', 'why', 'all', 'any', 'both', 'each', 'few',
            'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'not', 'only',
            'own', 'same', 'so', 'than', 'too', 'very', 'can', 'will', 'just',
            'should', 'now', 'with', 'that', 'from', 'they', 'specific',  # ADDED YOUR PROBLEM WORDS
            
            # Your existing stop words
            'speaker', 'their', 'about', 'discussion', 'strong', 'would', 'could',
            'there', 'which', 'other', 'these', 'those', 'being', 'having', 'doing',
            'through', 'during', 'before', 'after', 'above', 'below', 'between',
            'under', 'while', 'since', 'until', 'upon', 'regarding', 'following',
            'according', 'another', 'every', 'several', 'various', 'different',
            'important', 'because', 'however', 'therefore', 'moreover', 'furthermore',
            'nevertheless', 'otherwise', 'additionally', 'consequently',
            
            # More common words that aren't meaningful topics
            'this', 'that', 'these', 'those', 'them', 'then', 'than', 'here', 'there',
            'when', 'where', 'why', 'how', 'what', 'which', 'who', 'whom', 'whose',
            'been', 'have', 'has', 'had', 'having', 'do', 'does', 'did', 'doing',
            'will', 'would', 'could', 'should', 'may', 'might', 'must', 'shall'
        }

        # Expanded meaningful business terms
        meaningful_business_terms = {
            'project', 'team', 'meeting', 'client', 'budget', 'timeline', 'deadline',
            'strategy', 'marketing', 'sales', 'product', 'development', 'design',
            'analysis', 'research', 'planning', 'execution', 'implementation',
            'review', 'feedback', 'performance', 'metrics', 'goals', 'objectives',
            'results', 'outcomes', 'decisions', 'actions', 'tasks', 'followup',
            'stakeholder', 'manager', 'director', 'executive', 'department',
            'quarter', 'annual', 'weekly', 'monthly', 'progress', 'update',
            'business', 'company', 'organization', 'initiative', 'proposal',
            'solution', 'problem', 'challenge', 'opportunity', 'risk', 'issue',
            'agenda', 'minutes', 'summary', 'conclusion', 'recommendation',
            'priority', 'milestone', 'deliverable', 'resource', 'allocation',
            'cost', 'revenue', 'profit', 'loss', 'investment', 'return',
            'customer', 'user', 'market', 'competitor', 'industry', 'sector',
            'technology', 'innovation', 'digital', 'transformation', 'change',
            'leadership', 'management', 'collaboration', 'communication',
            'efficiency', 'productivity', 'quality', 'standard', 'process',
            'workflow', 'automation', 'optimization', 'improvement'
        }

        meetings = Meeting.query.filter_by(user_id=user_id).all()
        
        # Topic analysis from AI notes
        all_notes = []
        for meeting in meetings:
            if meeting.notes:
                try:
                    notes_data = json.loads(meeting.notes)
                    if isinstance(notes_data, dict):
                        # Combine summary and key points for better topic extraction
                        summary = notes_data.get('summary', '')
                        key_points = notes_data.get('key_points', [])
                        action_items = notes_data.get('action_items', [])
                        decisions = notes_data.get('decisions', [])
                        
                        # Create a comprehensive text for analysis
                        combined_text = f"{summary} {' '.join(key_points)} {' '.join(action_items)} {' '.join(decisions)}"
                        all_notes.append(combined_text)
                except:
                    pass
        
        # Enhanced word frequency analysis with MUCH better filtering
        word_freq = {}
        for note in all_notes:
            if isinstance(note, str):
                # Clean and tokenize text - only keep words that are likely meaningful
                words = re.findall(r'\b[a-zA-Z]{5,}\b', note.lower())  # Increased to 5+ characters
                for word in words:
                    # Skip if it's a stop word
                    if word in stop_words:
                        continue
                    
                    # Skip common verb endings and generic patterns
                    if word.endswith(('ing', 'ed', 'ly', 's', 'es', 'ment', 'tion', 'ity')):
                        # But keep if it's a meaningful business term
                        if word not in meaningful_business_terms:
                            continue
                    
                    word_freq[word] = word_freq.get(word, 0) + 1
        
        # Get top topics, prioritizing meaningful business terms
        top_topics = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]
        
        # If we still don't have good topics, use phrase extraction
        if len(top_topics) < 5 or max(freq for _, freq in top_topics) < 2:
            logger.info("Using phrase extraction as fallback")
            phrases_freq = {}
            for note in all_notes:
                if isinstance(note, str):
                    # Clean the text first
                    clean_note = re.sub(r'[^\w\s]', ' ', note)
                    words = clean_note.lower().split()
                    
                    # Create 2-3 word phrases
                    for i in range(len(words) - 1):
                        # 2-word phrases
                        phrase2 = f"{words[i]} {words[i+1]}"
                        if (len(phrase2) > 10 and 
                            not any(stop in phrase2.split() for stop in stop_words) and
                            not phrase2.endswith(('ing', 'ed', 'ly'))):
                            phrases_freq[phrase2] = phrases_freq.get(phrase2, 0) + 1
                        
                        # 3-word phrases (if available)
                        if i < len(words) - 2:
                            phrase3 = f"{words[i]} {words[i+1]} {words[i+2]}"
                            if (len(phrase3) > 15 and 
                                not any(stop in phrase3.split() for stop in stop_words)):
                                phrases_freq[phrase3] = phrases_freq.get(phrase3, 0) + 1
            
            # Get top phrases and add them
            top_phrases = sorted(phrases_freq.items(), key=lambda x: x[1], reverse=True)[:8]
            top_topics = [(phrase, freq) for phrase, freq in top_phrases]

        # Final filtering - remove any topics that are too generic
        final_topics = []
        for topic, freq in top_topics:
            topic_lower = topic.lower()
            # Skip if it contains any stop words or is too short
            if (len(topic) < 5 or 
                any(stop in topic_lower.split() for stop in stop_words) or
                topic_lower in stop_words):
                continue
            final_topics.append((topic, freq))
        
        # Take top 5 final topics
        final_topics = final_topics[:5]

        # Upload patterns
        upload_patterns = db.session.query(
            db.func.extract('dow', Meeting.upload_date).label('day_of_week'),
            db.func.count(Meeting.id).label('count')
        ).filter(
            Meeting.user_id == user_id
        ).group_by('day_of_week').all()
        
        days = ['Sunday', 'Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday']
        upload_by_day = [{"day": days[int(day)], "count": count} for day, count in upload_patterns]
        
        return jsonify({
            "total_meetings_analyzed": len(meetings),
            "top_topics": [{"topic": topic.title(), "frequency": freq} for topic, freq in final_topics],
            "upload_patterns": upload_by_day,
            "average_meetings_per_week": len(meetings) / max((datetime.now(timezone.utc) - min(m.upload_date for m in meetings)).days / 7, 1),
            "most_active_day": max(upload_by_day, key=lambda x: x['count']) if upload_by_day else None
        })
        
    except Exception as e:
        logger.error(f"Failed to get meeting insights: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/analytics/productivity", methods=["GET"])
@jwt_required()
def get_productivity_metrics():
    try:
        user_id = get_jwt_identity()
        
        # Calculate time saved (estimate 30 minutes per meeting vs manual notes)
        completed_meetings = Meeting.query.filter_by(user_id=user_id, status='completed').count()
        estimated_time_saved = completed_meetings * 30  # minutes
        
        # Action items completion (placeholder - you'd track this)
        total_action_items = 0
        for meeting in Meeting.query.filter_by(user_id=user_id, status='completed'):
            if meeting.notes:
                try:
                    notes_data = json.loads(meeting.notes)
                    action_items = notes_data.get('action_items', [])
                    total_action_items += len(action_items)
                except:
                    pass
        
        # Export usage
        export_activities = Activity.query.filter_by(
            user_id=user_id, 
            type='export'
        ).count()
        
        return jsonify({
            "estimated_time_saved_minutes": estimated_time_saved,
            "estimated_time_saved_hours": round(estimated_time_saved / 60, 1),
            "total_action_items_generated": total_action_items,
            "average_action_items_per_meeting": round(total_action_items / max(completed_meetings, 1), 1),
            "exports_count": export_activities,
            "productivity_score": min(100, (completed_meetings * 10 + total_action_items * 2 + export_activities * 5))
        })
        
    except Exception as e:
        logger.error(f"Failed to get productivity metrics: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/leaderboard", methods=["GET"])
@jwt_required()
def get_leaderboard():
    try:
        period = request.args.get('period', 'weekly')  # weekly, monthly, all_time
        
        # Calculate leaderboard scores
        # Strict Monthly Logic: Reset every month (start from 1st of current month)
        today = datetime.now(timezone.utc)
        start_date = datetime(today.year, today.month, 1)
        
        # Override period to always be monthly for now per requirements
        period = 'monthly'
        
        # Get top users by meeting completion
        leaderboard_data = db.session.query(
            User.id,
            User.full_name,
            User.email,
            db.func.count(Meeting.id).label('meetings_count'),
            db.func.sum(
                db.case(
                    (Meeting.status == 'completed', 1),
                    else_=0
                )
            ).label('completed_count')
        ).join(Meeting, User.id == Meeting.user_id).filter(
            Meeting.upload_date >= start_date,
            User.is_active == True
        ).group_by(User.id, User.full_name, User.email).order_by(
            db.desc('completed_count')
        ).limit(5).all()
        
        # Calculate scores (completed meetings + bonus for exports + bonus for consistency)
        ranked_users = []
        for rank, (user_id, full_name, email, total_meetings, completed_count) in enumerate(leaderboard_data, 1):
            # Simple scoring algorithm
            score = completed_count * 10
            
            # Add current user flag
            current_user_id = get_jwt_identity()
            is_current_user = (user_id == current_user_id)
            
            ranked_users.append({
                "rank": rank,
                "user_id": user_id,
                "full_name": full_name or email.split('@')[0],
                "email": email,
                "score": score,
                "completed_meetings": completed_count,
                "total_meetings": total_meetings,
                "is_current_user": is_current_user
            })
        
        return jsonify({
            "leaderboard": ranked_users,
            "period": period,
            "updated_at": datetime.now(timezone.utc).isoformat()
        })
        
    except Exception as e:
        logger.error(f"Failed to get leaderboard: {e}")
        return jsonify({"error": str(e)}), 500

# Helper function to log activities
def log_activity(user_id, activity_type, title, description, meeting_id=None, metadata=None):
    try:
        activity = Activity(
            user_id=user_id,
            type=activity_type,
            title=title,
            description=description,
            meeting_id=meeting_id,
            activity_metadata=json.dumps(metadata or {})  # Updated field name
        )
        db.session.add(activity)
        db.session.commit()
        logger.info(f"Logged activity: {activity_type} for user {user_id}")
        
        # Real-time update via Socket.IO
        try:
            socketio.emit('activity_update', {
                'user_id': user_id,
                'type': activity_type,
                'title': title,
                'description': description,
                'timestamp': datetime.now(timezone.utc).isoformat()
            })
        except Exception as socket_error:
            logger.warning(f"Failed to emit socket event: {socket_error}")
    except Exception as e:
        logger.error(f"Failed to log activity: {e}")
        db.session.rollback()

# Friends and Messaging Blueprint
try:
    print(f"DEBUG: Blueprint is {Blueprint}")
except NameError:
    print("DEBUG: Blueprint is NOT defined! Re-importing...")
    from flask import Blueprint

friends_bp = Blueprint('friends', __name__)

@friends_bp.route('/api/friends', methods=['GET'])
@jwt_required()
def get_friends():
    user_id = int(get_jwt_identity())
    
    # Fetch all non-deleted friendships involving this user
    # Also include blocked friendships so the user can see/manage blocked contacts
    friendships = Friendship.query.filter(
        ((Friendship.user_id == user_id) | (Friendship.friend_id == user_id)),
        Friendship.is_deleted == False  # CRITICAL: exclude soft-deleted friendships
    ).all()
    
    friends = []
    for f in friendships:
        other_id = f.friend_id if f.user_id == user_id else f.user_id
        
        other_user = db.session.get(User, other_id)
        if not other_user:
            continue

        # Determine block direction
        i_blocked_them = bool(f.is_blocked and f.blocked_by_id == user_id and f.user_id != f.friend_id)
        they_blocked_me = bool(f.is_blocked and f.blocked_by_id != user_id and f.blocked_by_id is not None and f.user_id != f.friend_id)

        # Check for unread messages
        unread_count = Message.query.filter_by(
            sender_id=other_id,
            receiver_id=user_id,
            is_read=False
        ).count()

        friends.append({
            "id": other_user.id,
            "name": other_user.full_name,
            "email": other_user.email,
            "image": other_user.image,
            "isFriend": True,
            "status": other_user.status,
            "lastSeen": other_user.last_seen.isoformat() if other_user.last_seen else None,
            "bio": other_user.bio,
            "isPinned": f.is_pinned,
            "isBlocked": i_blocked_them,
            "hasBlockedMe": they_blocked_me,
            "isMuted": f.is_muted,
            "isArchived": f.is_archived,
            "isFavourite": f.is_favourite,
            "unreadCount": unread_count
        })
    return jsonify(friends)

@friends_bp.route('/api/groups', methods=['GET'])
@jwt_required()
def get_groups():
    user_id = int(get_jwt_identity())
    # Get all groups where the user is a member
    memberships = GroupMember.query.filter_by(user_id=user_id).all()
    group_ids = {m.group_id for m in memberships}
    
    groups = []
    for g_id in group_ids:
        group = db.session.get(Group, g_id)
        if group:
            # Get all members; split into active (non-exited) for memberIds
            group_members = GroupMember.query.filter_by(group_id=group.id).all()
            active_members = [gm for gm in group_members if not gm.is_exited]
            member_ids = [gm.user_id for gm in active_members]
            admin_ids  = [gm.user_id for gm in active_members if gm.role == 'admin']

            # Per-user preferences from this user's membership row
            my_membership = next((gm for gm in group_members if gm.user_id == user_id), None)
            is_exited    = my_membership.is_exited    if my_membership else False
            is_pinned    = my_membership.is_pinned    if my_membership else False
            is_archived  = my_membership.is_archived  if my_membership else False
            is_favourite = my_membership.is_favourite if my_membership else False
            is_muted     = my_membership.is_muted     if my_membership else False

            # Unread messages for this user in this group
            read_msg_ids = db.session.query(MessageReceipt.message_id).filter(
                MessageReceipt.user_id == user_id,
                MessageReceipt.read_at != None
            ).subquery()
            unread_count = Message.query.filter(
                Message.group_id == group.id,
                Message.sender_id != user_id,
                ~Message.id.in_(read_msg_ids)
            ).count()

            groups.append({
                "id":          group.id,
                "name":        group.name,
                "image":       group.image,
                "isGroup":     True,
                "bio":         f"{len(active_members)} members",
                "description": group.description,
                "memberIds":   member_ids,
                "adminIds":    admin_ids,
                "creatorId":   group.created_by_id,
                "createdAt":   group.created_at.isoformat(),
                "isExited":    is_exited,
                "isPinned":    is_pinned,
                "isArchived":  is_archived,
                "isFavourite": is_favourite,
                "isMuted":     is_muted,
                "unreadCount": unread_count,
                "groupSettings": {
                    "onlyAdminsCanEditInfo": False,
                    "onlyAdminsCanAddMembers": False,
                    "onlyAdminsCanSendMessages": False
                }
            })
    return jsonify(groups)

@friends_bp.route('/api/groups', methods=['POST'])
@jwt_required()
def create_group():
    user_id = int(get_jwt_identity())
    data = request.json
    name = data.get('name')
    image = data.get('image')
    description = data.get('description', '')
    member_ids = data.get('member_ids', [])
    
    if not name:
        return jsonify({"error": "Group name is required"}), 400
        
    new_group = Group(
        name=name,
        image=image,
        description=description,
        created_by_id=user_id
    )
    db.session.add(new_group)
    db.session.flush() # Get ID before commit
    
    # Add creator as admin
    db.session.add(GroupMember(group_id=new_group.id, user_id=user_id, role='admin'))
    
    # Send invites to other members instead of adding them directly
    for m_id in member_ids:
        if int(m_id) != user_id:
            existing_invite = GroupInvite.query.filter_by(group_id=new_group.id, invitee_id=int(m_id), status='pending').first()
            if not existing_invite:
                db.session.add(GroupInvite(
                    group_id=new_group.id,
                    inviter_id=user_id,
                    invitee_id=int(m_id),
                    status='pending'
                ))
            
    db.session.commit()
    
    return jsonify({
        "id": new_group.id,
        "name": new_group.name,
        "success": True
    })

@friends_bp.route('/api/groups/invites', methods=['GET'])
@jwt_required()
def get_group_invites():
    user_id = int(get_jwt_identity())
    invites = GroupInvite.query.filter_by(invitee_id=user_id, status='pending').all()

    result = []
    for i in invites:
        group = db.session.get(Group, i.group_id)
        inviter = db.session.get(User, i.inviter_id)
        if not group or not inviter:
            continue
        member_count = GroupMember.query.filter_by(group_id=group.id).count()
        result.append({
            "id": i.id,
            "groupId": group.id,
            "groupName": group.name,
            "groupImage": group.image,
            "groupDescription": group.description,
            "inviterName": inviter.full_name,
            "inviterImage": inviter.image,
            "memberCount": member_count,
            "timestamp": i.created_at.isoformat()
        })
    return jsonify(result)


@friends_bp.route('/api/groups/<int:group_id>/invite', methods=['POST'])
@jwt_required()
def send_group_invite(group_id):
    inviter_id = int(get_jwt_identity())
    data = request.json
    user_ids = data.get('user_ids', [])

    if not user_ids:
        return jsonify({"error": "No users specified"}), 400

    group = db.session.get(Group, group_id)
    if not group:
        return jsonify({"error": "Group not found"}), 404

    # Check inviter is a member of the group
    inviter_membership = GroupMember.query.filter_by(group_id=group_id, user_id=inviter_id).first()
    if not inviter_membership:
        return jsonify({"error": "You are not a member of this group"}), 403

    sent = []
    skipped = []
    for uid in user_ids:
        uid = int(uid)
        # Skip if already a member
        if GroupMember.query.filter_by(group_id=group_id, user_id=uid).first():
            skipped.append({"userId": uid, "reason": "already_member"})
            continue
        # Skip if already has a pending invite
        if GroupInvite.query.filter_by(group_id=group_id, invitee_id=uid, status='pending').first():
            skipped.append({"userId": uid, "reason": "already_invited"})
            continue
        invite = GroupInvite(
            group_id=group_id,
            inviter_id=inviter_id,
            invitee_id=uid,
            status='pending',
            created_at=datetime.now(timezone.utc)
        )
        db.session.add(invite)
        sent.append(uid)

    db.session.commit()
    return jsonify({"success": True, "sent": sent, "skipped": skipped})


@friends_bp.route('/api/groups/invites/<int:invite_id>/accept', methods=['POST'])
@jwt_required()
def accept_group_invite(invite_id):
    user_id = int(get_jwt_identity())
    invite = db.session.get(GroupInvite, invite_id)

    if not invite or invite.invitee_id != user_id:
        return jsonify({"error": "Invite not found"}), 404
    if invite.status != 'pending':
        return jsonify({"error": "Invite already handled"}), 400

    invite.status = 'accepted'
    # Add as group member if not already
    existing = GroupMember.query.filter_by(group_id=invite.group_id, user_id=user_id).first()
    if not existing:
        db.session.add(GroupMember(group_id=invite.group_id, user_id=user_id, role='member'))
    db.session.commit()
    return jsonify({"success": True, "groupId": invite.group_id})


@friends_bp.route('/api/groups/invites/<int:invite_id>/reject', methods=['POST'])
@jwt_required()
def reject_group_invite(invite_id):
    user_id = int(get_jwt_identity())
    invite = db.session.get(GroupInvite, invite_id)

    if not invite or invite.invitee_id != user_id:
        return jsonify({"error": "Invite not found"}), 404
    if invite.status != 'pending':
        return jsonify({"error": "Invite already handled"}), 400

    invite.status = 'rejected'
    db.session.commit()
    return jsonify({"success": True})


@friends_bp.route('/api/groups/<int:group_id>/pending-invites', methods=['GET'])
@jwt_required()
def get_pending_invites_for_group(group_id):
    user_id = int(get_jwt_identity())
    # Only members can query this
    if not GroupMember.query.filter_by(group_id=group_id, user_id=user_id).first():
        return jsonify({"error": "Not a member"}), 403
    pending = GroupInvite.query.filter_by(group_id=group_id, status='pending').all()
    return jsonify({"pendingUserIds": [i.invitee_id for i in pending]})

@friends_bp.route('/api/groups/<int:group_id>/image', methods=['POST'])
@jwt_required()
def update_group_image(group_id):
    user_id = get_jwt_identity()
    data = request.json
    image_url = data.get('image')
    
    membership = GroupMember.query.filter_by(group_id=group_id, user_id=user_id, role='admin').first()
    if not membership:
        return jsonify({"error": "Only admins can change group image"}), 403
        
    group = db.session.get(Group, group_id)
    group.image = image_url
    db.session.commit()
    
    return jsonify({"success": True})

@friends_bp.route('/api/calls/log', methods=['GET'])
@jwt_required()
def get_call_logs():
    user_id = get_jwt_identity()
    logs = CallLog.query.filter(
        (CallLog.user_id == user_id)
    ).order_by(CallLog.timestamp.desc()).all()
    
    result = []
    for l in logs:
        other = db.session.get(User, l.other_user_id) if l.other_user_id else None
        result.append({
            "id":          l.id,
            "otherId":     l.other_user_id,          # numeric â€” used by frontend for matching
            "name":        other.full_name if other else "Unknown",
            "image":       other.image if other else None,
            "type":        l.type,                   # 'outgoing' | 'incoming' | 'missed'
            "isVideo":     l.is_video,
            "timestamp":   l.timestamp.isoformat(),  # UTC, no Z suffix
            "duration":    l.duration or 0,
        })
    return jsonify(result)


# â”€â”€ WebRTC Call Signaling â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
import threading as _call_threading
import time as _time

_active_calls = {}      # call_id -> call dict
_pending_signals = {}   # user_id -> list of signal dicts
_calls_lock = _call_threading.Lock()

CALL_RING_TIMEOUT = 45   # seconds before unanswered call is auto-cancelled
STALE_CALL_TTL   = 120   # seconds before we purge dead calls

def _push_signal(user_id: int, signal: dict):
    with _calls_lock:
        _pending_signals.setdefault(user_id, []).append(signal)

def _log_call_entry(user_id: int, other_user_id: int, call_type: str, is_video: bool, duration: int):
    try:
        log = CallLog(
            user_id=user_id,
            other_user_id=other_user_id,
            type=call_type,
            is_video=is_video,
            duration=duration
        )
        db.session.add(log)
        db.session.commit()
    except Exception:
        db.session.rollback()

def _auto_cancel_call(call_id: str):
    """Background timer: if call not answered within timeout, cancel it."""
    import time
    time.sleep(CALL_RING_TIMEOUT)
    with _calls_lock:
        call = _active_calls.get(call_id)
        if not call or call['state'] != 'ringing':
            return
        _active_calls.pop(call_id, None)
        caller_id  = call['caller_id']
        callee_id  = call['callee_id']
        is_video   = call['is_video']
    # Push cancelled to caller, missed to callee
    _push_signal(caller_id,  {'type': 'call_no_answer', 'call_id': call_id})
    _push_signal(callee_id,  {'type': 'call_cancelled',  'call_id': call_id, 'caller_id': caller_id, 'is_video': is_video})
    _log_call_entry(caller_id,  callee_id,  'outgoing', is_video, 0)
    _log_call_entry(callee_id,  caller_id,  'missed',   is_video, 0)

def _parse_callee_id(raw):
    if raw is None:
        return None, "callee_id is required"
    if isinstance(raw, str):
        for prefix in ('u_', 'g_'):
            if raw.startswith(prefix):
                raw = raw[len(prefix):]
                break
    try:
        return int(raw), None
    except (ValueError, TypeError):
        return None, f"Invalid callee_id: {raw}"

@friends_bp.route('/api/calls/initiate', methods=['POST'])
@jwt_required()
def call_initiate():
    caller_id = int(get_jwt_identity())
    data = request.json or {}
    callee_id, err = _parse_callee_id(data.get('callee_id'))
    if err:
        return jsonify({"error": err}), 400

    is_video  = bool(data.get('is_video', False))
    call_id   = f"{caller_id}_{callee_id}_{int(datetime.now(timezone.utc).timestamp())}"

    caller = db.session.get(User, caller_id)
    callee = db.session.get(User, callee_id)
    if not caller or not callee:
        return jsonify({"error": "User not found"}), 404

    with _calls_lock:
        _active_calls[call_id] = {
            'caller_id':    caller_id,
            'callee_id':    callee_id,
            'is_video':     is_video,
            'state':        'ringing',
            'offer':        None,
            'answer':       None,
            'started_at':   _time.time(),
            'connected_at': None,
        }

    _push_signal(callee_id, {
        'type':         'incoming_call',
        'call_id':      call_id,
        'caller_id':    caller_id,
        'caller_name':  caller.full_name,
        'caller_image': caller.image,
        'is_video':     is_video,
        'expires_at':   _time.time() + CALL_RING_TIMEOUT,
    })

    # Auto-cancel after timeout in background
    t = _call_threading.Thread(target=_auto_cancel_call, args=(call_id,), daemon=True)
    t.start()

    return jsonify({"call_id": call_id, "status": "ringing"})


@friends_bp.route('/api/calls/offer', methods=['POST'])
@jwt_required()
def call_offer():
    caller_id = int(get_jwt_identity())
    data      = request.json or {}
    call_id   = data.get('call_id')
    offer     = data.get('offer')
    with _calls_lock:
        call = _active_calls.get(call_id)
        if not call or call['caller_id'] != caller_id:
            return jsonify({"error": "Call not found"}), 404
        call['offer'] = offer
    _push_signal(call['callee_id'], {'type': 'offer', 'call_id': call_id, 'offer': offer})
    return jsonify({"success": True})


@friends_bp.route('/api/calls/answer', methods=['POST'])
@jwt_required()
def call_answer():
    callee_id = int(get_jwt_identity())
    data      = request.json or {}
    call_id   = data.get('call_id')
    accepted  = bool(data.get('accepted', False))
    answer    = data.get('answer')

    with _calls_lock:
        call = _active_calls.get(call_id)
        if not call or call['callee_id'] != callee_id:
            return jsonify({"error": "Call not found or already ended"}), 404
        if accepted:
            call['state']        = 'connected'
            call['connected_at'] = _time.time()
            call['answer']       = answer
        else:
            call['state'] = 'rejected'

    caller_id = call['caller_id']
    if accepted:
        _push_signal(caller_id, {'type': 'answer', 'call_id': call_id, 'answer': answer})
    else:
        _push_signal(caller_id, {'type': 'call_rejected', 'call_id': call_id})
        _log_call_entry(caller_id, callee_id, 'outgoing', call['is_video'], 0)
        _log_call_entry(callee_id, caller_id, 'missed',   call['is_video'], 0)
        with _calls_lock:
            _active_calls.pop(call_id, None)
    return jsonify({"success": True})


@friends_bp.route('/api/calls/ice', methods=['POST'])
@jwt_required()
def call_ice():
    user_id  = int(get_jwt_identity())
    data     = request.json or {}
    call_id  = data.get('call_id')
    candidate = data.get('candidate')
    with _calls_lock:
        call = _active_calls.get(call_id)
        if not call:
            return jsonify({"error": "Call not found"}), 404
        remote_id = call['callee_id'] if call['caller_id'] == user_id else call['caller_id']
    _push_signal(remote_id, {'type': 'ice_candidate', 'call_id': call_id, 'candidate': candidate})
    return jsonify({"success": True})


@friends_bp.route('/api/calls/cancel', methods=['POST'])
@jwt_required()
def call_cancel():
    """Caller cancels before callee answers (e.g. caller hangs up during ringing)."""
    caller_id = int(get_jwt_identity())
    data      = request.json or {}
    call_id   = data.get('call_id')
    with _calls_lock:
        call = _active_calls.pop(call_id, None)
    if not call:
        return jsonify({"success": True})
    callee_id = call['callee_id']
    is_video  = call['is_video']
    _push_signal(callee_id, {'type': 'call_cancelled', 'call_id': call_id, 'caller_id': caller_id, 'is_video': is_video})
    _log_call_entry(caller_id, callee_id, 'outgoing', is_video, 0)
    _log_call_entry(callee_id, caller_id, 'missed',   is_video, 0)
    return jsonify({"success": True})


@friends_bp.route('/api/calls/end', methods=['POST'])
@jwt_required()
def call_end():
    user_id  = int(get_jwt_identity())
    data     = request.json or {}
    call_id  = data.get('call_id')
    with _calls_lock:
        call = _active_calls.pop(call_id, None)
    if not call:
        return jsonify({"success": True})

    caller_id = call['caller_id']
    callee_id = call['callee_id']
    is_video  = call['is_video']
    state     = call.get('state', 'ringing')

    # If still ringing (never connected) â€” caller cancelled, callee missed
    if state == 'ringing':
        remote_id = callee_id if user_id == caller_id else caller_id
        _push_signal(remote_id, {'type': 'call_cancelled', 'call_id': call_id, 'caller_id': caller_id, 'is_video': is_video})
        _log_call_entry(caller_id, callee_id, 'outgoing', is_video, 0)
        _log_call_entry(callee_id, caller_id, 'missed',   is_video, 0)
        return jsonify({"success": True, "duration": 0})

    # Connected call
    duration = 0
    if call.get('connected_at'):
        duration = max(0, int(_time.time() - call['connected_at']))

    remote_id = callee_id if user_id == caller_id else caller_id
    _push_signal(remote_id, {'type': 'call_ended', 'call_id': call_id, 'duration': duration})

    _log_call_entry(caller_id, callee_id, 'outgoing', is_video, duration)
    _log_call_entry(callee_id, caller_id, 'incoming' if duration > 0 else 'missed', is_video, duration)
    return jsonify({"success": True, "duration": duration})


@friends_bp.route('/api/calls/log/clear', methods=['DELETE'])
@jwt_required()
def clear_call_logs():
    """Dev utility: wipe all call logs for current user so bad old entries don't confuse."""
    user_id = int(get_jwt_identity())
    try:
        CallLog.query.filter_by(user_id=user_id).delete()
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        db.session.rollback()
        return jsonify({"error": str(e)}), 500


@friends_bp.route('/api/calls/signals', methods=['GET'])
@jwt_required()
def get_call_signals():
    user_id = int(get_jwt_identity())
    now     = _time.time()
    with _calls_lock:
        raw = _pending_signals.pop(user_id, [])
    # Filter out stale incoming_call signals (call already cancelled/expired)
    filtered = []
    for sig in raw:
        if sig.get('type') == 'incoming_call':
            call_id = sig.get('call_id')
            expires = sig.get('expires_at', 0)
            # Only deliver if call still active AND not expired
            with _calls_lock:
                still_active = call_id in _active_calls
            if still_active and now < expires:
                filtered.append(sig)
            # else: silently drop it â€” call was cancelled or timed out
        else:
            filtered.append(sig)
    return jsonify(filtered)

@friends_bp.route('/api/broadcast-lists', methods=['GET'])
@jwt_required()
def get_broadcast_lists():
    user_id = int(get_jwt_identity())

    # Broadcasts the user OWNS
    owned_lists = BroadcastList.query.filter_by(user_id=user_id).all()

    # Broadcasts the user is a MEMBER/RECIPIENT of
    member_rows = BroadcastRecipient.query.filter_by(recipient_id=user_id).all()
    member_list_ids = {r.list_id for r in member_rows}
    member_lists = [BroadcastList.query.get(lid) for lid in member_list_ids
                    if BroadcastList.query.get(lid)]

    # Merge, deduplicate
    all_lists = list({l.id: l for l in owned_lists + member_lists}.values())

    result = []
    for l in all_lists:
        recipients = BroadcastRecipient.query.filter_by(list_id=l.id).all()
        last_msg = Message.query.filter_by(broadcast_id=l.id).order_by(Message.timestamp.desc()).first()

        # Build rich recipient details
        recipient_details = []
        for r in recipients:
            u = db.session.get(User, r.recipient_id)
            if u:
                recipient_details.append({
                    'id': r.recipient_id,
                    'name': u.full_name,
                    'image': u.image,
                    'email': u.email,
                    'role': r.role if r.role else 'member',
                })

        # Last message preview text
        lmt = None
        if last_msg:
            if last_msg.is_deleted: lmt = 'This message was deleted'
            elif last_msg.type == 'image': lmt = 'ðŸ“· Photo'
            elif last_msg.type == 'video': lmt = 'ðŸŽ¥ Video'
            elif last_msg.type == 'voice': lmt = 'ðŸŽ™ï¸ Voice message'
            elif last_msg.type == 'file': lmt = 'ðŸ“Ž File'
            else: lmt = last_msg.text

        is_owner = (l.user_id == user_id)
        my_row = next((r for r in recipients if r.recipient_id == user_id), None)
        my_role = 'owner' if is_owner else (my_row.role if my_row and my_row.role else 'member')

        creator = db.session.get(User, l.user_id)

        result.append({
            'id': l.id,
            'name': l.name,
            'image': l.image if hasattr(l, 'image') else None,
            'description': l.description if hasattr(l, 'description') else None,
            'creatorId': l.user_id,
            'creatorName': creator.full_name if creator else 'Unknown',
            'creatorImage': creator.image if creator else None,
            'isOwner': is_owner,
            'myRole': my_role,
            'recipients': [r.recipient_id for r in recipients],
            'recipientDetails': recipient_details,
            'memberCount': len(recipients),
            'lastMessage': lmt or 'No messages yet',
            'lastMessageType': last_msg.type if last_msg else None,
            'lastUsed': (last_msg.timestamp if last_msg else l.last_used).isoformat() if (last_msg or l.last_used) else None,
            'createdAt': l.created_at.isoformat() if l.created_at else None,
        })
    return jsonify(result)

@friends_bp.route('/api/broadcast-lists', methods=['POST'])
@jwt_required()
def create_broadcast_list():
    user_id = int(get_jwt_identity())
    data = request.get_json()
    if not data or not data.get('name') or not data.get('recipients'):
        return jsonify({'error': 'Name and at least one recipient are required'}), 400

    new_list = BroadcastList(user_id=user_id, name=data['name'],
                            description=data.get('description', ''))
    db.session.add(new_list)
    db.session.flush()

    for item in data['recipients']:
        # Support {id, role} objects OR plain ints/strings
        if isinstance(item, dict):
            raw_id = item.get('id', item.get('recipient_id', 0))
            role = item.get('role', 'member')
        else:
            raw_id = item
            role = 'member'
        try:
            clean_id = int(str(raw_id).replace('u_', '').replace('g_', ''))
        except:
            continue
        if clean_id == user_id:
            continue  # never add owner as recipient
        db.session.add(BroadcastRecipient(list_id=new_list.id, recipient_id=clean_id, role=role))

    db.session.commit()
    return jsonify({'id': new_list.id, 'name': new_list.name}), 201

@friends_bp.route('/api/broadcast-lists/<int:list_id>', methods=['DELETE'])
@jwt_required()
def delete_broadcast_list(list_id):
    user_id = int(get_jwt_identity())
    broadcast_list = BroadcastList.query.filter_by(id=list_id, user_id=user_id).first()
    
    if not broadcast_list:
        return jsonify({"error": "Broadcast list not found"}), 404
        
    # Delete recipients first
    BroadcastRecipient.query.filter_by(list_id=list_id).delete()
    db.session.delete(broadcast_list)
    db.session.commit()
    
    return jsonify({"message": "Broadcast list deleted successfully"}), 200

@friends_bp.route('/api/friends/requests', methods=['GET'])
@jwt_required()
def get_requests():
    user_id = get_jwt_identity()
    incoming = FriendRequest.query.filter_by(receiver_id=user_id, status='pending').all()
    outgoing = FriendRequest.query.filter_by(sender_id=user_id, status='pending').all()
    
    return jsonify({
        "incoming": [{
            "id": r.id,
            "sender": {
                "id": r.sender_id,
                "name": db.session.get(User, r.sender_id).full_name,
                "email": db.session.get(User, r.sender_id).email,
                "image": db.session.get(User, r.sender_id).image
            },
            "created_at": r.created_at.isoformat()
        } for r in incoming],
        "outgoing": [{
            "id": r.id,
            "receiver": {
                "id": r.receiver_id,
                "name": db.session.get(User, r.receiver_id).full_name,
                "email": db.session.get(User, r.receiver_id).email,
                "image": db.session.get(User, r.receiver_id).image
            },
            "created_at": r.created_at.isoformat()
        } for r in outgoing]
    })

@friends_bp.route('/api/friends/request', methods=['POST'])
@jwt_required()
def send_request():
    try:
        user_id = int(get_jwt_identity())
        data = request.json
        receiver_id = data.get('receiver_id')
        
        if not receiver_id:
            return jsonify({"error": "Missing receiver_id"}), 400
            
        receiver_id = int(receiver_id)
        
        if receiver_id == user_id:
            return jsonify({"error": "Cannot send request to yourself"}), 400
            
        logger.info(f"Sending friend request from {user_id} to {receiver_id}")
            
        # Check if already friends (ignore soft-deleted friendships)
        existing_friendship = Friendship.query.filter(
            ((Friendship.user_id == user_id) & (Friendship.friend_id == receiver_id)) |
            ((Friendship.user_id == receiver_id) & (Friendship.friend_id == user_id)),
            Friendship.is_deleted == False,
            Friendship.is_blocked == False
        ).first()
        
        if existing_friendship:
            return jsonify({"error": "Already friends"}), 400
            
        # Check if already pending
        existing_request = FriendRequest.query.filter_by(
            sender_id=user_id, receiver_id=receiver_id, status='pending'
        ).first()
        
        if existing_request:
            return jsonify({"error": "Request already pending"}), 400

        # Also check reverse pending request (they already sent one to us)
        reverse_request = FriendRequest.query.filter_by(
            sender_id=receiver_id, receiver_id=user_id, status='pending'
        ).first()
        if reverse_request:
            return jsonify({"error": "This user already sent you a friend request"}), 400
            
        new_request = FriendRequest(sender_id=user_id, receiver_id=receiver_id)
        db.session.add(new_request)
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        logger.error(f"Error in send_request: {e}")
        return jsonify({"error": str(e)}), 500

@friends_bp.route('/api/friends/request/accept', methods=['POST'])
@jwt_required()
def accept_request():
    try:
        user_id = int(get_jwt_identity())
        data = request.json
        request_id = data.get('request_id')
        
        logger.info(f"Accepting friend request: {request_id} for user: {user_id}")
        
        req = db.session.get(FriendRequest, request_id)
        if not req:
            logger.warning(f"Request {request_id} not found")
            return jsonify({"error": "Request not found"}), 404
            
        if req.receiver_id != user_id:
            logger.warning(f"User {user_id} is not the receiver {req.receiver_id} of request {request_id}")
            return jsonify({"error": "Unauthorized"}), 404
            
        if req.status != 'pending':
            logger.warning(f"Request {request_id} is already {req.status}")
            return jsonify({"error": "Request already processed"}), 400
            
        req.status = 'accepted'
        
        # Check if a friendship row already exists (could be soft-deleted or blocked)
        existing = Friendship.query.filter(
            ((Friendship.user_id == req.sender_id) & (Friendship.friend_id == req.receiver_id)) |
            ((Friendship.user_id == req.receiver_id) & (Friendship.friend_id == req.sender_id))
        ).first()
        
        if existing:
            # Reactivate soft-deleted or blocked friendship
            existing.is_deleted = False
            existing.is_blocked = False
            existing.blocked_by_id = None
        else:
            friendship = Friendship(user_id=req.sender_id, friend_id=req.receiver_id)
            db.session.add(friendship)
            
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        logger.error(f"Error in accept_request: {e}")
        return jsonify({"error": str(e)}), 500

@friends_bp.route('/api/friends/request/decline', methods=['POST'])
@jwt_required()
def decline_request():
    try:
        user_id = int(get_jwt_identity())
        data = request.json
        request_id = data.get('request_id')
        
        logger.info(f"Declining friend request: {request_id} for user: {user_id}")
        
        req = db.session.get(FriendRequest, request_id)
        if not req:
            return jsonify({"error": "Request not found"}), 404
            
        if req.receiver_id != user_id:
            return jsonify({"error": "Unauthorized"}), 404
            
        if req.status != 'pending':
            return jsonify({"error": "Request already processed"}), 400
            
        req.status = 'rejected'
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        logger.error(f"Error in decline_request: {e}")
        return jsonify({"error": str(e)}), 500

@friends_bp.route('/api/friends/request/cancel', methods=['POST'])
@jwt_required()
def cancel_request():
    try:
        user_id = int(get_jwt_identity())
        data = request.json
        request_id = data.get('request_id')
        
        logger.info(f"Canceling friend request: {request_id} for user: {user_id}")
        
        req = db.session.get(FriendRequest, request_id)
        if not req:
            return jsonify({"error": "Request not found"}), 404
            
        if req.sender_id != user_id:
            return jsonify({"error": "Unauthorized"}), 404
            
        if req.status != 'pending':
            return jsonify({"error": "Only pending requests can be cancelled"}), 400
            
        db.session.delete(req)
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        logger.error(f"Error in cancel_request: {e}")
        return jsonify({"error": str(e)}), 500

@friends_bp.route('/api/users/search', methods=['GET'])
@jwt_required()
def search_users():
    try:
        user_id = int(get_jwt_identity())
        q = request.args.get('q', '')
        if not q or len(q) < 3:
            return jsonify([])
        
        users = User.query.filter(
            (User.email.ilike(f"%{q}%")) | (User.full_name.ilike(f"%{q}%"))
        ).limit(10).all()
        
        result = []
        for u in users:
            if u.id == user_id:
                continue
                
            # Check if already friends or requested
            friendship = Friendship.query.filter(
                ((Friendship.user_id == user_id) & (Friendship.friend_id == u.id)) |
                ((Friendship.user_id == u.id) & (Friendship.friend_id == user_id))
            ).first()
            
            request_pending = FriendRequest.query.filter_by(
                sender_id=user_id, receiver_id=u.id, status='pending'
            ).first() is not None
            
            result.append({
                "id": u.id,
                "name": u.full_name,
                "email": u.email,
                "image": u.image,
                "isFriend": friendship is not None,
                "requestPending": request_pending,
                "bio": u.bio
            })
        return jsonify(result)
    except Exception as e:
        logger.error(f"Error in search_users: {e}")
        return jsonify({"error": str(e)}), 500

@friends_bp.route('/api/messages/<int:other_user_id>', methods=['GET'])
@jwt_required()
def get_messages(other_user_id):
    user_id = int(get_jwt_identity())
    messages = Message.query.filter(
        ((Message.sender_id == user_id) & (Message.receiver_id == other_user_id)) |
        ((Message.sender_id == other_user_id) & (Message.receiver_id == user_id))
    ).order_by(Message.timestamp.asc()).all()
    
    # Mark as read
    Message.query.filter_by(sender_id=other_user_id, receiver_id=user_id, is_read=False).update({"is_read": True})
    db.session.commit()
    
    result = []
    for m in messages:
        try:
            deleted_for = json.loads(m.deleted_for) if m.deleted_for else []
            starred_by = json.loads(m.is_starred_by) if m.is_starred_by else []
        except:
            deleted_for = []
            starred_by = []
            
        if user_id in deleted_for:
            continue
            
        # Get status
        status = 'sent'
        if m.sender_id == user_id:
            receipt = MessageReceipt.query.filter_by(message_id=m.id).first()
            if receipt:
                if receipt.read_at: status = 'read'
                elif receipt.delivered_at: status = 'delivered'
        else:
            if m.is_read: status = 'read'

        result.append({
            "id": m.id,
            "text": "This message was deleted" if m.is_deleted else m.text,
            "sender": "me" if m.sender_id == user_id else "other",
            "time": m.timestamp.isoformat() + 'Z',
            "type": m.type,
            "mediaUrl": m.media_url,
            "isRead": m.is_read,
            "status": status,
            "isDeleted": m.is_deleted,
            "isStarred": user_id in starred_by,
            "reaction": m.reaction,
            "reply_to_id": m.reply_to_id,
            "replyTo": {
                "id": rm.id,
                "text": rm.text,
                "sender": "me" if rm.sender_id == user_id else "other",
                "type": rm.type,
                "mediaUrl": rm.media_url
            } if m.reply_to_id and (rm := db.session.get(Message, m.reply_to_id)) else None
        })
        
    return jsonify(result)

@friends_bp.route('/api/messages/group/<int:group_id>', methods=['GET'])
@jwt_required()
def get_group_messages(group_id):
    user_id = int(get_jwt_identity())
    # Check if user is a member
    membership = GroupMember.query.filter_by(group_id=group_id, user_id=user_id).first()
    if not membership:
        return jsonify({"error": "Not a member of this group"}), 403
        
    messages = Message.query.filter_by(group_id=group_id).order_by(Message.timestamp.asc()).all()
    
    result = []
    for m in messages:
        try:
            deleted_for = json.loads(m.deleted_for) if m.deleted_for else []
            starred_by = json.loads(m.is_starred_by) if m.is_starred_by else []
        except:
            deleted_for = []
            starred_by = []
            
        if user_id in deleted_for:
            continue
            
        # Get status
        status = 'sent'
        if m.sender_id == user_id:
            all_receipts = MessageReceipt.query.filter_by(message_id=m.id).all()
            if all_receipts:
                if all(r.read_at for r in all_receipts): status = 'read'
                elif all(r.delivered_at for r in all_receipts): status = 'delivered'

        result.append({
            "id": m.id,
            "sender_id": m.sender_id,
            "sender": "me" if m.sender_id == user_id else "other",
            "text": "This message was deleted" if m.is_deleted else m.text,
            "type": m.type,
            "mediaUrl": m.media_url,
            "time": m.timestamp.isoformat() + 'Z',
            "isRead": m.is_read,
            "status": status,
            "isDeleted": m.is_deleted,
            "isStarred": user_id in starred_by,
            "reaction": m.reaction,
            "reply_to_id": m.reply_to_id,
            "replyTo": {
                "id": rm.id,
                "text": rm.text,
                "sender": "me" if rm.sender_id == user_id else "other",
                "type": rm.type,
                "mediaUrl": rm.media_url
            } if m.reply_to_id and (rm := db.session.get(Message, m.reply_to_id)) else None
        })
        
    return jsonify(result)

@friends_bp.route('/api/messages/broadcast/<int:list_id>', methods=['GET'])
@jwt_required()
def get_broadcast_messages(list_id):
    user_id = int(get_jwt_identity())

    # Allow owner OR any recipient member to read messages
    broadcast_list = BroadcastList.query.get(list_id)
    if not broadcast_list:
        return jsonify({'error': 'Broadcast list not found'}), 404

    is_owner = broadcast_list.user_id == user_id
    is_member = BroadcastRecipient.query.filter_by(list_id=list_id, recipient_id=user_id).first() is not None
    if not is_owner and not is_member:
        return jsonify({'error': 'Not authorized'}), 403

    messages = Message.query.filter_by(broadcast_id=list_id).order_by(Message.timestamp.asc()).all()

    result = []
    for m in messages:
        try:
            deleted_for = json.loads(m.deleted_for) if m.deleted_for else []
            starred_by = json.loads(m.is_starred_by) if m.is_starred_by else []
        except:
            deleted_for = []
            starred_by = []
        if user_id in deleted_for:
            continue

        status = 'sent'
        if m.sender_id == user_id:
            all_receipts = MessageReceipt.query.filter_by(message_id=m.id).all()
            if all_receipts:
                if all(r.read_at for r in all_receipts): status = 'read'
                elif all(r.delivered_at for r in all_receipts): status = 'delivered'

        # Get sender name for non-owner senders (admins can also send)
        sender_user = db.session.get(User, m.sender_id)

        result.append({
            'id': m.id,
            'sender_id': m.sender_id,
            'senderName': sender_user.full_name if sender_user else 'Unknown',
            'senderImage': sender_user.image if sender_user else None,
            'sender': 'me' if m.sender_id == user_id else 'other',
            'text': 'This message was deleted' if m.is_deleted else m.text,
            'type': m.type,
            'mediaUrl': m.media_url,
            'time': m.timestamp.isoformat() + 'Z',
            'isRead': m.is_read,
            'status': status,
            'isDeleted': m.is_deleted,
            'isStarred': user_id in starred_by,
            'reaction': m.reaction,
            'reply_to_id': m.reply_to_id,
            'replyTo': {
                "id": rm.id,
                "text": rm.text,
                "sender": "me" if rm.sender_id == user_id else "other",
                "type": rm.type,
                "mediaUrl": rm.media_url
            } if m.reply_to_id and (rm := db.session.get(Message, m.reply_to_id)) else None
        })

    return jsonify(result)

@friends_bp.route('/api/messages/react', methods=['POST'])
@jwt_required()
def react_to_message():
    user_id = int(get_jwt_identity())
    data = request.json or {}
    message_id = data.get('message_id')
    emoji = data.get('emoji')  # None means remove reaction

    if not message_id:
        return jsonify({'error': 'message_id is required'}), 400

    msg = db.session.get(Message, message_id)
    if not msg:
        return jsonify({'error': 'Message not found'}), 404

    # Toggle: if same emoji is sent again, remove it
    if msg.reaction == emoji:
        msg.reaction = None
    else:
        msg.reaction = emoji

    db.session.commit()
    return jsonify({'success': True, 'reaction': msg.reaction})

@friends_bp.route('/api/messages/send', methods=['POST'])
@jwt_required()
def send_message():
    user_id = int(get_jwt_identity())
    data = request.json
    text = data.get('text')
    receiver_id = data.get('receiver_id')
    group_id = data.get('group_id')
    broadcast_id = data.get('broadcast_id')
    msg_type = data.get('type', 'text')
    media_url = data.get('mediaUrl')
    reply_to_id = data.get('reply_to_id')
    
    if not any([receiver_id, group_id, broadcast_id]):
        return jsonify({"error": "No recipient specified"}), 400
        
    if broadcast_id:
        # Check sender is owner OR admin member
        bl = BroadcastList.query.get(broadcast_id)
        if not bl:
            return jsonify({'error': 'Broadcast list not found'}), 404
        is_owner = bl.user_id == user_id
        sender_row = BroadcastRecipient.query.filter_by(list_id=broadcast_id, recipient_id=user_id).first()
        is_admin = sender_row and sender_row.role == 'admin'
        if not is_owner and not is_admin:
            return jsonify({'error': 'Only the owner or admins can send broadcast messages'}), 403

        # Store ONE message linked to the broadcast â€” NO individual DM copies
        broadcast_msg = Message(
            sender_id=user_id,
            broadcast_id=broadcast_id,
            text=text,
            type=msg_type,
            media_url=media_url,
            reply_to_id=reply_to_id
        )
        db.session.add(broadcast_msg)
        bl.last_used = datetime.now(timezone.utc)
        db.session.commit()

        sender_user = db.session.get(User, user_id)
        return jsonify({
            'id': broadcast_msg.id,
            'sender': 'me',
            'sender_id': user_id,
            'senderName': sender_user.full_name if sender_user else 'Unknown',
            'senderImage': sender_user.image if sender_user else None,
            'text': broadcast_msg.text,
            'type': broadcast_msg.type,
            'mediaUrl': broadcast_msg.media_url,
            'time': broadcast_msg.timestamp.isoformat() + 'Z',
            'status': 'sent',
            'reply_to_id': broadcast_msg.reply_to_id,
            'replyTo': {
                "id": rm.id,
                "text": rm.text,
                "sender": "me" if rm.sender_id == user_id else "other",
                "type": rm.type,
                "mediaUrl": rm.media_url
            } if broadcast_msg.reply_to_id and (rm := db.session.get(Message, broadcast_msg.reply_to_id)) else None
        })

    # Standard send
    msg = Message(
        sender_id=user_id,
        receiver_id=receiver_id,
        group_id=group_id,
        text=text,
        type=msg_type,
        media_url=media_url,
        reply_to_id=reply_to_id
    )
    db.session.add(msg)
    db.session.flush() # Get msg.id

    # Create receipts
    if receiver_id:
        receipt = MessageReceipt(message_id=msg.id, user_id=receiver_id)
        db.session.add(receipt)
    elif group_id:
        members = GroupMember.query.filter_by(group_id=group_id, is_exited=False).all()
        for member in members:
            if member.user_id != user_id:
                receipt = MessageReceipt(message_id=msg.id, user_id=member.user_id)
                db.session.add(receipt)

    db.session.commit()
    
    return jsonify({
        "id": msg.id,
        "text": msg.text,
        "sender": "me",
        "time": msg.timestamp.isoformat() + 'Z',
        "type": msg.type,
        "mediaUrl": msg.media_url,
        "reply_to_id": msg.reply_to_id,
        "replyTo": {
            "id": rm.id,
            "text": rm.text,
            "sender": "me" if rm.sender_id == user_id else "other",
            "type": rm.type,
            "mediaUrl": rm.media_url
        } if msg.reply_to_id and (rm := db.session.get(Message, msg.reply_to_id)) else None,
        "status": 'sent'
    })

@friends_bp.route('/api/messages/receipt', methods=['POST'])
@jwt_required()
def update_message_receipt():
    user_id = int(get_jwt_identity())
    data = request.json
    message_id = data.get('message_id')
    status_type = data.get('type') # 'delivered', 'read', 'played'
    
    if not message_id or not status_type:
        return jsonify({"error": "Missing parameters"}), 400
        
    receipt = MessageReceipt.query.filter_by(message_id=message_id, user_id=user_id).first()
    
    # If it's a DM and the receipt doesn't exist yet (maybe legacy or broadcast)
    if not receipt:
        msg = db.session.get(Message, message_id)
        if msg and (msg.receiver_id == user_id or (msg.group_id and GroupMember.query.filter_by(group_id=msg.group_id, user_id=user_id).first())):
            receipt = MessageReceipt(message_id=message_id, user_id=user_id)
            db.session.add(receipt)
        else:
            return jsonify({"error": "Receipt not found or unauthorized"}), 404

    now = datetime.now(timezone.utc)
    if status_type == 'delivered' and not receipt.delivered_at:
        receipt.delivered_at = now
    elif status_type == 'read':
        if not receipt.delivered_at: receipt.delivered_at = now
        if not receipt.read_at: receipt.read_at = now
        # Also sync legacy is_read
        msg = db.session.get(Message, message_id)
        if msg and msg.receiver_id == user_id:
            msg.is_read = True
    elif status_type == 'played' and not receipt.played_at:
        if not receipt.delivered_at: receipt.delivered_at = now
        if not receipt.read_at: receipt.read_at = now
        receipt.played_at = now
        
    db.session.commit()
    return jsonify({"success": True})

@friends_bp.route('/api/messages/<int:msg_id>/info', methods=['GET'])
@jwt_required()
def get_message_info(msg_id):
    user_id = int(get_jwt_identity())
    msg = db.session.get(Message, msg_id)
    if not msg or msg.sender_id != user_id:
        return jsonify({"error": "Message not found or unauthorized"}), 404
        
    receipts = MessageReceipt.query.filter_by(message_id=msg_id).all()
    
    info = []
    for r in receipts:
        u = db.session.get(User, r.user_id)
        if u:
            # Respect receiver's read receipt privacy setting
            receipts_enabled = getattr(u, 'privacy_read_receipts', True)
            info.append({
                "user_id":      u.id,
                "name":         u.full_name,
                "image":        u.image,
                "delivered_at": r.delivered_at.isoformat() + 'Z' if r.delivered_at else None,
                "read_at":      (r.read_at.isoformat() + 'Z' if r.read_at else None) if receipts_enabled else None,
                "played_at":    (r.played_at.isoformat() + 'Z' if r.played_at else None) if receipts_enabled else None,
            })
            
    return jsonify({
        "id":       msg.id,
        "text":     msg.text,
        "type":     msg.type,
        "time":     msg.timestamp.isoformat() + 'Z',
        "receipts": info
    })

# --- NEW DATA PERSISTENCE ENDPOINTS ---

@friends_bp.route('/api/friends/<int:friend_id>/preferences', methods=['POST'])
@jwt_required()
def update_friend_preferences(friend_id):
    user_id = int(get_jwt_identity())
    data = request.json
    friendship = Friendship.query.filter(
        ((Friendship.user_id == user_id) & (Friendship.friend_id == friend_id)) |
        ((Friendship.user_id == friend_id) & (Friendship.friend_id == user_id))
    ).first()
    
    if not friendship:
        return jsonify({"error": "Friendship not found"}), 404
        
    if 'isPinned' in data:
        friendship.is_pinned = data['isPinned']
    if 'isArchived' in data:
        friendship.is_archived = data['isArchived']
    if 'isFavourite' in data:
        friendship.is_favourite = data['isFavourite']
    if 'isMuted' in data:
        friendship.is_muted = data['isMuted']
        
    db.session.commit()
    return jsonify({"success": True})

@friends_bp.route('/api/groups/<int:group_id>/preferences', methods=['POST'])
@jwt_required()
def update_group_preferences(group_id):
    user_id = int(get_jwt_identity())
    data = request.json
    membership = GroupMember.query.filter_by(group_id=group_id, user_id=user_id).first()
    
    if not membership:
        return jsonify({"error": "Group membership not found"}), 404
        
    if 'isPinned' in data:
        membership.is_pinned = data['isPinned']
    if 'isArchived' in data:
        membership.is_archived = data['isArchived']
    if 'isFavourite' in data:
        membership.is_favourite = data['isFavourite']
    if 'isMuted' in data:
        membership.is_muted = data['isMuted']
        
    db.session.commit()
    return jsonify({"success": True})

@friends_bp.route('/api/groups/<int:group_id>/exit', methods=['POST'])
@jwt_required()
def exit_group(group_id):
    user_id = int(get_jwt_identity())
    membership = GroupMember.query.filter_by(group_id=group_id, user_id=user_id).first()
    
    if not membership:
        return jsonify({"error": "Not a member"}), 404
        
    membership.is_exited = True
    db.session.commit()
    return jsonify({"success": True})

@friends_bp.route('/api/messages/<int:msg_id>/delete', methods=['POST'])
@jwt_required()
def delete_message(msg_id):
    user_id = int(get_jwt_identity())
    data = request.json
    del_type = data.get('type', 'me')  # 'me' or 'everyone'
    
    msg = db.session.get(Message, msg_id)
    if not msg:
        return jsonify({"error": "Message not found"}), 404
        
    if del_type == 'everyone':
        if msg.sender_id != user_id:
            return jsonify({"error": "Cannot delete others' messages for everyone"}), 403
        msg.is_deleted = True
    else:
        try:
            deleted_for = json.loads(msg.deleted_for) if msg.deleted_for else []
        except:
            deleted_for = []
        if user_id not in deleted_for:
            deleted_for.append(user_id)
            msg.deleted_for = json.dumps(deleted_for)
            
    db.session.commit()
    return jsonify({"success": True})

@friends_bp.route('/api/messages/clear', methods=['POST'])
@jwt_required()
def clear_chat():
    user_id = int(get_jwt_identity())
    data = request.json
    chat_id = data.get('chatId')
    
    if not chat_id:
        return jsonify({"error": "No chat specified"}), 400
        
    # parse chatId format: "u_##", "g_##", or "b_##"
    if chat_id.startswith('u_'):
        other_id = int(chat_id[2:])
        messages = Message.query.filter(
            ((Message.sender_id == user_id) & (Message.receiver_id == other_id)) |
            ((Message.sender_id == other_id) & (Message.receiver_id == user_id))
        ).all()
    elif chat_id.startswith('g_'):
        g_id = int(chat_id[2:])
        messages = Message.query.filter_by(group_id=g_id).all()
    elif chat_id.startswith('b_'):
        b_id = int(chat_id[2:])
        messages = Message.query.filter_by(broadcast_id=b_id).all()
    else:
        return jsonify({"error": "Invalid chat format"}), 400
        
    for msg in messages:
        try:
            deleted = json.loads(msg.deleted_for) if msg.deleted_for else []
        except:
            deleted = []
        if user_id not in deleted:
            deleted.append(user_id)
            msg.deleted_for = json.dumps(deleted)
            
    db.session.commit()
    return jsonify({"success": True})

@friends_bp.route('/api/messages/<int:msg_id>/star', methods=['POST'])
@jwt_required()
def toggle_star(msg_id):
    user_id = int(get_jwt_identity())
    msg = db.session.get(Message, msg_id)
    if not msg:
        return jsonify({"error": "Message not found"}), 404
        
    try:
        starred = json.loads(msg.is_starred_by) if msg.is_starred_by else []
    except:
        starred = []
        
    if user_id in starred:
        starred.remove(user_id)
    else:
        starred.append(user_id)
        
    msg.is_starred_by = json.dumps(starred)
    db.session.commit()
    return jsonify({"success": True})

# â”€â”€ DELETE /api/friends/<id> â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
@friends_bp.route('/api/friends/<int:friend_id>', methods=['DELETE'])
@jwt_required()
def delete_friend(friend_id):
    """Soft-delete a friendship (removes from friend list on both sides)."""
    user_id = int(get_jwt_identity())
    friendship = Friendship.query.filter(
        ((Friendship.user_id == user_id) & (Friendship.friend_id == friend_id)) |
        ((Friendship.user_id == friend_id) & (Friendship.friend_id == user_id))
    ).first()
    if friendship:
        friendship.is_deleted = True
        # Also clear block state so it doesn't linger
        friendship.is_blocked = False
        friendship.blocked_by_id = None
        db.session.commit()
    return jsonify({"success": True})

# â”€â”€ POST /api/friends/block â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
@friends_bp.route('/api/friends/block', methods=['POST'])
@jwt_required()
def block_user():
    """Block a user. Creates a friendship record if one does not exist."""
    user_id = int(get_jwt_identity())
    data = request.json
    target_id = data.get('user_id')

    if not target_id:
        return jsonify({"error": "user_id is required"}), 400

    target_id = int(target_id)

    friendship = Friendship.query.filter(
        ((Friendship.user_id == user_id) & (Friendship.friend_id == target_id)) |
        ((Friendship.user_id == target_id) & (Friendship.friend_id == user_id))
    ).first()

    if not friendship:
        friendship = Friendship(
            user_id=user_id,
            friend_id=target_id,
            is_blocked=True,
            blocked_by_id=user_id
        )
        db.session.add(friendship)
    else:
        friendship.is_deleted = False
        friendship.is_blocked = True
        friendship.blocked_by_id = user_id

    db.session.commit()
    return jsonify({"success": True})

# â”€â”€ POST /api/friends/unblock â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
@friends_bp.route('/api/friends/unblock', methods=['POST'])
@jwt_required()
def unblock_user():
    """Unblock a previously blocked user."""
    user_id = int(get_jwt_identity())
    data = request.json
    target_id = data.get('user_id')

    if not target_id:
        return jsonify({"error": "user_id is required"}), 400

    target_id = int(target_id)

    friendship = Friendship.query.filter(
        ((Friendship.user_id == user_id) & (Friendship.friend_id == target_id)) |
        ((Friendship.user_id == target_id) & (Friendship.friend_id == user_id))
    ).first()

    if not friendship:
        # Nothing to unblock - treat as success
        return jsonify({"success": True})

    # If not blocked at all, nothing to do
    if not friendship.is_blocked:
        return jsonify({"success": True, "message": "User is not blocked"})

    # Allow unblock if this user was the blocker, or if blocked_by_id is NULL (legacy rows)
    blocked_by = friendship.blocked_by_id
    if blocked_by is not None and int(blocked_by) != user_id:
        return jsonify({"error": "You did not block this user"}), 403

    friendship.is_blocked = False
    friendship.blocked_by_id = None
    db.session.commit()
    logger.info(f"User {user_id} unblocked user {target_id}")
    return jsonify({"success": True})

# â”€â”€ DELETE /api/messages/<other_user_id>/clear â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
@friends_bp.route('/api/messages/<int:other_user_id>/clear', methods=['DELETE'])
@jwt_required()
def clear_chat_delete(other_user_id):
    """Delete all messages between current user and another user."""
    user_id = int(get_jwt_identity())
    Message.query.filter(
        ((Message.sender_id == user_id) & (Message.receiver_id == other_user_id)) |
        ((Message.sender_id == other_user_id) & (Message.receiver_id == user_id))
    ).delete(synchronize_session=False)
    db.session.commit()
    return jsonify({"success": True})

# â”€â”€ DELETE /api/messages/group/<group_id>/clear â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
@friends_bp.route('/api/messages/group/<int:group_id>/clear', methods=['DELETE'])
@jwt_required()
def clear_group_chat_delete(group_id):
    """Delete all messages in a group chat."""
    user_id = int(get_jwt_identity())
    membership = GroupMember.query.filter_by(group_id=group_id, user_id=user_id).first()
    if not membership:
        return jsonify({"error": "Not a member of this group"}), 403
    Message.query.filter_by(group_id=group_id).delete(synchronize_session=False)
    db.session.commit()
    return jsonify({"success": True})

# â”€â”€â”€ REAL-TIME PRESENCE SYSTEM â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
# In-memory store: { user_id: { 'last_beat': datetime, 'typing_to': int|None,
#                               'typing_at': datetime|None, 'activity': str } }
# 'activity' can be 'online' | 'typing' | 'recording'
import threading
_presence_lock = threading.Lock()
_presence: dict = {}          # user_id -> dict
_ONLINE_TIMEOUT  = 35         # seconds without heartbeat â†’ offline
_TYPING_TIMEOUT  = 8          # seconds without typing ping â†’ stop showing typing

def _user_is_online(uid: int) -> bool:
    p = _presence.get(uid)
    if not p: return False
    return (datetime.now(timezone.utc) - p['last_beat']).total_seconds() < _ONLINE_TIMEOUT

def _user_activity(uid: int) -> str:
    """Return 'recording' | 'typing' | 'online' | 'offline'."""
    if not _user_is_online(uid):
        return 'offline'
    p = _presence[uid]
    if p.get('activity') in ('typing', 'recording') and p.get('typing_at'):
        age = (datetime.now(timezone.utc) - p['typing_at']).total_seconds()
        if age < _TYPING_TIMEOUT:
            return p['activity']
    return 'online'

@friends_bp.route('/api/presence/heartbeat', methods=['POST'])
@jwt_required()
def presence_heartbeat():
    """Client calls every 20 s to signal the tab is open."""
    user_id = int(get_jwt_identity())
    now = datetime.now(timezone.utc)
    with _presence_lock:
        p = _presence.setdefault(user_id, {})
        p['last_beat'] = now
        if 'activity' not in p:
            p['activity'] = 'online'
    # Persist online status + last_seen to DB (non-blocking best-effort)
    try:
        user = db.session.get(User, user_id)
        if user:
            user.status    = 'online'
            user.last_seen = now
            db.session.commit()
    except Exception:
        db.session.rollback()
    return jsonify({'ok': True})

@friends_bp.route('/api/presence/typing', methods=['POST'])
@jwt_required()
def presence_typing():
    """Client calls while user is typing or recording a voice note."""
    user_id = int(get_jwt_identity())
    data     = request.json or {}
    activity = data.get('activity', 'typing')   # 'typing' | 'recording'
    now      = datetime.now(timezone.utc)
    with _presence_lock:
        p = _presence.setdefault(user_id, {})
        p['last_beat'] = now
        p['typing_at'] = now
        p['activity']  = activity
        p['typing_to'] = data.get('to_user_id')  # friend they're typing to
    return jsonify({'ok': True})

@friends_bp.route('/api/presence/offline', methods=['POST'])
@jwt_required()
def presence_offline():
    """Client calls on logout / beforeunload."""
    user_id = int(get_jwt_identity())
    now     = datetime.now(timezone.utc)
    with _presence_lock:
        _presence.pop(user_id, None)
    try:
        user = db.session.get(User, user_id)
        if user:
            user.status    = 'offline'
            user.last_seen = now
            db.session.commit()
    except Exception:
        db.session.rollback()
    return jsonify({'ok': True})

@friends_bp.route('/api/presence/statuses', methods=['GET'])
@jwt_required()
def presence_statuses():
    """Return presence info for all friends of this user (fast in-memory read)."""
    user_id = int(get_jwt_identity())
    # Collect friend IDs
    friendships = Friendship.query.filter(
        ((Friendship.user_id == user_id) | (Friendship.friend_id == user_id)),
        Friendship.is_deleted == False
    ).all()
    friend_ids = []
    for f in friendships:
        other = f.friend_id if f.user_id == user_id else f.user_id
        friend_ids.append(other)

    # Group members
    memberships = GroupMember.query.filter_by(user_id=user_id, is_exited=False).all()
    for m in memberships:
        others = GroupMember.query.filter(
            GroupMember.group_id == m.group_id,
            GroupMember.user_id != user_id
        ).all()
        friend_ids.extend(o.user_id for o in others)
    friend_ids = list(set(friend_ids))

    result = {}
    with _presence_lock:
        for fid in friend_ids:
            activity = _user_activity(fid)
            p        = _presence.get(fid, {})
            # Check if this friend is typing TO this user
            typing_to = p.get('typing_to')
            is_typing_to_me = (typing_to == user_id) and activity in ('typing', 'recording')
            if activity == 'offline':
                # Fall back to DB last_seen
                u = db.session.get(User, fid)
                last_seen = u.last_seen.isoformat() + 'Z' if u and u.last_seen else None
            else:
                last_seen = datetime.now(timezone.utc).isoformat() + 'Z'
            result[str(fid)] = {
                'status':         'online' if activity != 'offline' else 'offline',
                'activity':       activity,          # 'online'|'typing'|'recording'|'offline'
                'typing_to_me':   is_typing_to_me,
                'last_seen':      last_seen,
            }
    return jsonify(result)

@friends_bp.route('/api/friends/unread-counts', methods=['GET'])
@jwt_required()
def get_unread_counts():
    """Lightweight endpoint â€” returns {friend_id: unread_count} map."""
    user_id = int(get_jwt_identity())
    counts  = {}
    # DM unread
    friendships = Friendship.query.filter(
        ((Friendship.user_id == user_id) | (Friendship.friend_id == user_id)),
        Friendship.is_deleted == False
    ).all()
    for f in friendships:
        other_id = f.friend_id if f.user_id == user_id else f.user_id
        count = Message.query.filter_by(
            sender_id=other_id, receiver_id=user_id, is_read=False
        ).count()
        counts[f'u_{other_id}'] = count
    # Group unread
    memberships = GroupMember.query.filter_by(user_id=user_id, is_exited=False).all()
    for m in memberships:
        count = Message.query.filter(
            Message.group_id == m.group_id,
            Message.sender_id != user_id,
            ~Message.id.in_(
                db.session.query(MessageReceipt.message_id).filter(
                    MessageReceipt.user_id == user_id,
                    MessageReceipt.read_at != None
                )
            )
        ).count()
        counts[f'g_{m.group_id}'] = count
    return jsonify(counts)


# â”€â”€â”€ CHAT MEDIA UPLOAD â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
import uuid as _uuid

@friends_bp.route('/api/chat/upload', methods=['POST'])
@jwt_required()
def chat_media_upload():
    """Upload voice/image/file blobs for chat messages. Returns a persistent server URL."""
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    f = request.files['file']

    ext = os.path.splitext(f.filename or '')[1].lower()
    mime = f.mimetype or ''
    # Assign sensible extension for recorded blobs that arrive with no extension
    if not ext or ext == '.bin':
        if 'webm' in mime or 'ogg' in mime: ext = '.webm'
        elif 'mp4'  in mime: ext = '.mp4'
        elif 'wav'  in mime: ext = '.wav'
        elif 'mpeg' in mime or 'mp3' in mime: ext = '.mp3'
        elif 'image' in mime: ext = '.jpg'
        else: ext = '.webm'   # default for MediaRecorder blobs

    unique_name = f"chat_{_uuid.uuid4().hex}{ext}"
    save_dir    = os.path.join(app.config['UPLOAD_FOLDER'], 'chat')
    os.makedirs(save_dir, exist_ok=True)
    f.save(os.path.join(save_dir, unique_name))

    if mime.startswith('audio') or ext in ('.webm','.mp3','.ogg','.wav','.m4a','.opus'):
        media_type = 'voice'
    elif mime.startswith('video') or ext in ('.mp4','.mov','.avi','.mkv'):
        media_type = 'video'
    elif mime.startswith('image') or ext in ('.jpg','.jpeg','.png','.gif','.webp'):
        media_type = 'image'
    else:
        media_type = 'file'

    return jsonify({"url": f"/uploads/chat/{unique_name}", "type": media_type, "filename": f.filename or unique_name})


# â”€â”€â”€ PRIVACY SETTINGS â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
@friends_bp.route('/api/user/privacy', methods=['GET'])
@jwt_required()
def get_privacy():
    user_id = int(get_jwt_identity())
    user = db.session.get(User, user_id)
    if not user:
        return jsonify({"error": "User not found"}), 404
    return jsonify({
        "lastSeen":     getattr(user, 'privacy_last_seen',     'everyone'),
        "profilePhoto": getattr(user, 'privacy_profile_photo', 'everyone'),
        "about":        getattr(user, 'privacy_about',         'everyone'),
        "readReceipts": getattr(user, 'privacy_read_receipts', True),
    })

@friends_bp.route('/api/user/privacy', methods=['POST'])
@jwt_required()
def update_privacy():
    user_id = int(get_jwt_identity())
    user = db.session.get(User, user_id)
    if not user:
        return jsonify({"error": "User not found"}), 404
    data = request.json or {}
    valid = {'everyone', 'contacts', 'nobody'}
    if 'lastSeen'     in data and data['lastSeen']     in valid: user.privacy_last_seen     = data['lastSeen']
    if 'profilePhoto' in data and data['profilePhoto'] in valid: user.privacy_profile_photo = data['profilePhoto']
    if 'about'        in data and data['about']        in valid: user.privacy_about         = data['about']
    if 'readReceipts' in data: user.privacy_read_receipts = bool(data['readReceipts'])
    db.session.commit()
    return jsonify({"success": True})


app.register_blueprint(friends_bp)

# serve_chat_media moved above serve_upload to avoid route shadowing

if __name__ == "__main__":
    # Create necessary directories
    os.makedirs("uploads", exist_ok=True)
    os.makedirs("outputs", exist_ok=True)
    
    # Startup info
    print("Starting TalkToText Pro Backend Server...")
    print("=" * 60)
    print("Upload folder:", app.config['UPLOAD_FOLDER'])
    print("Database:", app.config['SQLALCHEMY_DATABASE_URI'])
    print("JWT Secret configured:", bool(app.config['JWT_SECRET_KEY']))
    print("Gemini API configured:", bool(gemini_api_key))
    print("AssemblyAI configured:", bool(os.getenv("ASSEMBLYAI_API_KEY")))
    print("SocketIO Meeting System: Enabled")
    
    print("\nEmail Configuration:")
    print(f"   SendGrid API: {'Set' if os.getenv('SENDGRID_API_KEY') else 'Not set'}")
    print(f"   SMTP Username: {'Set' if os.getenv('SENDGRID_FROM_EMAIL') else 'Not set'}")
    print(f"   SMTP Password: {'Set' if os.getenv('SMTP_PASSWORD') else 'Not set'}")
    
    print("\nAPI Endpoints Available:")
    endpoints = [
        "POST /api/auth/register      - User registration",
        "POST /api/auth/login         - User login", 
        "GET  /api/auth/check         - Token validation",
        "GET  /api/auth/profile       - Get user profile",
        "PUT  /api/auth/profile       - Update user profile",
        "PUT  /api/auth/profile/password - Change password",
        "DELETE /api/auth/profile    - Delete account",
        "POST /api/auth/verify-credentials - Live credential check",
        "POST /api/upload             - File upload",
        "POST /api/process/<id>       - Start processing",
        "GET  /api/processing-status/<id> - Check processing status",
        "GET  /api/meetings           - List user meetings",
        "GET  /api/meetings/<id>      - Get specific meeting",
        "PUT  /api/meetings/<id>/notes - Update meeting notes",
        "DELETE /api/meetings/<id>    - Delete meeting",
        "POST /api/translate          - Text translation",
        "GET  /api/export/<id>/<format> - Export notes (PDF/Word)",
        "POST /api/chat               - AI chat assistant",
        "POST /api/send-email         - Send email with attachments",
        "GET  /api/stats              - User statistics",
        "GET  /api/health             - Detailed health check",
        # New meeting endpoints
        "POST /api/meetings/create    - Create new meeting room",
        "GET  /api/meetings/<id>      - Get meeting details",
        "POST /api/meetings/<id>/join - Join meeting",
        "POST /api/meetings/<id>/recordings - Save meeting recording",
        "GET  /api/meetings/<id>/transcript - Get meeting transcript",
        "GET  /api/meetings/<id>/analytics - Get meeting analytics",
        "GET  /api/meetings/user/recent - Get user's recent meetings",
        "GET  /api/meetings/public    - Get public meetings",
        # WebSocket endpoints
        "WS   /socket.io              - Real-time meeting WebSocket"
    ]
    
    for endpoint in endpoints:
        print(f"   {endpoint}")
    
    print("\nKey Features Added:")
    features = [
        "Real-time video/audio meetings with WebRTC",
        "Interactive quiz system with leaderboards",
        "AI-powered meeting assistant",
        "Ghost replay of past moments",
        "Command-based meeting control",
        "Screen sharing and recording",
        "Real-time chat with reactions",
        "Participant management",
        "Meeting analytics and insights",
        "Automatic cleanup of inactive meetings",
        "Redis-based session management",
        "Secure password-protected meetings",
        "Public/private meeting rooms"
    ]
    
    for feature in features:
        print(f"   {feature}")
    
    print("\n" + "=" * 60)
    print("Server starting on http://0.0.0.0:5000")
    print("CORS enabled for all origins")
    print("SocketIO enabled for real-time communication")
    print("=" * 60)
    
    # Run the Flask app with SocketIO
    with app.app_context():
        db.create_all()
        print("Database tables verified/created")

    socketio.run(app, debug=True, host="0.0.0.0", port=5000, allow_unsafe_werkzeug=True)